"""
KimFam Hub — Document Embedding Pipeline
Canonical ingestion script. Always reference KimFam-RAG-Ingestion.md before modifying.
"""

import os
import sys
import json
import hashlib
import re
import logging
from pathlib import Path
from datetime import date

import mammoth
import pdfplumber
import tiktoken
import chromadb
# from openai import OpenAI  # removed: switched to Gemini

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger("embed_documents")

# Paths — derive from this file's location so staging/local embed their own tree
# (matches main.py's DOCS_DIR = Path(__file__).parent/"docs").
BASE_DIR = Path(__file__).resolve().parent
DOCS_DIR = BASE_DIR / "docs"
CHROMA_DIR = BASE_DIR / "data" / "chroma"
MANIFEST_PATH = BASE_DIR / "data" / "embed_manifest.json"
COLLECTION_NAME = "kimfam_docs"

# Chunking
CHUNK_TOKENS = 500
OVERLAP_TOKENS = 50
EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"  # local, 384-dim, no API

# doc_type detection from folder name
FOLDER_TO_DOCTYPE = {
    "minutes": "minutes",
    "governance": "constitution",
    "projects": "proposal",
    "financial": "report",
    "receipts": "receipt",
}

# client = OpenAI(...)  # removed: switched to Gemini
enc = tiktoken.get_encoding("cl100k_base")


def slugify(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    return text.strip("-")


def md5(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
        return result.value
    elif suffix == ".pdf":
        with pdfplumber.open(path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    elif suffix == ".pptx":
        # Presentations: pull every text frame + table cell, slide by slide.
        from pptx import Presentation
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            out.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            out.append(" | ".join(cells))
        return "\n".join(out)
    elif suffix == ".xlsx":
        # Spreadsheets: flatten each sheet to "sheet | row values" lines.
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    out.append(" | ".join(vals))
        return "\n".join(out)
    return ""


def chunk_text(text: str) -> list[str]:
    tokens = enc.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = start + CHUNK_TOKENS
        chunk_tokens = tokens[start:end]
        chunks.append(enc.decode(chunk_tokens))
        if end >= len(tokens):
            break
        start += CHUNK_TOKENS - OVERLAP_TOKENS
    return [c.strip() for c in chunks if c.strip()]


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Embed locally via sentence-transformers. Returns 384-dim vectors."""
    from sentence_transformers import SentenceTransformer
    log.info(f"Loading embedder: {EMBED_MODEL}")
    model = SentenceTransformer(EMBED_MODEL)
    vectors = model.encode(texts, normalize_embeddings=True, show_progress_bar=False)
    return [[float(x) for x in v] for v in vectors]


def extract_date(filename: str) -> str:
    # Try to extract a date like May_24_2026 or 2024 from filename
    patterns = [
        r"(\w+)_(\d{1,2})_(\d{4})",  # Month_DD_YYYY
        r"(\w+)\s+(\d{1,2})\s+(\d{4})",  # Month DD YYYY
        r"(\d{4})",  # Just a year
    ]
    month_map = {
        "jan": "01", "feb": "02", "mar": "03", "apr": "04",
        "may": "05", "jun": "06", "jul": "07", "aug": "08",
        "sep": "09", "oct": "10", "nov": "11", "dec": "12",
    }
    m = re.search(r"(\w+)[_\s](\d{1,2})[_\s](\d{4})", filename, re.IGNORECASE)
    if m:
        month = month_map.get(m.group(1).lower()[:3])
        if month:
            return f"{m.group(3)}-{month}-{int(m.group(2)):02d}"
    m = re.search(r"(\d{4})", filename)
    if m:
        return m.group(1)
    return str(date.today())


def get_doc_type(path: Path) -> str:
    folder = path.parent.name
    return FOLDER_TO_DOCTYPE.get(folder, "document")


def load_manifest() -> dict:
    if MANIFEST_PATH.exists():
        with open(MANIFEST_PATH) as f:
            return json.load(f)
    return {}


def save_manifest(manifest: dict):
    MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(MANIFEST_PATH, "w") as f:
        json.dump(manifest, f, indent=2)


def run():
    CHROMA_DIR.mkdir(parents=True, exist_ok=True)
    chroma = chromadb.PersistentClient(path=str(CHROMA_DIR))
    collection = chroma.get_or_create_collection(COLLECTION_NAME, metadata={"hnsw:space": "cosine"})

    manifest = load_manifest()

    # Find all supported documents in docs/ (docx, pdf, pptx, xlsx). Skip Office
    # lock/temp files (~$...) and archived prior versions in _versions/ (those are
    # hidden in the Documents nav and must not pollute RAG search).
    source_files = [
        p for ext in ("*.docx", "*.pdf", "*.pptx", "*.xlsx")
        for p in DOCS_DIR.rglob(ext)
        if not p.name.startswith("~") and "_versions" not in p.parts
    ]
    on_disk = {str(p.relative_to(BASE_DIR)) for p in source_files}

    files_processed = 0
    chunks_added = 0
    files_skipped = 0

    # --- Handle deleted files ---
    deleted_keys = [k for k in manifest if k not in on_disk]
    for key in deleted_keys:
        old_ids = manifest[key].get("chunk_ids", [])
        if old_ids:
            collection.delete(ids=old_ids)
            log.info(f"Deleted {len(old_ids)} chunks for removed file: {key}")
        del manifest[key]

    # --- Process each file ---
    for fpath in source_files:
        rel = str(fpath.relative_to(BASE_DIR))
        file_hash = md5(fpath)

        if rel in manifest and manifest[rel]["md5"] == file_hash:
            files_skipped += 1
            continue

        log.info(f"Processing: {rel}")
        text = extract_text(fpath)
        if not text.strip():
            log.warning(f"  No text extracted from {rel}, skipping")
            continue

        chunks = chunk_text(text)
        if not chunks:
            continue

        slug = slugify(fpath.stem)
        doc_type = get_doc_type(fpath)
        doc_date = extract_date(fpath.name)
        chunk_ids = [f"{slug}_{i}" for i in range(len(chunks))]

        # Delete old chunks if updating
        if rel in manifest:
            old_ids = manifest[rel].get("chunk_ids", [])
            if old_ids:
                collection.delete(ids=old_ids)

        embeddings = embed_texts(chunks)
        metadatas = [
            {
                "source": rel.replace("docs/", ""),
                "doc_type": doc_type,
                "date": doc_date,
                "chunk_index": i,
            }
            for i in range(len(chunks))
        ]

        collection.upsert(
            ids=chunk_ids,
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadatas,
        )

        manifest[rel] = {
            "md5": file_hash,
            "embedded_at": str(date.today()),
            "chunk_ids": chunk_ids,
            "chunk_count": len(chunks),
        }

        files_processed += 1
        chunks_added += len(chunks)
        log.info(f"  -> {len(chunks)} chunks embedded ({doc_type}, {doc_date})")

    save_manifest(manifest)
    log.info(
        f"Done. {files_processed} files processed, {chunks_added} chunks added, "
        f"{files_skipped} files skipped (unchanged)."
    )
    return files_processed, chunks_added


if __name__ == "__main__":
    run()
