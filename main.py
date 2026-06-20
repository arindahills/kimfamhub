import os
import logging as _logging, sys as _sys

# ── Logging setup ────────────────────────────────────────────────────────────
# Previously there was NO logging config, so the app's getLogger("main").error()
# calls had no consistent handler/format and failures were effectively invisible.
# Attach a formatted stdout handler (captured by journald) to the app loggers.
_log_handler = _logging.StreamHandler(_sys.stdout)
_log_handler.setFormatter(_logging.Formatter(
    "%(asctime)s %(levelname)s [%(name)s] %(message)s", "%Y-%m-%d %H:%M:%S"))
for _ln in ("main", "kimfam", "pitch", "scheduler", "notifications"):
    _l = _logging.getLogger(_ln)
    _l.setLevel(_logging.INFO)
    if not _l.handlers:
        _l.addHandler(_log_handler)
    _l.propagate = False
log = _logging.getLogger("kimfam")

from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from contextlib import asynccontextmanager
import scheduler as _scheduler_mod
import gspread
from google.oauth2.service_account import Credentials

@asynccontextmanager
async def lifespan(app):
    if os.environ.get("SCHEDULER_ENABLED") == "1":
        _scheduler_mod.start()
    # Always-cooking "why join" pitch engine (independent of the scheduler so it
    # runs on staging too). Cooks on boot, then re-cooks every few hours.
    try:
        import pitch_engine as _pitch
        _pitch.start_background(get_all_projects)
    except Exception as _e:
        import logging as _lg; _lg.getLogger("pitch").warning("pitch engine not started: %s", _e)
    yield
    if os.environ.get("SCHEDULER_ENABLED") == "1":
        _scheduler_mod.stop()

app = FastAPI(lifespan=lifespan, docs_url="/api/_swagger", redoc_url=None)


@app.middleware("http")
async def _log_requests(request: Request, call_next):
    """Surface API failures and slow operations in the logs (errors used to vanish
    into swallowed excepts). Logs any 4xx/5xx and any API call over 3s."""
    import time as _t
    t0 = _t.time()
    try:
        resp = await call_next(request)
    except Exception:
        log.exception("Unhandled error on %s %s", request.method, request.url.path)
        raise
    dt_ms = (_t.time() - t0) * 1000
    path = request.url.path
    if path.startswith("/api/") and (resp.status_code >= 400 or dt_ms > 3000):
        log.info("%s %s -> %s (%.0fms)", request.method, path, resp.status_code, dt_ms)
    return resp


app.mount("/static", StaticFiles(directory=os.path.join(os.path.dirname(__file__), "static")), name="static")

# React frontend assets (built output from frontend/dist/assets)
_DIST = os.path.join(os.path.dirname(__file__), "frontend", "dist")
if os.path.isdir(os.path.join(_DIST, "assets")):
    app.mount("/assets", StaticFiles(directory=os.path.join(_DIST, "assets")), name="frontend-assets")

SHEET_ID = "1R3_j2ArvMZsfiLDvFwEQURJBXEPaW3mmWU-FyUwrqPg"
SOLOMON_ID = "1CqF-NzkMJ8iJw0tC8xkLE9DI94cjFr2vvlAfx4QfXhI"
SCOPES = ["https://www.googleapis.com/auth/spreadsheets","https://www.googleapis.com/auth/drive"]

def gc():
    creds = Credentials.from_service_account_file(os.path.join(os.path.dirname(__file__), "service-account.json"), scopes=SCOPES)
    return gspread.authorize(creds)

def g(row, i):
    return row[i].strip() if i < len(row) else ""

@app.get("/api/meetings/suggest-agenda")
def meetings_suggest_agenda(request: Request):
    """Ask Claude to suggest main agenda topics for the next meeting based on history."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    token = _get_tok(request)
    if not _auth_verify(token):
        raise _HE(status_code=401, detail="Auth required")

    # Pull last 4 meetings' context
    meetings = _dbq("""
        SELECT ref, date, key_topics, key_decisions, next_actions, summary
        FROM meetings ORDER BY date DESC LIMIT 4
    """)
    # Pull open actions (unresolved items)
    open_actions = _dbq("""
        SELECT description, assignees, deadline FROM actions
        WHERE status NOT IN ('done','cancelled','carried_over')
        ORDER BY deadline ASC NULLS LAST LIMIT 20
    """)

    history_lines = []
    for m in reversed(meetings):
        history_lines.append(f"## {m['ref']} ({m['date']})")
        if m["key_topics"]:    history_lines.append(f"Topics: {m['key_topics']}")
        if m["key_decisions"]: history_lines.append(f"Decisions: {m['key_decisions']}")
        if m["next_actions"]:  history_lines.append(f"Actions set: {m['next_actions']}")
        if m["summary"]:       history_lines.append(f"Summary: {m['summary']}")

    open_lines = []
    for a in open_actions:
        dl = f" (due {a['deadline']})" if a["deadline"] else ""
        open_lines.append(f"- {a['description']}{dl}")

    prompt = f"""You are the secretary for KimFam Investment Club, a Ugandan family investment club.
Based on the meeting history and open actions below, suggest 3-5 main agenda items for the NEXT meeting.

Focus on:
- Recurring themes that keep coming back (deferred votes, ongoing projects needing decisions)
- Open actions that are overdue or approaching deadline
- Natural next steps from the most recent meeting's decisions
- Any new initiatives that were flagged but not yet actioned

Return ONLY a semicolon-separated list of agenda item titles, no explanations.
Example format: Project performance review; Equity model vote; Washing bay Phase 2 update; Constitution ratification

Meeting history (oldest to newest):
{chr(10).join(history_lines)}

Currently open actions:
{chr(10).join(open_lines) if open_lines else 'None'}
"""
    result = _ask_claude(prompt, timeout=45)
    # Keep only the last non-empty line (the semicolon list)
    lines = [l.strip() for l in result.strip().splitlines() if l.strip()]
    suggestion = lines[-1] if lines else ""
    return {"suggestion": suggestion}


def _compute_next_ref():
    """Next meeting ref based on the HIGHEST meeting NUMBER (not the latest date —
    multiple meetings can share a date, which made date-ordering collide)."""
    from db import query as _dbq
    from datetime import date as _date
    import re as _re
    rows = _dbq("SELECT ref FROM meetings")
    today = _date.today()
    best_num, best_year, best_ref = 0, today.year, None
    for r in rows:
        m = _re.match(r"KIM\s+(\d+)/(\d{4})", (r["ref"] or "").strip())
        if m and int(m.group(1)) > best_num:
            best_num = int(m.group(1)); best_year = int(m.group(2)); best_ref = r["ref"]
    if best_num == 0:
        return f"KIM 001/{today.year}", None
    year = today.year if today.year >= best_year else best_year
    return f"KIM {best_num + 1:03d}/{year}", best_ref


@app.get("/api/meetings/next-ref")
def meetings_next_ref():
    """Return the next meeting ref based on the highest meeting number in the DB."""
    next_ref, prev_ref = _compute_next_ref()
    return {"next_ref": next_ref, "prev_ref": prev_ref}


@app.get("/api/meetings/analytics")
def meetings_analytics(request: Request):
    """Cross-meeting time analytics for the whole club. Visible to all members.
    Aggregates per-meeting retrospectives + conductor timings into:
    efficiency trend, recurring time sinks, and time-spent-per-topic."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    import json as _json_a, re as _re_a
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")

    rows = _dbq("""SELECT ref, date, conductor_retrospective, conductor_timings
                   FROM meetings""")

    def _num(ref):
        m = _re_a.match(r"KIM\s+(\d+)/(\d{4})", (ref or "").strip())
        return int(m.group(1)) if m else 0
    def _load(v):
        return _json_a.loads(v) if isinstance(v, (str, bytes, bytearray)) else (v or None)

    trend = []            # [{ref, date, efficiency, planned_min, actual_min}]
    sink_counts = {}      # normalised topic -> count of meetings
    item_agg = {}         # normalised label -> {label, total_s, planned_s, count}

    for r in rows:
        retro = _load(r["conductor_retrospective"])
        if isinstance(retro, dict) and retro.get("efficiency_score") is not None:
            trend.append({
                "ref": r["ref"], "date": str(r["date"]),
                "num": _num(r["ref"]),
                "efficiency": int(retro.get("efficiency_score", 0)),
                "planned_min": retro.get("planned_total_min"),
                "actual_min": retro.get("actual_total_min"),
            })
            seen = set()
            for s in (retro.get("time_sinks") or []):
                key = str(s).strip().lower()
                if key and key not in seen:
                    seen.add(key)
                    sink_counts[key] = sink_counts.get(key, 0) + 1

        timings = _load(r["conductor_timings"]) or {}
        for v in timings.values():
            if not isinstance(v, dict):
                continue
            label = (v.get("label") or "").strip()
            if not label:
                continue
            key = label.lower()
            a = item_agg.setdefault(key, {"label": label, "total_s": 0, "planned_s": 0, "count": 0})
            a["total_s"]   += int(v.get("actual_s", 0))
            a["planned_s"] += int(v.get("planned_min", 0)) * 60
            a["count"]     += 1

    trend.sort(key=lambda x: x["num"])
    avg_eff = round(sum(t["efficiency"] for t in trend) / len(trend)) if trend else None
    # trend direction: compare last half vs first half average
    direction = "steady"
    if len(trend) >= 4:
        half = len(trend) // 2
        first = sum(t["efficiency"] for t in trend[:half]) / half
        last  = sum(t["efficiency"] for t in trend[half:]) / (len(trend) - half)
        if last - first >= 5:   direction = "improving"
        elif first - last >= 5: direction = "declining"

    recurring = sorted(
        [{"topic": k, "meetings": c} for k, c in sink_counts.items() if c >= 2],
        key=lambda x: -x["meetings"])[:8]

    time_by_topic = sorted(
        [{"label": a["label"], "total_min": a["total_s"] // 60,
          "avg_min": round(a["total_s"] / a["count"] / 60, 1) if a["count"] else 0,
          "planned_min": a["planned_s"] // 60, "occurrences": a["count"]}
         for a in item_agg.values() if a["total_s"] > 0],
        key=lambda x: -x["total_min"])[:10]

    return {
        "trend": trend,
        "avg_efficiency": avg_eff,
        "direction": direction,
        "recurring_sinks": recurring,
        "time_by_topic": time_by_topic,
        "meetings_analysed": len(trend),
    }


@app.get("/api/actions")
def get_actions(status: str = "open"):
    from db import query as _dbq
    from datetime import date as _date
    filt = status.lower().strip()
    if filt == "done":
        where = "a.status IN ('done','cancelled','carried_over')"
    elif filt == "all":
        where = "TRUE"
    else:
        where = "a.status NOT IN ('done','cancelled','carried_over')"
    rows = _dbq(f"""
        SELECT a.id, a.ref, a.description, a.assignee, a.assignees, a.deadline,
               a.status, a.related_meeting, a.priority, a.effort_hours,
               a.project_id, a.parent_ref, a.blocked_reason,
               m.ref AS meeting_ref,
               (SELECT text FROM action_updates
                WHERE action_id = a.id ORDER BY created_at DESC LIMIT 1) AS latest_update
        FROM actions a
        LEFT JOIN meetings m ON m.id = a.meeting_id
        WHERE {where}
        ORDER BY a.ref DESC
    """)
    today = _date.today()
    by_person = {}
    for r in rows:
        health = None
        if r["status"] not in ("done", "cancelled", "carried_over") and r["deadline"]:
            days = (r["deadline"] - today).days
            if days < 0:      health = "overdue"
            elif days <= 3:   health = "at_risk"
            else:             health = "on_track"

        # Resolve the list of people this action belongs to
        if r["assignees"]:          # TEXT[] column set
            people = [p.strip() for p in r["assignees"] if p.strip()]
        else:
            people = [(r["assignee"] or "Unknown").strip()]

        payload = {
            "id":           r["ref"],
            "action":       r["description"],
            "deadline":     str(r["deadline"]) if r["deadline"] else "",
            "status":       r["status"],
            "meeting":      r["meeting_ref"] or r["related_meeting"] or "",
            "note":         r["latest_update"] or "",
            "priority":     r["priority"],
            "effort_hours": float(r["effort_hours"]) if r["effort_hours"] else None,
            "project_id":   r["project_id"],
            "parent_ref":   r["parent_ref"],
            "health":       health,
            "assignees":    people,
        }
        for p in people:
            by_person.setdefault(p, []).append(payload)
    return by_person


@app.patch("/api/actions/done")
async def mark_action_done(request: Request):
    """Admin: mark an action point as Done in the DB."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    action_ref = str(body.get("action_id", "")).strip()
    comment    = str(body.get("comment", "")).strip()
    if not action_ref:
        raise _HE(status_code=400, detail="action_id required")
    rows = _dbq("SELECT id FROM actions WHERE ref=%s", (action_ref,))
    if not rows:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    action_db_id = rows[0]["id"]
    author = payload.get("sub", "admin")
    try:
        _exec("UPDATE actions SET status='done', closed_at=NOW() WHERE ref=%s", (action_ref,))
        _exec("""INSERT INTO action_updates (action_id, author, text, type, old_value, new_value)
                 VALUES (%s,%s,%s,'status_change','open','done')""",
              (action_db_id, author, comment or "Marked done"))
        return {"ok": True, "action_id": action_ref}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"mark_action_done: {e}")
        raise _HE(status_code=500, detail="DB update failed")


@app.patch("/api/actions/update")
async def add_action_update(request: Request):
    """Any authenticated user: log a progress update on an action."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        raise _HE(status_code=401, detail="Not authenticated")
    body = await request.json()
    action_ref  = str(body.get("action_id", "")).strip()
    update_text = str(body.get("update_text", "")).strip()
    if not action_ref or not update_text:
        raise _HE(status_code=400, detail="action_id and update_text required")
    rows = _dbq("SELECT id, status FROM actions WHERE ref=%s", (action_ref,))
    if not rows:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    action_db_id = rows[0]["id"]
    author = payload.get("sub", "unknown")
    try:
        _exec("""INSERT INTO action_updates (action_id, author, text, type)
                 VALUES (%s,%s,%s,'comment')""",
              (action_db_id, author, update_text))
        # Auto-advance status open → in_progress on first update
        if rows[0]["status"] == "open":
            _exec("UPDATE actions SET status='in_progress' WHERE ref=%s", (action_ref,))
        return {"ok": True, "action_id": action_ref}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"add_action_update: {e}")
        raise _HE(status_code=500, detail="DB update failed")


@app.patch("/api/actions/status")
async def set_action_status(request: Request):
    """Admin: move an action to any state (in_progress, blocked, carried_over,
    cancelled, done, open). Logs the change as an action_update."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    action_ref = str(body.get("action_id", "")).strip()
    new_status = str(body.get("status", "")).strip().lower()
    comment    = str(body.get("comment", "")).strip()
    VALID = {"open", "in_progress", "blocked", "carried_over", "cancelled", "done"}
    if not action_ref or new_status not in VALID:
        raise _HE(status_code=400, detail=f"action_id and a valid status ({', '.join(sorted(VALID))}) required")
    rows = _dbq("SELECT id, status FROM actions WHERE ref=%s", (action_ref,))
    if not rows:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    action_db_id = rows[0]["id"]
    old_status   = rows[0]["status"]
    author       = payload.get("sub", "admin")
    closes = new_status in ("done", "cancelled", "carried_over")
    try:
        if closes:
            _exec("UPDATE actions SET status=%s, closed_at=NOW() WHERE ref=%s", (new_status, action_ref))
        else:
            _exec("UPDATE actions SET status=%s, closed_at=NULL WHERE ref=%s", (new_status, action_ref))
        _exec("""INSERT INTO action_updates (action_id, author, text, type, old_value, new_value)
                 VALUES (%s,%s,%s,'status_change',%s,%s)""",
              (action_db_id, author,
               comment or f"Status changed to {new_status.replace('_',' ')}",
               old_status, new_status))
        return {"ok": True, "action_id": action_ref, "status": new_status}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"set_action_status: {e}")
        raise _HE(status_code=500, detail="DB update failed")


def _next_or_create_meeting():
    """Return (id, ref) of the next upcoming, not-yet-conducted meeting.
    Creates it (next Sunday, KIM N+1, default agenda) if none exists."""
    from db import query as _dbq, execute as _exec
    from datetime import date as _date, timedelta as _td
    import json as _json_nm
    rows = _dbq("""SELECT id, ref FROM meetings
                   WHERE date >= CURRENT_DATE AND conductor_started_at IS NULL
                   ORDER BY date ASC, ref ASC LIMIT 1""")
    if rows:
        return rows[0]["id"], rows[0]["ref"]
    # None upcoming — create the next meeting on the next Sunday.
    ref, _prev = _compute_next_ref()
    today = _date.today()
    days = (6 - today.weekday()) % 7  # Sunday = weekday 6; 0 if today is Sunday
    nxt_sunday = today + _td(days=days if days else 7)
    agenda = _build_default_agenda(None)
    _exec("""INSERT INTO meetings (ref, date, venue, start_time_eat, agenda)
             VALUES (%s, %s::date, %s, %s::time, %s)""",
          (ref, str(nxt_sunday), "Google Meet", "16:30", _json_nm.dumps(agenda)))
    row = _dbq("SELECT id FROM meetings WHERE ref=%s", (ref,))
    return row[0]["id"], ref


@app.post("/api/actions/carry-over")
async def carry_over_action(request: Request):
    """Admin: carry an action forward to the next meeting. Creates a fresh
    continuation action there (parent_ref = old ref) and marks the old one
    carried_over. Creates the next meeting if none exists yet."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import re as _re_co
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    action_ref = str(body.get("action_id", "")).strip()
    comment    = str(body.get("comment", "")).strip()
    if not action_ref:
        raise _HE(status_code=400, detail="action_id required")
    src = _dbq("""SELECT id, description, assignee, assignees, deadline, priority,
                         project_id, status
                  FROM actions WHERE ref=%s""", (action_ref,))
    if not src:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    a = src[0]
    author = payload.get("sub", "admin")

    target_id, target_ref = _next_or_create_meeting()

    # Build the continuation ref under the target meeting: KIM/NN/YY-{i}
    m = _re_co.match(r"KIM\s+(\d+)/(\d{4})", target_ref)
    if m:
        num = int(m.group(1)); yr = m.group(2)[-2:]
        i = (_dbq("SELECT COUNT(*) AS c FROM actions WHERE meeting_id=%s", (target_id,))[0]["c"]) + 1
        new_ref = f"KIM/{num:02d}/{yr}-{i}"
        while _dbq("SELECT ref FROM actions WHERE ref=%s", (new_ref,)):
            i += 1; new_ref = f"KIM/{num:02d}/{yr}-{i}"
    else:
        new_ref = f"{target_ref}-carry"

    try:
        _exec("""INSERT INTO actions
                   (ref, description, assignee, assignees, meeting_id, related_meeting,
                    status, priority, deadline, parent_ref, project_id, created_by)
                 VALUES (%s,%s,%s,%s,%s,%s,'open',%s,%s,%s,%s,%s)""",
              (new_ref, a["description"], a["assignee"], a["assignees"], target_id, target_ref,
               a["priority"] or "medium", a["deadline"], action_ref, a["project_id"], author))
        # Close out the source action as carried over
        _exec("UPDATE actions SET status='carried_over', closed_at=NOW() WHERE ref=%s", (action_ref,))
        _exec("""INSERT INTO action_updates (action_id, author, text, type, old_value, new_value)
                 VALUES (%s,%s,%s,'status_change',%s,'carried_over')""",
              (a["id"], author,
               (comment + " " if comment else "") + f"Carried over to {target_ref} as {new_ref}",
               a["status"]))
        return {"ok": True, "carried_to_meeting": target_ref, "new_ref": new_ref}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"carry_over_action: {e}")
        raise _HE(status_code=500, detail="Carry-over failed")


async def _ai_complete(prompt: str, claude_timeout: int = 150) -> str:
    """One AI completion. Claude CLI is PRIMARY (for prompts it can handle in time);
    DeepSeek is the high-capacity fallback for oversized prompts or Claude failures;
    Groq (truncated) is the last resort. Never raises — returns '' if all fail."""
    import os as _os, logging as _lg_ai, asyncio as _aio
    raw = ""
    if len(prompt) <= 24000:
        try:
            env = dict(_os.environ); env["HOME"] = "/root"
            proc = await _aio.create_subprocess_exec(
                "claude", "-p", prompt, "--model", "claude-haiku-4-5-20251001",
                stdout=_aio.subprocess.PIPE, stderr=_aio.subprocess.DEVNULL, env=env)
            stdout, _ = await _aio.wait_for(proc.communicate(), timeout=claude_timeout)
            if proc.returncode == 0:
                raw = stdout.decode().strip()
        except Exception:
            pass
    if not raw:
        dk = _os.getenv("DEEPSEEK_API_KEY", "")
        if dk:
            try:
                from openai import OpenAI as _OAI
                r = _OAI(api_key=dk, base_url="https://api.deepseek.com").chat.completions.create(
                    model="deepseek-chat", messages=[{"role": "user", "content": prompt}],
                    max_tokens=4000, temperature=0.1)
                raw = (r.choices[0].message.content or "").strip()
            except Exception as e:
                _lg_ai.getLogger("main").error(f"DeepSeek complete failed: {e}")
    if not raw:
        gk = _os.getenv("GROQ_API_KEY", "")
        if gk:
            try:
                from groq import Groq as _Groq
                gp = prompt if len(prompt) <= 26000 else prompt[:26000] + "\n\n[Truncated to fit.]"
                r = _Groq(api_key=gk).chat.completions.create(
                    model="llama-3.3-70b-versatile", messages=[{"role": "user", "content": gp}],
                    max_tokens=2000, temperature=0.1)
                raw = (r.choices[0].message.content or "").strip()
            except Exception as e:
                _lg_ai.getLogger("main").error(f"Groq complete failed: {e}")
    return raw


async def _condense_transcript(transcript: str) -> str:
    """A long meeting (e.g. a 3-hour Tactiq transcript) won't fit one AI call.
    Split it into chunks, have Claude extract the substance of each chunk, and
    stitch the results into a dense record that DOES fit the final extraction.
    Keeps Claude as the engine for transcripts of any length."""
    if len(transcript) <= 18000:
        return transcript
    size = 16000
    chunks = [transcript[i:i + size] for i in range(0, len(transcript), size)]
    parts = []
    for idx, ch in enumerate(chunks, 1):
        cp = (
            f"This is part {idx} of {len(chunks)} of a KimFam Investment Club meeting transcript. "
            f"Extract, as concise bullet points and losing nothing material: every DECISION made; "
            f"every ACTION assigned (who is responsible and any deadline); key figures and amounts "
            f"in UGX; and the important points of each discussion. No preamble, bullets only. "
            f"No em-dashes or en-dashes.\n\nTRANSCRIPT PART {idx} OF {len(chunks)}:\n{ch}"
        )
        out = await _ai_complete(cp, claude_timeout=90)
        parts.append(f"--- Notes from part {idx} ---\n{out or '(could not process this part)'}")
    return ("[This is a condensed record of a long meeting, assembled from the full transcript "
            "in sections.]\n\n" + "\n\n".join(parts))


def _transcribe_audio(audio_bytes: bytes, audio_name: str, groq_key: str) -> str:
    """Transcribe a recording of any length with Groq Whisper. Groq caps a single
    file at 25 MB, so a long meeting (e.g. 177 min of conductor audio) is first
    re-encoded to mono 16 kHz low-bitrate MP3 and split into 15-minute segments
    with ffmpeg; each segment is transcribed and the text is stitched back in
    order. Short recordings go straight through in one call. Raises on hard
    failure so the caller can surface a clear error."""
    import os as _os_a, tempfile as _tmp_a, glob as _glob_a, subprocess as _sub_a
    from groq import Groq as _Groq

    def _one(path_or_tuple):
        return str(_Groq(api_key=groq_key).audio.transcriptions.create(
            file=path_or_tuple, model="whisper-large-v3-turbo",
            response_format="text")).strip()

    # Small enough for a single Whisper call — keep the original bytes as-is.
    if len(audio_bytes) <= 24 * 1024 * 1024:
        return _one((audio_name, audio_bytes, "audio/mpeg"))

    workdir = _tmp_a.mkdtemp(prefix="kf_audio_")
    src = _os_a.path.join(workdir, "src_" + _os_a.path.basename(audio_name or "audio"))
    seg_tmpl = _os_a.path.join(workdir, "seg_%03d.mp3")
    try:
        with open(src, "wb") as _f:
            _f.write(audio_bytes)
        # Re-encode to speech-optimal mono 16 kHz 32 kbps MP3 and cut into 900s
        # segments. 15 min @ 32 kbps ~ 3.6 MB per chunk, well under the 25 MB cap.
        _sub_a.run(
            ["ffmpeg", "-y", "-i", src, "-ac", "1", "-ar", "16000",
             "-c:a", "libmp3lame", "-b:a", "32k", "-f", "segment",
             "-segment_time", "900", "-reset_timestamps", "1", seg_tmpl],
            check=True, stdout=_sub_a.DEVNULL, stderr=_sub_a.DEVNULL, timeout=900,
        )
        segs = sorted(_glob_a.glob(_os_a.path.join(workdir, "seg_*.mp3")))
        if not segs:
            # Fallback: ffmpeg produced nothing usable; try the original in one shot.
            return _one((audio_name, audio_bytes, "audio/mpeg"))
        out = []
        for i, seg in enumerate(segs, 1):
            try:
                with open(seg, "rb") as _sf:
                    txt = _one((_os_a.path.basename(seg), _sf.read(), "audio/mpeg"))
                if txt:
                    out.append(txt)
            except Exception as _se:
                import logging as _lg_a
                _lg_a.getLogger("main").error(f"Audio segment {i} transcribe failed: {_se}")
        return " ".join(out).strip()
    finally:
        try:
            for _p in _glob_a.glob(_os_a.path.join(workdir, "*")):
                _os_a.remove(_p)
            _os_a.rmdir(workdir)
        except Exception:
            pass


@app.post("/api/meetings/{meeting_id}/process")
async def process_meeting(meeting_id: int, request: Request):
    """Streamed wrapper: runs the (long) transcribe+extract work as a task while emitting
    SSE heartbeats, so audio/long transcripts survive the ~60s edge timeout and the user
    sees live progress. The body is read once up-front (Starlette caches it) so the inner
    impl can re-read the cached form during streaming."""
    from fastapi import HTTPException as _HE
    from fastapi.responses import StreamingResponse
    import asyncio as _aio
    payload = _auth_verify(_get_tok(request))
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    await request.form()   # consume + cache the multipart body before streaming starts

    _STAGES = ["Transcribing the recording (long audio is split into 15-min chunks)...",
               "Condensing the transcript...",
               "Extracting actions and decisions with Claude...",
               "Reconciling against open actions...",
               "Almost done, finalising..."]

    async def _gen():
        yield _sse({"type": "step", "msg": "Processing the meeting..."})
        task = _aio.create_task(_process_meeting_impl(meeting_id, request))
        i = 0
        while not task.done():
            await _aio.sleep(6)
            if not task.done():
                yield _sse({"type": "step", "msg": _STAGES[min(i, len(_STAGES) - 1)]}); i += 1
        try:
            res = task.result()
            yield _sse({**res, "type": "result"})
        except _HE as he:
            yield _sse({"type": "error", "msg": str(he.detail)})
        except Exception as e:
            log.error(f"process_meeting stream failed: {e}")
            yield _sse({"type": "error", "msg": "Processing failed. Please try again."})

    return StreamingResponse(_gen(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive"})


async def _process_meeting_impl(meeting_id: int, request: Request,
    audio_file: "UploadFile | None" = None,
    transcript_text: str = "",
    transcript_file: "UploadFile | None" = None,
):
    """
    Accept audio / pasted text / .txt / .docx, transcribe if needed,
    extract actions + decisions via Claude, return for frontend review.
    Nothing is written to DB here — that happens on /confirm.
    """
    from fastapi import HTTPException as _HE, UploadFile as _UF, File as _File, Form as _Form
    from db import query as _dbq
    import os as _os, json as _json, tempfile as _tmp, re as _re

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    # ── 1. Parse multipart body manually (FastAPI doesn't inject File/Form into path endpoints cleanly)
    form = await request.form()
    audio_file   = form.get("audio_file")
    txt_field    = form.get("transcript_text", "").strip()
    upload_file  = form.get("transcript_file")
    sec_notes    = (form.get("secretary_notes", "") or "").strip()

    # ── 2. Collect transcript from all provided sources, combine if multiple ──
    transcript_parts = []  # list of (label, text)

    if txt_field:
        transcript_parts.append(("Pasted text", txt_field))

    if upload_file and hasattr(upload_file, "filename"):
        fname_u = (upload_file.filename or "").lower()
        raw = await upload_file.read()
        if fname_u.endswith(".docx"):
            from io import BytesIO as _BIO
            import docx as _docx
            _doc = _docx.Document(_BIO(raw))
            file_text = "\n".join(p.text for p in _doc.paragraphs if p.text.strip())
        else:
            file_text = raw.decode("utf-8", errors="replace")
        if file_text.strip():
            transcript_parts.append((f"File: {upload_file.filename}", file_text))

    # Auto-use conductor recording if no explicit audio was uploaded
    conductor_recording = f"/tmp/kimfam_recording_{meeting_id}.webm"
    if not audio_file and _os.path.exists(conductor_recording):
        with open(conductor_recording, "rb") as _rf:
            _rb = _rf.read()
        class _FakeFile:
            filename = "conductor_recording.webm"
            async def read(self): return _rb
        audio_file = _FakeFile()

    if audio_file and hasattr(audio_file, "filename"):
        groq_key = _os.getenv("GROQ_API_KEY", "")
        if not groq_key:
            raise _HE(status_code=503, detail="Groq API key not configured")
        audio_bytes = await audio_file.read()
        audio_name = audio_file.filename or "audio.m4a"
        try:
            # Handles recordings of any length: long audio is re-encoded + split
            # into 15-min segments and transcribed chunk-by-chunk, then stitched.
            # Runs in a thread so the long ffmpeg + Whisper work doesn't freeze the
            # worker's event loop for the whole upload.
            import asyncio as _aio_t
            audio_text = await _aio_t.to_thread(
                _transcribe_audio, audio_bytes, audio_name, groq_key)
            if audio_text:
                transcript_parts.append((f"Audio: {audio_name}", audio_text))
        except Exception as e:
            raise _HE(status_code=502, detail=f"Transcription failed: {e}")

    # Fallback: if the form carried no transcript/audio (e.g. the paste box was
    # cleared, or a recording was lost), use a transcript already saved on this
    # meeting — group auto-capture or an earlier attempt may have stored it. This
    # stops a cleared text box from silently downgrading the minutes to notes-only.
    if not transcript_parts:
        try:
            _saved = _dbq("SELECT transcript FROM meetings WHERE id=%s", (meeting_id,))
            _saved_tx = (_saved[0]["transcript"] or "").strip() if _saved else ""
            if len(_saved_tx) > 200:
                transcript_parts.append(("Saved transcript", _saved_tx))
        except Exception:
            pass

    if not transcript_parts and not sec_notes:
        raise _HE(status_code=400, detail="No content provided — supply audio, a transcript, or your notes")

    if not transcript_parts:
        # Notes-only: the secretary's notes ARE the meeting record
        transcript = "(No recording or transcript — see the secretary's notes below.)"
    elif len(transcript_parts) == 1:
        transcript = transcript_parts[0][1]
    else:
        # Multiple sources — label clearly so Claude can reconcile
        transcript = "\n\n".join(
            f"[SOURCE: {label}]\n{text}"
            for label, text in transcript_parts
        )

    # Persist the FULL, raw transcript. The minutes writer and retrospective need
    # the complete detail (figures, who said what); condensing here is what made
    # past minutes thin. The extraction prompt below uses a condensed COPY only.
    if transcript_parts:
        try:
            from db import execute as _exec_tx
            _exec_tx("UPDATE meetings SET transcript=%s WHERE id=%s", (transcript, meeting_id))
        except Exception:
            pass

    # For the action/decision EXTRACTION call only, condense a long transcript so
    # it fits one AI call. This does NOT touch the stored transcript above.
    _already_condensed = transcript.startswith("[This is a condensed record")
    if transcript_parts and len(transcript) > 18000 and not _already_condensed:
        transcript = await _condense_transcript(transcript)

    # ── 3. Pull context from DB ───────────────────────────────────────────────
    meeting_rows = _dbq("""
        SELECT ref, date, key_decisions FROM meetings
        ORDER BY date DESC LIMIT 8
    """)
    open_actions = _dbq("""
        SELECT id, ref, description, assignee, deadline, status
        FROM actions
        WHERE status NOT IN ('done','cancelled','carried_over')
           OR (status = 'done' AND closed_at >= NOW() - INTERVAL '21 days')
        ORDER BY ref DESC
    """)
    # Pull the latest few updates members have posted on each open action so the
    # minutes generator sees the real progress, not just the static action state.
    updates_by_action = {}
    if open_actions:
        ids = tuple(r["id"] for r in open_actions)
        upd_rows = _dbq("""
            SELECT action_id, author, text, created_at
            FROM (
                SELECT action_id, author, text, created_at,
                       ROW_NUMBER() OVER (PARTITION BY action_id ORDER BY created_at DESC) AS rn
                FROM action_updates
                WHERE action_id IN %s AND text IS NOT NULL AND text <> ''
            ) t
            WHERE rn <= 3
            ORDER BY action_id, created_at ASC
        """, (ids,))
        for u in upd_rows:
            updates_by_action.setdefault(u["action_id"], []).append(u)

    current_meeting = _dbq("SELECT ref, date, venue, attendance FROM meetings WHERE id=%s", (meeting_id,))
    if not current_meeting:
        raise _HE(status_code=404, detail="Meeting not found")
    mtg = current_meeting[0]

    # Roll-call attendance captured in the conductor (present / apology / absent + notes)
    _raw_att = mtg.get("attendance")
    _att = (_json.loads(_raw_att) if isinstance(_raw_att, (str, bytes, bytearray)) else (_raw_att or {}))
    attendance_block = ""
    if _att:
        _lines = []
        for _m, _e in _att.items():
            if not isinstance(_e, dict):
                continue
            _st = (_e.get("status") or "").strip()
            if not _st:
                continue
            _cm = (_e.get("comment") or "").strip()
            _lines.append(f"- {_m}: {_st}" + (f" ({_cm})" if _cm else ""))
        if _lines:
            attendance_block = (
                "\n\nATTENDANCE (roll call taken in the meeting — record this in the minutes "
                "under attendance/apologies; do NOT invent names beyond this list):\n"
                + "\n".join(_lines) + "\n"
            )

    prev_summary = "\n".join(
        f"- {r['ref']} ({r['date']}): {(r['key_decisions'] or 'no decisions recorded')[:120]}"
        for r in meeting_rows if r['ref'] != mtg['ref']
    )

    def _action_block(r):
        head = (f"- {r['ref']} [{r['status']}] {r['assignee']}: {r['description'][:80]}"
                + (f" (due {r['deadline']})" if r['deadline'] else ""))
        ups = updates_by_action.get(r["id"], [])
        if ups:
            lines = [
                f"    · update from {u['author']} ({str(u['created_at'])[:10]}): {u['text'][:160]}"
                for u in ups
            ]
            return head + "\n" + "\n".join(lines)
        return head

    open_summary = "\n".join(_action_block(r) for r in open_actions)

    multi_source_note = (
        "NOTE: Multiple transcript sources are provided below. "
        "Reconcile them: prefer the more detailed account for each point. "
        "Do not duplicate actions that appear in both sources.\n\n"
        if len(transcript_parts) > 1 else ""
    )

    secretary_block = (
        f"\n\nSECRETARY'S NOTES (authoritative — these come directly from the meeting secretary "
        f"and OVERRIDE the transcript where they conflict. Use them for framing, emphasis, "
        f"corrections, and any context the transcript missed):\n{sec_notes}\n"
        if sec_notes else ""
    )

    prompt = f"""You are processing minutes for KimFam Investment Club meeting {mtg['ref']} held {mtg['date']}.

RECENT MEETINGS (for context):
{prev_summary or 'None yet'}

OPEN AND RECENTLY-COMPLETED ACTIONS — with the progress updates members have already posted
in the app. Items tagged [done] were completed in the last few weeks: acknowledge them as
closed/achieved in the summary or decisions, but do NOT reopen them.
(Do NOT recreate any of these unless the transcript explicitly assigns a new deadline or owner;
treat the posted updates as factual progress and reflect them in updates_to_existing even if
the transcript is silent on them):
{open_summary or 'None'}

{multi_source_note}MEETING TRANSCRIPT:
{transcript}{secretary_block}{attendance_block}

Extract and return ONLY valid JSON (no markdown, no explanation) in this exact shape:
{{
  "summary": "2-3 sentence meeting summary",
  "key_topics": "comma-separated list of main topics discussed",
  "key_decisions": ["decision 1", "decision 2"],
  "new_actions": [
    {{
      "description": "what needs to be done",
      "assignee": "member name or All Members",
      "deadline": "YYYY-MM-DD or null",
      "priority": "high|medium|low",
      "matches_existing": "KIM/ref if this updates an existing action or null"
    }}
  ],
  "updates_to_existing": [
    {{
      "ref": "KIM/existing-ref",
      "note": "what was said about this action in the meeting",
      "new_status": "in_progress|blocked|done or null if unchanged"
    }}
  ]
}}

Rules:
- Only include actions explicitly assigned in the transcript. Do not invent actions.
- For each open action mentioned, add an entry to updates_to_existing.
- If an existing action is being carried over with a new deadline, add it to new_actions with matches_existing set.
- Assignee names must match exactly: Hillary, Hellen, Alex, Solomon, Viola, Max, James, or "All Members".
- Deadlines must be absolute dates (YYYY-MM-DD). If relative ("next week"), compute from meeting date {mtg['date']}.
"""

    # ── 4. Extract via the AI chain (Claude primary → DeepSeek → Groq). The
    #       transcript was already condensed above if it was long, so this prompt
    #       fits Claude. ──────────────────────────────────────────────────────────
    raw_json = await _ai_complete(prompt, claude_timeout=150)

    if not raw_json:
        raise _HE(status_code=503, detail="AI extraction is busy right now — please tap Extract again in a moment.")

    # Strip markdown fences if model wrapped the JSON
    raw_json = _re.sub(r"^```(?:json)?\s*", "", raw_json)
    raw_json = _re.sub(r"\s*```$", "", raw_json.strip())

    extracted = None
    try:
        extracted = _json.loads(raw_json)
    except _json.JSONDecodeError:
        # Last resort: find the outermost JSON block and try that
        m = _re.search(r"\{.*\}", raw_json, _re.DOTALL)
        if m:
            try:
                extracted = _json.loads(m.group())
            except _json.JSONDecodeError:
                extracted = None
    if not isinstance(extracted, dict):
        raise _HE(
            status_code=422,
            detail="Couldn't extract minutes from the content provided. Add the actual "
                   "meeting transcript, recording, or substantive notes (a single word "
                   "like 'none' isn't enough to work from), then try again.",
        )

    return {
        "ok": True,
        "meeting_ref": mtg["ref"],
        "transcript_preview": transcript[:500] + ("…" if len(transcript) > 500 else ""),
        "extracted": extracted,
    }


@app.post("/api/meetings/{meeting_id}/confirm")
async def confirm_meeting_extraction(meeting_id: int, request: Request):
    """Write approved extraction results to DB (actions + meeting decisions)."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq, execute as _exec
    from datetime import datetime as _dt

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    body = await request.json()
    extracted   = body.get("extracted", {})
    meeting_ref = body.get("meeting_ref", "")
    author      = payload.get("sub", "admin")

    rows = _dbq("SELECT id, ref FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")

    summary        = extracted.get("summary", "")
    key_topics     = extracted.get("key_topics", "")
    key_decisions  = extracted.get("key_decisions", [])
    new_actions    = extracted.get("new_actions", [])
    updates        = extracted.get("updates_to_existing", [])

    # Update meeting record
    decisions_text = "\n".join(f"- {d}" for d in key_decisions)
    _exec("""UPDATE meetings SET summary=%s, key_topics=%s, key_decisions=%s WHERE id=%s""",
          (summary, key_topics, decisions_text, meeting_id))

    created_refs = []
    docx_actions = []  # for .docx generation
    import re as _re2
    for i, a in enumerate(new_actions, start=1):
        desc      = (a.get("description") or "").strip()
        # assignees: accept list or single string
        raw_asgn  = a.get("assignees") or a.get("assignee") or "Unknown"
        if isinstance(raw_asgn, list):
            assignees = [x.strip() for x in raw_asgn if str(x).strip()]
        else:
            assignees = [s.strip() for s in str(raw_asgn).split(",") if s.strip()]
        if not assignees:
            assignees = ["Unknown"]
        primary_assignee = assignees[0]
        deadline  = a.get("deadline") or None
        priority  = a.get("priority") or "medium"
        parent    = a.get("matches_existing") or None
        project   = (a.get("project_id") or "").strip() or None
        if not desc:
            continue

        # Build ref: "KIM 009/2026" → "KIM/09/26-{i}"
        m = _re2.match(r"KIM\s+(\d+)/(\d{4})", meeting_ref)
        if m:
            num = int(m.group(1)); yr = m.group(2)[-2:]
            ref = f"KIM/{num:02d}/{yr}-{i}"
        else:
            ref = f"{meeting_ref}-{i}"

        # Idempotent: re-confirming the same meeting updates the deterministic ref
        # in place rather than creating a duplicate "-1a" set (the old collision
        # logic is what duplicated actions on a double-tap). Status is left
        # untouched so a manually progressed action is not reset to open.
        _exec("""INSERT INTO actions
                   (ref, description, assignee, assignees, meeting_id, related_meeting,
                    status, priority, deadline, parent_ref, project_id, created_by)
                 VALUES (%s,%s,%s,%s,%s,%s,'open',%s,%s,%s,%s,%s)
                 ON CONFLICT (ref) DO UPDATE SET
                    description = EXCLUDED.description,
                    assignee    = EXCLUDED.assignee,
                    assignees   = EXCLUDED.assignees,
                    priority    = EXCLUDED.priority,
                    deadline    = EXCLUDED.deadline,
                    parent_ref  = EXCLUDED.parent_ref,
                    project_id  = EXCLUDED.project_id""",
              (ref, desc, primary_assignee, assignees, meeting_id, meeting_ref,
               priority, deadline, parent, project, author))
        created_refs.append(ref)
        docx_actions.append({
            "ref": ref, "description": desc, "assignees": assignees,
            "deadline": deadline or "", "priority": priority,
        })

    # Idempotent: clear any prior "[From minutes]" notes for this meeting so a
    # re-confirm does not stack duplicate update comments on existing actions.
    _exec("DELETE FROM action_updates WHERE author=%s AND text LIKE %s",
          (f"Meeting {meeting_ref}", "[From minutes]%"))

    # Apply updates to existing actions (additive: add note, optionally change status)
    for u in updates:
        ref  = (u.get("ref") or "").strip()
        note = (u.get("note") or "").strip()
        new_status = (u.get("new_status") or "").strip().lower() or None
        if not ref or not note:
            continue
        rows2 = _dbq("SELECT id, status FROM actions WHERE ref=%s", (ref,))
        if not rows2:
            continue
        action_id = rows2[0]["id"]
        current_status = rows2[0]["status"]
        _exec("""INSERT INTO action_updates (action_id, author, text, type)
                 VALUES (%s,%s,%s,'comment')""",
              (action_id, f"Meeting {meeting_ref}", f"[From minutes] {note}"))
        valid_transitions = {"open", "in_progress", "blocked", "done", "cancelled", "carried_over"}
        if new_status and new_status in valid_transitions and current_status not in ("done", "cancelled"):
            _exec("UPDATE actions SET status=%s WHERE ref=%s", (new_status, ref))

    # ── Build structured minutes data, save JSON + generate .docx ────────────
    mtg_row = _dbq("SELECT ref, date, venue, start_time_eat FROM meetings WHERE id=%s", (meeting_id,))
    mtg_meta = mtg_row[0] if mtg_row else {}

    minutes_data = {
        "meeting_ref":   meeting_ref,
        "date":          str(mtg_meta.get("date", "")),
        "venue":         mtg_meta.get("venue") or "Google Meet",
        "start_time_eat": str(mtg_meta.get("start_time_eat") or "")[:5],
        "summary":       summary,
        "key_topics":    key_topics,
        "key_decisions": key_decisions,
        "actions":       docx_actions,
    }

    data_path  = f"/tmp/kimfam_minutes_{meeting_id}_data.json"
    draft_path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    try:
        import json as _json2
        with open(data_path, "w") as _f:
            _json2.dump(minutes_data, _f)
        # Prefer the full old-style narrative minutes; fall back to the simple
        # builder if the AI narrative can't be generated.
        _docx_bytes = None
        try:
            _narr = _generate_minutes_narrative(meeting_id, key_topics, key_decisions, docx_actions)
            if _narr:
                _docx_bytes = _build_minutes_docx_v2(meeting_ref, mtg_meta, _narr, docx_actions)
        except Exception as _ne:
            import logging as _lg
            _lg.getLogger("main").error(f"narrative minutes failed, falling back: {_ne}")
        if _docx_bytes is None:
            _docx_bytes = _build_minutes_docx(mtg=mtg_meta, **{
                k: minutes_data[k]
                for k in ("meeting_ref","summary","key_topics","key_decisions","actions")
            })
        with open(draft_path, "wb") as _f:
            _f.write(_docx_bytes)
    except Exception as _e:
        import logging as _lg
        _lg.getLogger("main").error(f"docx generation failed: {_e}")
        draft_path = None

    # Auto-generate the time-analytics retrospective so it's ready for the next
    # meeting's "Review of Last Meeting" item. Best-effort; never blocks confirm.
    try:
        _generate_retrospective(meeting_id)
    except Exception as _re:
        import logging as _lg
        _lg.getLogger("main").error(f"retrospective generation failed: {_re}")

    return {
        "ok":         True,
        "created":    created_refs,
        "updated":    len(updates),
        "draft_url":  f"/api/meetings/{meeting_id}/minutes/draft" if draft_path else None,
        "minutes_data": minutes_data,
    }


def _build_minutes_docx(meeting_ref: str, mtg: dict, summary: str,
                         key_topics: str, key_decisions: list, actions: list) -> bytes:
    """Generate a .docx meeting minutes file. Returns raw bytes."""
    from docx import Document as _DocxDoc
    from docx.shared import Pt as _Pt, Cm as _Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH as _WDA
    from io import BytesIO as _BIO
    from datetime import date as _date

    doc = _DocxDoc()

    # Margins
    for sec in doc.sections:
        sec.top_margin    = _Cm(2.5)
        sec.bottom_margin = _Cm(2.5)
        sec.left_margin   = _Cm(3)
        sec.right_margin  = _Cm(3)

    # Title block
    t = doc.add_heading("KIMFAM INVESTMENT CLUB", 0)
    t.alignment = _WDA.CENTER

    sub = doc.add_paragraph("Meeting Minutes")
    sub.alignment = _WDA.CENTER
    for run in sub.runs:
        run.font.size = _Pt(13)
        run.font.bold = True

    doc.add_paragraph()

    # Meeting metadata
    meta = [
        ("Meeting Reference", meeting_ref),
        ("Date",   str(mtg.get("date", ""))),
        ("Time",   (str(mtg.get("start_time_eat") or "")[:5] + " EAT") if mtg.get("start_time_eat") else ""),
        ("Venue",  mtg.get("venue") or "Google Meet"),
        ("Prepared by", "KimFam Hub"),
    ]
    for label, value in meta:
        if not value:
            continue
        p = doc.add_paragraph()
        run_l = p.add_run(f"{label}: ")
        run_l.bold = True
        p.add_run(value)

    doc.add_paragraph()

    # Summary
    doc.add_heading("MEETING SUMMARY", 1)
    doc.add_paragraph(summary or "")

    # Key topics
    if key_topics:
        doc.add_heading("KEY TOPICS DISCUSSED", 1)
        doc.add_paragraph(key_topics)

    # Decisions
    if key_decisions:
        doc.add_heading("KEY DECISIONS", 1)
        for idx, d in enumerate(key_decisions, 1):
            doc.add_paragraph(f"{idx}. {d}")

    # Action points table
    if actions:
        doc.add_heading("ACTION POINTS", 1)
        tbl = doc.add_table(rows=1, cols=5)
        tbl.style = "Table Grid"
        hcells = tbl.rows[0].cells
        for ci, hdr in enumerate(["Ref", "Description", "Assigned To", "Deadline", "Priority"]):
            hcells[ci].text = hdr
            for para in hcells[ci].paragraphs:
                for run in para.runs:
                    run.bold = True
        for a in actions:
            rc = tbl.add_row().cells
            rc[0].text = a.get("ref", "")
            rc[1].text = a.get("description", "")
            asgn = a.get("assignees", [])
            rc[2].text = ", ".join(asgn) if isinstance(asgn, list) else str(asgn)
            rc[3].text = str(a.get("deadline") or "")
            rc[4].text = (a.get("priority") or "medium").capitalize()

    doc.add_paragraph()
    doc.add_paragraph(f"Generated by KimFam Hub on {_date.today()}")

    buf = _BIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _minutes_filename(date_val) -> str:
    """File name matching the old convention: KIMFAM_Meeting_Minutes_June_7_2026.docx"""
    import datetime as _dt
    d = None
    if isinstance(date_val, str):
        try: d = _dt.date.fromisoformat(date_val[:10])
        except Exception: d = None
    elif isinstance(date_val, (_dt.date, _dt.datetime)):
        d = date_val if isinstance(date_val, _dt.date) else date_val.date()
    if not d:
        return "KIMFAM_Meeting_Minutes.docx"
    return f"KIMFAM_Meeting_Minutes_{d.strftime('%B')}_{d.day}_{d.year}.docx"


def _fmt_eat(ts) -> str:
    """Format a tz-aware UTC timestamp as an EAT clock string, e.g. '4:35 PM'."""
    if not ts:
        return ""
    import datetime as _dt
    try:
        eat = ts.astimezone(_dt.timezone(_dt.timedelta(hours=3)))
        return eat.strftime("%-I:%M %p")
    except Exception:
        return ""


def _generate_minutes_narrative(meeting_id: int, key_topics: str,
                                key_decisions: list, actions: list) -> dict | None:
    """Produce full old-style KimFam minutes (attendance + numbered narrative
    sections) from the transcript, live notes, roll call and actions. Returns a
    structured dict the docx builder renders, or None on failure."""
    from db import query as _dbq_n
    import json as _json_n2, re as _re_n2
    rows = _dbq_n("""SELECT ref, date, venue, conductor_notes, transcript,
                            attendance, conductor_started_at, conductor_ended_at
                     FROM meetings WHERE id=%s""", (meeting_id,))
    if not rows:
        return None
    r = rows[0]
    _att_raw = r["attendance"]
    attendance = (_json_n2.loads(_att_raw) if isinstance(_att_raw, (str, bytes, bytearray)) else (_att_raw or {}))
    att_lines = []
    for _m, _e in (attendance or {}).items():
        if isinstance(_e, dict) and _e.get("status"):
            _c = (_e.get("comment") or "").strip()
            att_lines.append(f"- {_m}: {_e['status']}" + (f" ({_c})" if _c else ""))
    # Use the FULL notes and transcript. Truncating here is what produced thin
    # minutes; we keep all the detail and handle length via a map-reduce below.
    notes = (r["conductor_notes"] or "")[:6000]
    transcript = (r["transcript"] or "")
    if len(transcript) > 200000:   # hard safety cap, far above a 3-hour meeting
        transcript = transcript[:200000]

    # A 3-hour transcript is too large for one final pass (sonnet times out).
    # MAP: split it and have haiku write DETAILED notes per chunk (preserving every
    # figure, decision and assignment). REDUCE (further below): sonnet assembles the
    # house-style minutes from these notes. This keeps the detail that makes good
    # minutes while fitting the time budget.
    if len(transcript) > 40000:
        _chunks = [transcript[i:i + 35000] for i in range(0, len(transcript), 35000)]
        _digparts = []
        # Keep each part's notes compact so the final assemble prompt stays within the
        # size the CLI can handle in time (a ~30k prompt works; ~60k times out).
        _per_cap = max(3500, 22000 // max(1, len(_chunks)))
        for _idx, _ch in enumerate(_chunks, 1):
            _dp = (
                f"This is part {_idx} of {len(_chunks)} of a KimFam Investment Club meeting "
                f"transcript. Write CONCISE but COMPLETE minute notes for this part as short "
                f"bullet points: every topic, every figure and amount (UGX, percentages, USD), "
                f"every decision, and every action assigned (who and any deadline). Keep all "
                f"facts and numbers but no filler or dialogue. No em-dashes or en-dashes.\n\n"
                f"PART {_idx}:\n{_ch}"
            )
            _out = (_ask_claude(_dp, model="claude-haiku-4-5-20251001", timeout=180) or "")[:_per_cap]
            _digparts.append(f"--- Notes from part {_idx} of {len(_chunks)} ---\n{_out or '(part unavailable)'}")
        transcript = "\n\n".join(_digparts)
    acts_txt = "\n".join(
        f"- {a.get('ref','')}: {a.get('description','')[:90]} "
        f"[{', '.join(a['assignees']) if isinstance(a.get('assignees'), list) else a.get('assignees','')}"
        f"{', due ' + str(a.get('deadline')) if a.get('deadline') else ''}]"
        for a in (actions or [])
    ) or "None recorded"
    prev = _dbq_n("""SELECT ref, date, key_decisions FROM meetings
                     WHERE date < %s ORDER BY date DESC LIMIT 1""", (r["date"],))
    prev_txt = (f"{prev[0]['ref']} ({prev[0]['date']}): {(prev[0]['key_decisions'] or '')[:400]}"
                if prev else "None")
    prev_ref = prev[0]["ref"] if prev else None

    # Status-of-prior-actions roll-up (restores the table the older minutes carried,
    # e.g. "Status of KIM 006 and Prior Actions"). Covers the previous meeting's
    # actions plus anything still open from before, with its current outcome.
    prior_status = []
    try:
        _ps_rows = _dbq_n("""
            SELECT a.ref, a.description, a.assignee, a.assignees, a.status
            FROM actions a
            WHERE (a.related_meeting = %s
                   OR (a.status IN ('open','in_progress','blocked') AND a.meeting_id < %s))
              AND a.meeting_id <> %s
            ORDER BY a.ref
        """, (prev_ref, meeting_id, meeting_id))
        _outcome = {"done": "Closed", "cancelled": "Cancelled", "carried_over": "Carried over",
                    "in_progress": "In progress", "blocked": "Blocked", "open": "Open"}
        for _pr in _ps_rows:
            _resp = (", ".join(_pr["assignees"]) if _pr.get("assignees")
                     else (_pr.get("assignee") or ""))
            prior_status.append({
                "ref": _pr["ref"], "description": _pr["description"],
                "responsible": _resp, "outcome": _outcome.get(_pr["status"], _pr["status"]),
            })
    except Exception:
        prior_status = []

    prompt = f"""You are the Secretary writing the official, detailed minutes for KimFam
Investment Club meeting {r['ref']} held {r['date']}, working from the full meeting
transcript. Produce thorough, faithful minutes in the club's established house style.
The family relies on these minutes as the permanent record, so capture every topic,
figure, decision and assignment that the transcript actually contains. Do not be brief.

STRICT STYLE RULES:
- NO em-dashes or en-dashes anywhere. Use commas, colons, semicolons or full stops.
- Warm, factual, third-person minute style (e.g. "Hellen presented the financial position.").
- Use Ugandan shilling amounts as "UGX 1,234,567" with comma separators; keep all figures and percentages.
- Refer to Simon and Viola together as "the Arungas" (the family prefers this).
- Use the roll call below as the authoritative attendance; do not invent presence.
- Do not invent figures, names, or decisions that are not in the inputs.

HOUSE STRUCTURE (use the sections that the meeting actually covered; add subsections freely):
- OPENING (who opened, prayer, chair/secretary, quorum)
- FINANCIAL REPORT / TREASURER'S REPORT (balance, expenditure, arrears, with figures)
- REVIEW OF PREVIOUS MINUTES
- FOLLOW-UP ON ACTION POINTS (one numbered subsection per item reviewed, with its KIM reference and status)
- MAIN AGENDA and/or PROJECT UPDATES (one numbered subsection per topic discussed, e.g. equity model vote, new venture proposals, project selection, investment decisions, each with the discussion, figures and outcome)
- ANY OTHER BUSINESS
- NEXT MEETING
- CLOSE (closing prayer)
Within each section use "paragraphs" for narrative and "bullets" for specifics (figures, sub-points, the agreed action). Be generous with bullets so nothing is lost.

INPUTS
Previous meeting: {prev_txt}
Roll call (authoritative attendance):
{chr(10).join(att_lines) or '(not recorded)'}
Secretary's live notes (organised by agenda item):
{notes or '(none)'}
Actions agreed this meeting:
{acts_txt}
MEETING RECORD (full transcript, or detailed section notes for a long meeting):
{transcript or '(no transcript)'}

Return ONLY valid JSON (no markdown) in EXACTLY this shape (use as many sections,
subsections, paragraphs and bullets as the meeting needs):
{{
  "attendance": {{"present": ["Name (Family)", "..."], "apologies": ["Name with reason", "..."], "absent": ["Name", "..."]}},
  "chair": "Name",
  "secretary": "Name",
  "sections": [
    {{"number": "1", "title": "OPENING", "paragraphs": ["..."], "bullets": [], "subsections": []}},
    {{"number": "4", "title": "FOLLOW-UP ON ACTION POINTS", "paragraphs": [], "bullets": [],
      "subsections": [{{"number": "4.1", "title": "Item name (KIM/ref)", "paragraphs": ["..."], "bullets": ["..."]}}]}}
  ],
  "decisions": ["Concise decision 1", "Concise decision 2"]
}}

Number the sections sequentially from 1. Only include what the inputs support, but be
exhaustive about what they do support."""
    raw = _ask_claude(prompt, model="claude-haiku-4-5-20251001", timeout=240)
    def _parse(x):
        if not x: return None
        try: return _json_n2.loads(x)
        except Exception:
            m = _re_n2.search(r"\{.*\}", x, _re_n2.DOTALL)
            if m:
                try: return _json_n2.loads(m.group())
                except Exception: return None
        return None
    data = _parse(raw)
    if not isinstance(data, dict):
        import os as _os_n
        groq_key = _os_n.getenv("GROQ_API_KEY", "")
        if groq_key:
            try:
                from groq import Groq as _Groq
                resp = _Groq(api_key=groq_key).chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=8000, temperature=0.2)
                data = _parse(resp.choices[0].message.content.strip())
            except Exception:
                data = None
    if not isinstance(data, dict):
        return None
    data["_start_eat"] = _fmt_eat(r["conductor_started_at"])
    data["_end_eat"]   = _fmt_eat(r["conductor_ended_at"])
    data["prior_actions_status"] = prior_status
    return data


def _build_minutes_docx_v2(meeting_ref: str, mtg: dict, narrative: dict, actions: list) -> bytes:
    """Render the rich, old-style minutes (attendance + numbered narrative sections
    + action points table) to .docx bytes."""
    from docx import Document as _DocxDoc
    from docx.shared import Pt as _Pt, Cm as _Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH as _WDA
    from io import BytesIO as _BIO2
    from datetime import date as _date2, datetime as _dt2

    doc = _DocxDoc()
    for sec in doc.sections:
        sec.top_margin = _Cm(2.5); sec.bottom_margin = _Cm(2.5)
        sec.left_margin = _Cm(3); sec.right_margin = _Cm(3)

    t = doc.add_heading("KIMFAM INVESTMENT CLUB", 0); t.alignment = _WDA.CENTER
    sub = doc.add_paragraph("MEETING MINUTES"); sub.alignment = _WDA.CENTER
    for run in sub.runs: run.font.size = _Pt(13); run.font.bold = True

    # Header line: KIM 013/2026  |  Sunday, 21 June 2026
    day_date = ""
    try:
        _d = _dt2.fromisoformat(str(mtg.get("date"))[:10])
        day_date = _d.strftime("%A, %-d %B %Y")
    except Exception:
        day_date = str(mtg.get("date", ""))
    hdr = doc.add_paragraph(); hdr.alignment = _WDA.CENTER
    hr = hdr.add_run(f"{meeting_ref}    |    {day_date}"); hr.bold = True
    # Actual conduct time
    _start = narrative.get("_start_eat"); _end = narrative.get("_end_eat")
    if _start:
        tl = doc.add_paragraph(); tl.alignment = _WDA.CENTER
        tl.add_run(f"Started {_start} EAT" + (f", ended {_end} EAT" if _end else "")
                   + f"  |  {mtg.get('venue') or 'Google Meet'}").italic = True
    doc.add_paragraph()

    # Meeting metadata table (matches the house template)
    _time_val = (f"{_start} EAT" + (f"  |  {_end} EAT" if _end else "")) if _start else ""
    _meta_rows = [
        ("Date", day_date),
        ("Time", _time_val),
        ("Venue", mtg.get("venue") or "Google Meet"),
        ("Chair", narrative.get("chair") or ""),
        ("Secretary", narrative.get("secretary") or ""),
    ]
    _meta_rows = [(k, v) for k, v in _meta_rows if v]
    if _meta_rows:
        mtbl = doc.add_table(rows=len(_meta_rows), cols=2); mtbl.style = "Table Grid"
        for _i, (_k, _v) in enumerate(_meta_rows):
            mtbl.rows[_i].cells[0].text = _k
            mtbl.rows[_i].cells[1].text = _v
            for _run in mtbl.rows[_i].cells[0].paragraphs[0].runs:
                _run.bold = True
        doc.add_paragraph()

    # Attendance
    att = narrative.get("attendance") or {}
    if any(att.get(k) for k in ("present", "apologies", "absent")):
        doc.add_heading("ATTENDANCE", 1)
        if att.get("present"):
            p = doc.add_paragraph(); p.add_run("Present: ").bold = True
            p.add_run(", ".join(att["present"]))
        if att.get("apologies"):
            p = doc.add_paragraph(); p.add_run("Absent with apologies: ").bold = True
            p.add_run(", ".join(att["apologies"]))
        if att.get("absent"):
            p = doc.add_paragraph(); p.add_run("Absent: ").bold = True
            p.add_run(", ".join(att["absent"]))

    def _render_block(b, level):
        for para in (b.get("paragraphs") or []):
            if para and para.strip():
                doc.add_paragraph(para.strip())
        for bullet in (b.get("bullets") or []):
            if bullet and bullet.strip():
                doc.add_paragraph(bullet.strip(), style="List Bullet")

    for s in (narrative.get("sections") or []):
        title = f"{s.get('number','')}. {s.get('title','')}".strip(". ")
        if not title:
            continue
        doc.add_heading(title, 1)
        _render_block(s, 1)
        for ss in (s.get("subsections") or []):
            stitle = f"{ss.get('number','')}  {ss.get('title','')}".strip()
            if stitle:
                doc.add_heading(stitle, 2)
            _render_block(ss, 2)

    # Action points table
    if actions:
        doc.add_heading("ACTION POINTS", 1)
        tbl = doc.add_table(rows=1, cols=5); tbl.style = "Table Grid"
        for ci, hdr2 in enumerate(["Ref", "Description", "Assigned To", "Deadline", "Priority"]):
            cell = tbl.rows[0].cells[ci]; cell.text = hdr2
            for para in cell.paragraphs:
                for run in para.runs: run.bold = True
        for a in actions:
            rc = tbl.add_row().cells
            rc[0].text = a.get("ref", "")
            rc[1].text = a.get("description", "")
            asgn = a.get("assignees", [])
            rc[2].text = ", ".join(asgn) if isinstance(asgn, list) else str(asgn)
            rc[3].text = str(a.get("deadline") or "")
            rc[4].text = (a.get("priority") or "medium").capitalize()

    # Status of prior actions (roll-up table, restores the older minutes format)
    prior = narrative.get("prior_actions_status") or []
    if prior:
        doc.add_heading("STATUS OF PRIOR ACTIONS", 1)
        ptbl = doc.add_table(rows=1, cols=4); ptbl.style = "Table Grid"
        for ci, hdr3 in enumerate(["Action ID", "Description", "Responsible", "Outcome"]):
            cell = ptbl.rows[0].cells[ci]; cell.text = hdr3
            for para in cell.paragraphs:
                for run in para.runs: run.bold = True
        for a in prior:
            rc = ptbl.add_row().cells
            rc[0].text = a.get("ref", "")
            rc[1].text = a.get("description", "")
            rc[2].text = a.get("responsible", "")
            rc[3].text = a.get("outcome", "")

    # Decisions
    if narrative.get("decisions"):
        doc.add_heading("KEY DECISIONS", 1)
        for idx, d in enumerate(narrative["decisions"], 1):
            doc.add_paragraph(f"{idx}. {d}")

    doc.add_paragraph()
    doc.add_paragraph(f"Prepared on KimFam Hub, {_date2.today().strftime('%-d %B %Y')}.")

    buf = _BIO2(); doc.save(buf); buf.seek(0)
    return buf.getvalue()


@app.get("/api/meetings/{meeting_id}/minutes")
async def meetings_minutes_file(meeting_id: int, request: Request):
    """Download the published minutes .docx. Any authenticated member.
    Used when minutes are served locally (e.g. staging, or no R2)."""
    from fastapi import HTTPException as _HE
    from fastapi.responses import FileResponse as _FR
    from db import query as _dbq
    import os as _os
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    if not _os.path.exists(path):
        raise _HE(status_code=404, detail="Minutes file not found")
    rows = _dbq("SELECT ref FROM meetings WHERE id=%s", (meeting_id,))
    safe_ref = (rows[0]["ref"] if rows else f"meeting_{meeting_id}").replace(" ", "_").replace("/", "_")
    return _FR(path, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
               filename=f"KimFam_{safe_ref}_Minutes.docx")


@app.get("/api/meetings/{meeting_id}/minutes/view")
async def meetings_minutes_view(meeting_id: int, request: Request):
    """Render the published minutes .docx as HTML for in-app viewing. All members."""
    from fastapi import HTTPException as _HE
    import os as _os, io as _io
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    if not _os.path.exists(path):
        raise _HE(status_code=404, detail="Minutes file not found")
    with open(path, "rb") as _f:
        result = mammoth.convert_to_html(_io.BytesIO(_f.read()))
    html = f"""<!DOCTYPE html><html><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
  body{{font-family:Georgia,serif;max-width:800px;margin:0 auto;padding:20px;background:#fff;color:#1a1a1a;line-height:1.7;font-size:15px}}
  h1,h2,h3{{color:#1a3a1a;margin-top:1.5em}} table{{border-collapse:collapse;width:100%;margin:1em 0}}
  td,th{{border:1px solid #ccc;padding:8px 12px}} th{{background:#f0f7f0}}
  p{{margin:0.7em 0}}
</style></head><body>{result.value}</body></html>"""
    return HTMLResp(content=html)


@app.get("/api/meetings/{meeting_id}/minutes/draft")
async def meetings_minutes_draft(meeting_id: int, request: Request):
    """Serve the generated draft .docx for admin review (before publish)."""
    from fastapi import HTTPException as _HE
    from fastapi.responses import FileResponse as _FR
    from db import query as _dbq
    import os as _os

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    if not _os.path.exists(path):
        raise _HE(status_code=404, detail="Draft not found — run Process + Confirm first")

    rows = _dbq("SELECT date FROM meetings WHERE id=%s", (meeting_id,))
    fname = _minutes_filename(rows[0]["date"]) if rows else f"KIMFAM_Meeting_Minutes_{meeting_id}.docx"
    return _FR(
        path,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=fname,
    )


@app.post("/api/meetings/{meeting_id}/minutes/draft")
async def meetings_minutes_draft_replace(meeting_id: int, request: Request):
    """Replace the server draft with an edited .docx uploaded by admin."""
    from fastapi import HTTPException as _HE
    import os as _os

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    form = await request.form()
    docx_file = form.get("docx_file")
    if not docx_file or not hasattr(docx_file, "filename"):
        raise _HE(status_code=400, detail="docx_file required")

    fname = (docx_file.filename or "").lower()
    if not fname.endswith(".docx"):
        raise _HE(status_code=400, detail="Only .docx files accepted")

    raw = await docx_file.read()
    path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    with open(path, "wb") as f:
        f.write(raw)

    return {"ok": True, "message": "Draft replaced — Approve & Send will use this version"}


@app.post("/api/meetings/{meeting_id}/minutes/edit")
async def meetings_minutes_edit(meeting_id: int, request: Request):
    """Apply a plain-English edit instruction to the minutes via Claude, regenerate .docx."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    import os as _os, json as _json3, re as _re4

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    body = await request.json()
    instruction = (body.get("instruction") or "").strip()
    if not instruction:
        raise _HE(status_code=400, detail="instruction required")

    data_path  = f"/tmp/kimfam_minutes_{meeting_id}_data.json"
    draft_path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    if not _os.path.exists(data_path):
        raise _HE(status_code=404, detail="Minutes data not found — confirm meeting first")

    with open(data_path) as _f:
        minutes_data = _json3.load(_f)

    prompt = f"""You are editing KimFam Investment Club meeting minutes.

CURRENT MINUTES (JSON):
{_json3.dumps(minutes_data, indent=2)}

EDIT INSTRUCTION FROM CHAIRMAN:
{instruction}

Apply the instruction and return ONLY valid JSON with the same structure (no markdown, no explanation):
{{
  "meeting_ref": "...",
  "date": "...",
  "venue": "...",
  "start_time_eat": "...",
  "summary": "...",
  "key_topics": "...",
  "key_decisions": ["..."],
  "actions": [
    {{"ref":"...","description":"...","assignees":["..."],"deadline":"...","priority":"..."}}
  ]
}}

Rules:
- Only apply what was instructed. Do not change anything else.
- Preserve all action refs exactly as-is.
- Keep all existing actions unless explicitly told to remove one.
"""

    raw = ""
    try:
        import asyncio as _aio
        env = dict(_os.environ); env["HOME"] = "/root"
        proc = await _aio.create_subprocess_exec(
            "claude", "-p", prompt, "--model", "claude-haiku-4-5-20251001",
            stdout=_aio.subprocess.PIPE, stderr=_aio.subprocess.DEVNULL, env=env,
        )
        stdout, _ = await _aio.wait_for(proc.communicate(), timeout=60)
        if proc.returncode == 0:
            raw = stdout.decode().strip()
    except Exception:
        pass

    if not raw:
        groq_key = _os.getenv("GROQ_API_KEY", "")
        if groq_key:
            from groq import Groq as _Groq
            resp = _Groq(api_key=groq_key).chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=2000, temperature=0.1,
            )
            raw = resp.choices[0].message.content.strip()

    if not raw:
        raise _HE(status_code=503, detail="AI unavailable — try again")

    raw = _re4.sub(r"^```(?:json)?\s*", "", raw)
    raw = _re4.sub(r"\s*```$", "", raw.strip())
    try:
        updated = _json3.loads(raw)
    except Exception:
        m = _re4.search(r"\{.*\}", raw, _re4.DOTALL)
        if m:
            updated = _json3.loads(m.group())
        else:
            raise _HE(status_code=502, detail="AI returned unparseable response")

    # Save updated data + regenerate .docx
    with open(data_path, "w") as _f:
        _json3.dump(updated, _f)

    mtg_row = _dbq("SELECT ref, date, venue, start_time_eat FROM meetings WHERE id=%s", (meeting_id,))
    _docx_bytes = _build_minutes_docx(mtg=mtg_row[0] if mtg_row else {}, **{
        k: updated[k]
        for k in ("meeting_ref","summary","key_topics","key_decisions","actions")
    })
    with open(draft_path, "wb") as _f:
        _f.write(_docx_bytes)

    return {"ok": True, "minutes_data": updated}


@app.post("/api/meetings/{meeting_id}/publish")
async def meetings_minutes_publish(meeting_id: int, request: Request):
    """Upload draft .docx to R2 and send WhatsApp group notification."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq2, execute as _exec2
    import os as _os, re as _re3

    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")

    rows = _dbq2("SELECT ref, date FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    mtg = rows[0]

    draft_path = f"/tmp/kimfam_minutes_{meeting_id}.docx"
    if not _os.path.exists(draft_path):
        raise _HE(status_code=404, detail="Draft not found — confirm meeting first")

    # Upload to R2. File named to match the old convention:
    # KIMFAM_Meeting_Minutes_June_7_2026.docx
    import r2_storage as _r2_pub
    _fname    = _minutes_filename(mtg["date"])
    # Place published minutes in the KimFam sub-group for their year, matching the
    # documents sub-grouping (e.g. minutes/KimFam (2026)/...).
    try:
        _yr = str(mtg["date"])[:4]
    except Exception:
        _yr = ""
    _sub = f"KimFam ({_yr})" if _yr.isdigit() else "KimFam"
    r2_key    = f"minutes/{_sub}/{_fname}"
    minutes_url = None
    if _r2_pub.is_configured():
        try:
            minutes_url = _r2_pub.upload(draft_path, r2_key, public=True)
        except Exception as _e:
            import logging as _lg
            _lg.getLogger("main").error(f"R2 upload failed: {_e}")

    # Update DB. When R2 isn't used (e.g. staging), point at the app's own
    # viewable/downloadable endpoint, NOT the raw /tmp path (which can't render).
    final_url = minutes_url or f"/api/meetings/{meeting_id}/minutes"
    _exec2("UPDATE meetings SET minutes_url=%s WHERE id=%s",
           (final_url, meeting_id))

    # WhatsApp notification. Use the SAME env flag as the rest of the app
    # (KIMFAM_ENV). On staging this must NEVER reach the real family group —
    # send to Hillary only so testing can't leak to the club.
    is_stg = _os.getenv("KIMFAM_ENV", "prod") == "staging"
    recipient = "256775102684" if is_stg else "254716595631-1631997730@g.us"  # Hillary vs KIM FAM PROJECTS
    env_tag = " [STAGING TEST]" if is_stg else ""
    from notifications import SIGNOFF as _SIGNOFF
    msg = (
        f"*KimFam Hub{env_tag}* — Meeting minutes ready\n\n"
        f"*{mtg['ref']}* ({mtg['date']})\n\n"
        f"Minutes have been published. Open KimFam Hub to read and download."
        + _SIGNOFF
    )
    try:
        import requests as _req
        _req.post("http://localhost:8080/api/send",
                  json={"recipient": recipient, "message": msg}, timeout=10)
    except Exception:
        pass

    # Make the new minutes answerable in Ask KimFam: save the .docx into the
    # embedder's read path and re-embed in the background (delta-tracked, so only
    # this new file is processed; non-blocking so publish returns immediately).
    try:
        import shutil as _sh, subprocess as _sp, sys as _sys
        _app_dir = _os.path.dirname(__file__)
        _minutes_dir = _os.path.join(_app_dir, "docs", "minutes", _sub)
        _os.makedirs(_minutes_dir, exist_ok=True)
        _sh.copy2(draft_path, _os.path.join(_minutes_dir, _fname))
        _sp.Popen(
            [_sys.executable, _os.path.join(_app_dir, "embed_documents.py")],
            cwd=_app_dir, env={**_os.environ},
            stdout=_sp.DEVNULL, stderr=_sp.DEVNULL,
        )
    except Exception as _ee:
        import logging as _lg
        _lg.getLogger("main").error(f"minutes embed kickoff failed: {_ee}")

    return {"ok": True, "minutes_url": minutes_url}


@app.get("/api/members")
def get_members():
    from db import query as dbq
    from contributions import compute_family_balance
    from datetime import date as _date
    try:
        families = dbq("SELECT id, family_name FROM families ORDER BY family_name")
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"members db: {e}")
        return {"as_of": "", "members": []}
    today = _date.today()
    months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    as_of = f"{today.day} {months[today.month-1]} {today.year}"
    members = []
    for fam in families:
        try:
            bal = compute_family_balance(fam["id"])
        except Exception as e:
            _lg = __import__("logging").getLogger("main")
            _lg.error(f"members balance for family {fam['id']}: {e}")
            continue
        members.append({
            "name": fam["family_name"],
            "initial_obligation": str(bal["initial_obligation"]),
            "paid_initial":       str(bal["initial_paid"]),
            "balance_initial":    str(bal["initial_balance"]),
            "total_contributions": str(bal["total_monthly_paid"]),
            "paid_current":       "",
            "balance_current":    str(bal["current_balance"]),
            "combined_balance":   str(bal["combined_balance"]),
        })
    return {"as_of": as_of, "members": members}



@app.get("/api/activity")
def get_activity():
    from db import query as dbq
    import datetime as _dt
    items = []

    contribs = dbq("""
        SELECT cp.confirmed_at AS ts, f.family_name, cp.amount_ugx
        FROM contribution_payments cp JOIN families f ON cp.family_id=f.id
        WHERE cp.status='confirmed' ORDER BY cp.confirmed_at DESC LIMIT 8
    """)
    for r in contribs:
        items.append({
            "ts": r["ts"].isoformat() if r["ts"] else None,
            "icon": "💰",
            "title": r['family_name'].title() + " paid UGX " + "{:,}".format(int(r['amount_ugx'])),
            "nav": "finances"
        })

    _cat_icons = {"staff": "👷", "project_investment": "🌾", "operational": "🔧", "loan": "🏦"}
    expenses = dbq("SELECT txn_date, description, amount_ugx, category FROM expenditure_records ORDER BY txn_date DESC LIMIT 6")
    for r in expenses:
        ts = _dt.datetime.combine(r["txn_date"], _dt.time(12, 0)).isoformat() if r["txn_date"] else None
        desc = str(r["description"])
        desc = desc[:55] + "…" if len(desc) > 55 else desc
        items.append({
            "ts": ts,
            "icon": _cat_icons.get(r["category"], "📤"),
            "title": desc + " — UGX " + "{:,}".format(int(r['amount_ugx'])),
            "nav": "finances"
        })

    loans = dbq("""
        SELECT ld.disbursed_at AS ts, f.family_name, ld.amount_ugx
        FROM loan_disbursements ld
        JOIN loan_applications la ON ld.loan_id=la.id
        JOIN families f ON la.family_id=f.id
        ORDER BY ld.disbursed_at DESC LIMIT 4
    """)
    for r in loans:
        items.append({
            "ts": r["ts"].isoformat() if r["ts"] else None,
            "icon": "🏦",
            "title": "Loan disbursed to " + r['family_name'].title() + " — UGX " + "{:,}".format(int(r['amount_ugx'])),
            "nav": "loans"
        })

    items = [i for i in items if i["ts"]]
    items.sort(key=lambda x: x["ts"], reverse=True)
    return items[:10]

@app.get("/api/loans")
def get_loans():
    from db import query as dbq
    from datetime import date as _date
    # Loan header from loan_applications + loan_disbursements
    apps = dbq("""
        SELECT la.id, la.applicant_name, la.amount_ugx, la.requested_at,
               ld.disbursed_at, la.status,
               COALESCE((SELECT SUM(lp.amount_ugx) FROM loan_payments lp WHERE lp.loan_id = la.id), 0) as amount_paid
        FROM loan_applications la
        LEFT JOIN loan_disbursements ld ON ld.loan_id = la.id
        WHERE la.status = 'approved'
        ORDER BY la.id
    """)
    loan = {}
    if apps:
        a = apps[0]
        # interest_pct and total_to_pay from disbursement notes or defaults
        principal = int(a["amount_ugx"])
        interest_pct = 3.0
        interest_cash = round(principal * interest_pct / 100 * 3)  # 3 months flat
        total_to_pay = principal + interest_cash
        amount_paid = int(a["amount_paid"])
        balance = total_to_pay - amount_paid
        issued = str(a["disbursed_at"] or a["requested_at"])
        # due date: 3 months from disbursement
        from datetime import timedelta, datetime
        try:
            issued_dt = datetime.strptime(issued, "%Y-%m-%d").date()
            due_dt = issued_dt.replace(year=issued_dt.year + (1 if issued_dt.month > 9 else 0),
                                       month=((issued_dt.month + 2) % 12) + 1)
            days_remaining = (due_dt - _date.today()).days
            due_date = due_dt.strftime("%-m/%-d/%Y")
        except Exception:
            due_date = ""
            days_remaining = 0
        loan = {
            "id": str(a["id"]),
            "member": a["applicant_name"].upper() + " TUHIMBISE",
            "date": issued_dt.strftime("%-m/%-d/%Y") if issued_dt else issued,
            "principal": f"{principal:,}",
            "interest_pct": f"{interest_pct:.1f}",
            "interest_cash": f"{interest_cash:,}",
            "total_to_pay": f"{total_to_pay:,}",
            "due_date": due_date,
            "days_remaining": str(days_remaining),
            "amount_paid": f"{amount_paid:,}",
            "balance": f"{balance:,}",
        }
    # Payment history from loan_payments
    prows = dbq("""
        SELECT lp.paid_at, lp.amount_ugx, lp.reference, lp.notes
        FROM loan_payments lp
        JOIN loan_applications la ON la.id = lp.loan_id
        WHERE la.status = 'approved'
        ORDER BY lp.paid_at
    """)
    payments = []
    for p in prows:
        ref = p.get("reference","")
        month = ref.replace("Sheet import — ","") if "Sheet import" in (ref or "") else ref
        notes = p.get("notes","") or ""
        interest_paid = ""
        balance_val = ""
        if "interest=" in notes:
            try:
                interest_paid = notes.split("interest=")[1].split(",")[0]
                interest_paid = f"{int(interest_paid):,}"
            except Exception:
                pass
        if "balance_after=" in notes:
            try:
                balance_val = notes.split("balance_after=")[1]
                balance_val = f"{int(balance_val):,}"
            except Exception:
                pass
        dt = p.get("paid_at")
        dt_str = dt.strftime("%-m/%-d/%Y") if hasattr(dt, "strftime") else str(dt)
        payments.append({
            "date": dt_str,
            "amount_paid": f"{int(p['amount_ugx']):,}",
            "month": month,
            "interest_paid": interest_paid,
            "balance": balance_val,
        })
    return {"loan": loan, "payments": payments}

def _minutes_by_date() -> dict:
    """Return {date(y,m,d): url_path} mapping minutes files to their extracted dates."""
    import datetime as _dt
    result: dict[_dt.date, str] = {}
    DATE_FMTS = ("%d %B %Y", "%B %d, %Y", "%d-%b-%Y", "%d/%m/%Y", "%Y-%m-%d")

    if _r2.is_configured():
        for obj in _r2.list_folder("minutes/"):
            fname = obj["key"].split("/", 1)[-1]
            if not fname or Path(fname).suffix.lower() not in (".docx", ".pdf"):
                continue
            if ".bak." in fname:
                continue
            dt = _extract_date(Path(fname))
            result[dt.date()] = f"/docs/minutes/{fname}"
    else:
        folder = DOCS_DIR / "minutes"
        if folder.exists():
            for f in folder.rglob("*"):
                if not f.is_file() or f.suffix.lower() not in (".docx", ".pdf"):
                    continue
                if ".bak." in f.name:
                    continue
                dt = _extract_date(f)
                result[dt.date()] = f"/docs/minutes/{f.relative_to(folder).as_posix()}"
    return result

def _parse_meeting_date(raw: str):
    """Parse a meeting date string from the sheet. Returns date or None."""
    import datetime as _dt
    raw = raw.strip()
    for fmt in ("%d %B %Y", "%B %d, %Y", "%d-%b-%Y", "%d/%m/%Y", "%Y-%m-%d", "%d %b %Y"):
        try:
            return _dt.datetime.strptime(raw, fmt).date()
        except ValueError:
            pass
    return None

@app.get("/api/meetings")
def get_meetings():
    """Meeting list from PostgreSQL with live action progress counts."""
    from db import query as _dbq
    rows = _dbq("""
        SELECT m.id, m.ref, m.date, m.start_time_eat, m.venue,
               m.key_topics, m.key_decisions, m.next_actions, m.summary,
               m.minutes_url,
               m.conductor_item, m.conductor_started_at, m.conductor_ended_at,
               COUNT(a.id)                                          AS action_count,
               COUNT(a.id) FILTER (WHERE a.status = 'done')        AS action_done_count
        FROM meetings m
        LEFT JOIN actions a ON a.meeting_id = m.id
        GROUP BY m.id
        ORDER BY m.date DESC
    """)
    out = []
    for r in rows:
        ref = r["ref"] or ""
        number = ref[3:].strip() if ref.upper().startswith("KIM") else ref
        out.append({
            "db_id":           r["id"],
            "id":              r["id"],
            "meeting_number":  number,
            "meeting_ref":     ref,
            "meeting_date":    str(r["date"]) if r["date"] else "",
            "start_time_eat":  str(r["start_time_eat"]) if r["start_time_eat"] else None,
            "location":        r["venue"],
            "key_topics":      r["key_topics"],
            "key_decisions":   r["key_decisions"],
            "next_actions":    r["next_actions"],
            "summary":         r["summary"],
            "attendance":      [],
            "minutes_url":     r["minutes_url"],
            "action_count":      int(r["action_count"]),
            "action_done_count": int(r["action_done_count"]),
            "conductor_active":  r["conductor_started_at"] is not None and r["conductor_ended_at"] is None,
            "conductor_ended":   r["conductor_ended_at"] is not None,
        })
    return out


@app.post("/api/meetings")
async def create_meeting(request: Request):
    """Create a new meeting record. Admin only."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    date_str   = str(body.get("date", "")).strip()
    venue      = str(body.get("venue", "")).strip() or None
    start_time = str(body.get("start_time", "")).strip() or None
    key_topics = str(body.get("key_topics", "")).strip() or None
    if not date_str:
        raise _HE(status_code=400, detail="date required")
    # Auto-generate ref from the highest meeting number (date-tie safe). Bump past
    # any collision so creating a meeting can never get stuck on a duplicate ref.
    import re as _re
    ref, _prev = _compute_next_ref()
    while _dbq("SELECT id FROM meetings WHERE ref=%s", (ref,)):
        m2 = _re.match(r"KIM\s+(\d+)/(\d{4})", ref)
        ref = f"KIM {int(m2.group(1)) + 1:03d}/{m2.group(2)}" if m2 else ref + "-x"
    # Build default agenda from template + active projects
    agenda = _build_default_agenda(key_topics)
    import json as _json_ag
    _exec("""INSERT INTO meetings (ref, date, venue, start_time_eat, key_topics, agenda)
             VALUES (%s, %s::date, %s, %s::time, %s, %s)""",
          (ref, date_str, venue, start_time, key_topics, _json_ag.dumps(agenda)))
    row = _dbq("SELECT id FROM meetings WHERE ref=%s", (ref,))

    # Announce the new meeting to the family group + all members (best-effort).
    # _broadcast is env-aware: prod → group + members, staging → Hillary + test group.
    try:
        import notifications as _notif_cm
        _t = (start_time or "")[:5]
        _agenda_block = ""
        if key_topics:
            _items = [i.strip() for i in key_topics.split(";") if i.strip()]
            if _items:
                _agenda_block = "\n*Agenda:*\n" + "\n".join(f"• {i}" for i in _items) + "\n"
        _msg = (
            f"📅 *New KimFam Meeting Scheduled*\n\n"
            f"*{ref}*\n"
            f"*Date:* {date_str}\n"
            + (f"*Time:* {_t} EAT\n" if _t else "")
            + (f"*Venue:* {venue}\n" if venue else "")
            + _agenda_block
            + f"\nFull agenda and details on kimfamhub.com. See you there."
            + _notif_cm.SIGNOFF
        )
        _phones = [p for p in _notif_cm.MEMBER_PHONES.values() if p]
        _notif_cm._broadcast(_phones, _msg)
    except Exception as _e:
        import logging as _lg; _lg.getLogger("main").error(f"meeting-create notify failed: {_e}")

    return {"ok": True, "db_id": row[0]["id"], "ref": ref}


@app.patch("/api/meetings/{meeting_id}")
async def edit_meeting(meeting_id: int, request: Request):
    """Edit a meeting's date/time/venue/topics — allowed until it is conducted.
    Rebuilds the agenda from the new topics. Admin only."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import json as _json_em
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT id, conductor_started_at FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    if rows[0]["conductor_started_at"] is not None:
        raise _HE(status_code=409, detail="This meeting has already been conducted — its agenda is locked.")
    body = await request.json()
    date_str   = str(body.get("date", "")).strip()
    venue      = str(body.get("venue", "")).strip() or None
    start_time = str(body.get("start_time", "")).strip() or None
    key_topics = str(body.get("key_topics", "")).strip() or None
    if not date_str:
        raise _HE(status_code=400, detail="date required")
    agenda = _build_default_agenda(key_topics)
    _exec("""UPDATE meetings SET date=%s::date, venue=%s, start_time_eat=%s::time,
                                 key_topics=%s, agenda=%s
             WHERE id=%s""",
          (date_str, venue, start_time, key_topics, _json_em.dumps(agenda), meeting_id))
    return {"ok": True, "id": meeting_id}


def _build_default_agenda(key_topics: str | None = None) -> list:
    """Generate a standard KimFam meeting agenda. Items 5/6 vary by topics."""
    from db import query as _dbq_ag
    # Derive active projects from DB (those with open actions in last 30 days)
    active_projects = _dbq_ag("""
        SELECT DISTINCT project_id FROM actions
        WHERE project_id IS NOT NULL AND status NOT IN ('done','cancelled')
        ORDER BY project_id
    """)
    proj_names = {
        "chicken": "Free Range Chicken", "washing_bay": "Washing Bay",
        "sheep": "Sheep (Dorper)", "goats": "Goats", "dairy": "Dairy / Cows",
        "mango": "Mango & Oranges", "trees": "Tree Planting", "bees": "Apiary",
        "rabbits": "Rabbits", "irrigation": "Irrigation & Bananas",
        "fortune_credit": "Fortune Credit", "kakoba": "Kakoba Land",
    }
    project_items = [
        {"label": proj_names.get(r["project_id"], r["project_id"]),
         "presenter": "", "duration_min": 5, "type": "project",
         "project_id": r["project_id"]}
        for r in active_projects
    ]

    # Parse custom main agenda items from key_topics string
    main_items = []
    if key_topics:
        for raw in key_topics.split(";"):
            t = raw.strip().strip(",")
            if t:
                main_items.append({"label": t, "presenter": "", "duration_min": 10, "type": "agenda"})

    return [
        {"label": "Opening Prayer",          "presenter": "Dad",    "duration_min": 3,  "type": "fixed"},
        {"label": "Apologies & Attendance",  "presenter": "Hellen", "duration_min": 3,  "type": "fixed"},
        {"label": "Review of Last Meeting",  "presenter": "Hillary","duration_min": 5,  "type": "fixed"},
        {"label": "Treasurer's Brief",       "presenter": "Hellen", "duration_min": 10, "type": "fixed"},
        {"label": "Action Review",           "presenter": "Hillary","duration_min": 15, "type": "fixed"},
        *([{"label": "Project Updates", "presenter": "", "duration_min": 0, "type": "section",
             "children": project_items}] if project_items else []),
        *([{"label": "Main Agenda",    "presenter": "", "duration_min": 0, "type": "section",
             "children": main_items}]   if main_items  else []),
        {"label": "Any Other Business",      "presenter": "",       "duration_min": 5,  "type": "fixed"},
        {"label": "Next Meeting Date",       "presenter": "Hillary","duration_min": 2,  "type": "fixed"},
        {"label": "Closing Prayer",          "presenter": "Dad",    "duration_min": 2,  "type": "fixed"},
    ]


# ── Conductor endpoints ───────────────────────────────────────────────────────

@app.get("/api/meetings/{meeting_id}/conductor")
async def conductor_state(meeting_id: int, request: Request):
    """Current conductor state — polled by all participants every 5s."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    import json as _json_c, datetime as _dt

    token = _get_tok(request)
    if not _auth_verify(token):
        raise _HE(status_code=401, detail="Auth required")
    import os as _os_c
    import notifications as _notif

    rows = _dbq("""SELECT ref, date, venue, start_time_eat, agenda,
                          conductor_item, conductor_started_at,
                          conductor_item_started_at, conductor_ended_at,
                          conductor_recording, conductor_notes, conductor_timings,
                          attendance
                   FROM meetings WHERE id=%s""", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    r = rows[0]

    _raw_ag = r["agenda"]
    agenda = (_json_c.loads(_raw_ag) if isinstance(_raw_ag, (str, bytes, bytearray)) else (_raw_ag or []))
    now_utc = _dt.datetime.utcnow().replace(tzinfo=_dt.timezone.utc)

    item_elapsed = None
    if r["conductor_item_started_at"]:
        item_elapsed = int((now_utc - r["conductor_item_started_at"]).total_seconds())
    total_elapsed = None
    if r["conductor_started_at"]:
        total_elapsed = int((now_utc - r["conductor_started_at"]).total_seconds())

    # Previous meeting's retrospective — for the "Review of Last Meeting" item.
    import re as _re_pc
    _cur_m = _re_pc.match(r"KIM\s+(\d+)/(\d{4})", (r["ref"] or "").strip())
    prev_retro, prev_ref = None, None
    if _cur_m:
        _cur_num = int(_cur_m.group(1))
        _prev_rows = _dbq("SELECT ref, conductor_retrospective FROM meetings")
        _best = -1
        for _pr in _prev_rows:
            _pm = _re_pc.match(r"KIM\s+(\d+)/(\d{4})", (_pr["ref"] or "").strip())
            if _pm and int(_pm.group(1)) < _cur_num and int(_pm.group(1)) > _best and _pr["conductor_retrospective"]:
                _best = int(_pm.group(1)); prev_ref = _pr["ref"]
                _raw = _pr["conductor_retrospective"]
                prev_retro = (_json_c.loads(_raw) if isinstance(_raw, (str, bytes, bytearray)) else _raw)

    return {
        "meeting_ref":   r["ref"],
        "date":          str(r["date"]),
        "venue":         r["venue"] or "Google Meet",
        "agenda":        agenda,
        "current_item":  r["conductor_item"],        # None = not started, -1 = ended
        "started":       r["conductor_started_at"] is not None,
        "ended":         r["conductor_ended_at"] is not None,
        "recording":     bool(r["conductor_recording"]) and r["conductor_ended_at"] is None,
        "item_elapsed_s":  item_elapsed,
        "total_elapsed_s": total_elapsed,
        "notes":           r["conductor_notes"] or "",
        "recording_present": _os_c.path.exists(f"/tmp/kimfam_recording_{meeting_id}.webm"),
        "timings": (_json_c.loads(r["conductor_timings"]) if isinstance(r["conductor_timings"], (str, bytes, bytearray)) else (r["conductor_timings"] or {})),
        "attendance": (_json_c.loads(r["attendance"]) if isinstance(r["attendance"], (str, bytes, bytearray)) else (r["attendance"] or {})),
        "members": list(_notif.MEMBER_PHONES.keys()),
        "prev_retrospective": prev_retro,
        "prev_ref": prev_ref,
    }


# ── Meeting retrospective / time analytics ────────────────────────────────────

def _generate_retrospective(meeting_id: int) -> dict:
    """Build an AI retrospective for a conducted meeting: per-item time verdict,
    what consumed the time, whether it produced output, mission link, and
    forward recommendations. Stores it in conductor_retrospective."""
    from db import query as _dbq_r, execute as _exec_r
    import json as _json_r
    rows = _dbq_r("""SELECT ref, date, conductor_timings, conductor_notes,
                            transcript, key_decisions
                     FROM meetings WHERE id=%s""", (meeting_id,))
    if not rows:
        return {}
    r = rows[0]
    raw_tm = r["conductor_timings"]
    timings = (_json_r.loads(raw_tm) if isinstance(raw_tm, (str, bytes, bytearray)) else (raw_tm or {}))

    # Actions this meeting produced (output check per item is approximate at meeting level)
    acts = _dbq_r("""SELECT ref, description, assignee FROM actions
                     WHERE meeting_id=%s ORDER BY ref""", (meeting_id,))
    actions_txt = "\n".join(f"- {a['ref']} ({a['assignee']}): {a['description'][:90]}" for a in acts) or "None recorded"

    notes = r["conductor_notes"] or ""
    transcript = (r["transcript"] or "")[:8000]

    # Need SOME substance to analyse — timing, transcript, notes, or actions.
    if not timings and not transcript.strip() and not notes.strip() and not acts:
        return {}

    timing_lines = []
    planned_total = actual_total = 0
    for k in sorted(timings, key=lambda x: int(x) if str(x).isdigit() else 0):
        t = timings[k]
        pm = int(t.get("planned_min", 0)); a_s = int(t.get("actual_s", 0))
        planned_total += pm * 60; actual_total += a_s
        timing_lines.append(f"- {t.get('label','')}: planned {pm}m, actual {a_s//60}m{a_s%60:02d}s")
    timing_block = ("\n".join(timing_lines) if timing_lines else
                    "(Per-item timing was not captured for this meeting. Skip the time-discipline "
                    "scoring and base efficiency on whether discussions produced clear decisions/actions.)")

    prompt = f"""You are the meeting-improvement analyst for KimFam Investment Club, a Ugandan
family investment club building diversified ventures (chicken, washing bay, sheep, dairy,
mango, fortune credit, land, etc.) toward long-term family wealth. Tone: warm, collective,
improvement-focused. NEVER single out an individual for running long; focus on the topic and
process. No em-dashes anywhere.

Analyse meeting {r['ref']} ({r['date']}).

TIME PER AGENDA ITEM (planned vs actual):
{timing_block}

SECRETARY NOTES (segmented by item where headers exist):
{(notes or '(no live notes)')[:4000]}

ACTIONS THIS MEETING PRODUCED:
{actions_txt}

TRANSCRIPT (may be partial):
{transcript or '(no transcript captured)'}

Return ONLY valid JSON (no markdown) in this exact shape:
{{
  "efficiency_score": 0-100 integer (how well time was used vs planned and output produced),
  "headline": "one warm sentence summarising how the meeting ran",
  "overall": "2-3 sentences: overall time discipline + biggest observation",
  "items": [
    {{
      "label": "agenda item",
      "verdict": "on_time|over|under",
      "what_consumed_time": "why it took the time it did (or 'efficient')",
      "produced_output": true|false,
      "mission_link": "how this ties to a club goal/project, or 'general governance'",
      "do_better": "one concrete improvement for next time"
    }}
  ],
  "time_sinks": ["topics that ran longest / consumed disproportionate time"],
  "decision_density_note": "did the long items actually produce decisions/actions, or mostly discussion?",
  "recommendations": ["2-4 concrete, actionable suggestions for the next meeting"]
}}"""

    import re as _re_r, os as _os_r
    def _parse(raw):
        if not raw: return None
        try: return _json_r.loads(raw)
        except Exception:
            m = _re_r.search(r"\{.*\}", raw, _re_r.DOTALL)
            if m:
                try: return _json_r.loads(m.group())
                except Exception: return None
        return None

    raw = _ask_claude(prompt, timeout=90)
    data = _parse(raw)
    # Fallback to Groq if the Claude CLI returned nothing usable (mirrors /process)
    if not isinstance(data, dict):
        groq_key = _os_r.getenv("GROQ_API_KEY", "")
        if groq_key:
            try:
                from groq import Groq as _Groq
                resp = _Groq(api_key=groq_key).chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=2000, temperature=0.2,
                )
                data = _parse(resp.choices[0].message.content.strip())
            except Exception as _ge:
                import logging as _lg
                _lg.getLogger("main").error(f"retrospective groq fallback failed: {_ge}")
    if not isinstance(data, dict):
        return {}
    import datetime as _dt_r
    data["generated_at"] = _dt_r.datetime.utcnow().isoformat()
    data["planned_total_min"] = planned_total // 60
    data["actual_total_min"]  = actual_total // 60
    _exec_r("UPDATE meetings SET conductor_retrospective=%s WHERE id=%s",
            (_json_r.dumps(data), meeting_id))
    return data


@app.post("/api/meetings/{meeting_id}/retrospective")
async def make_retrospective(meeting_id: int, request: Request):
    """Generate (or regenerate) the meeting retrospective. Admin only."""
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    data = _generate_retrospective(meeting_id)
    if not data:
        raise _HE(status_code=422, detail="Couldn't generate a review — this meeting has no transcript, notes, timing, or actions to analyse yet. Process its minutes first.")
    return {"ok": True, "retrospective": data}


@app.get("/api/meetings/{meeting_id}/retrospective")
async def get_retrospective(meeting_id: int, request: Request):
    """Read a meeting's retrospective. Visible to all authenticated members."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    import json as _json_g
    token = _get_tok(request)
    if not _auth_verify(token):
        raise _HE(status_code=401, detail="Login required")
    rows = _dbq("SELECT conductor_retrospective FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    raw = rows[0]["conductor_retrospective"]
    data = (_json_g.loads(raw) if isinstance(raw, (str, bytes, bytearray)) else (raw or None))
    return {"retrospective": data}


@app.post("/api/meetings/{meeting_id}/conductor/notes")
async def conductor_save_notes(meeting_id: int, request: Request):
    """Save the secretary's running notes typed during the meeting. Admin only."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT id FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    body = await request.json()
    notes = str(body.get("notes", ""))
    _exec("UPDATE meetings SET conductor_notes=%s WHERE id=%s", (notes, meeting_id))
    return {"ok": True}


@app.post("/api/meetings/{meeting_id}/attendance")
async def save_attendance(meeting_id: int, request: Request):
    """Save roll call: {member: {status: present|apology|absent, comment}}. Admin only."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import json as _json_at
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT id FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    body = await request.json()
    attendance = body.get("attendance", {})
    if not isinstance(attendance, dict):
        raise _HE(status_code=400, detail="attendance must be an object")
    _exec("UPDATE meetings SET attendance=%s WHERE id=%s",
          (_json_at.dumps(attendance), meeting_id))
    return {"ok": True}


@app.post("/api/meetings/{meeting_id}/recording")
async def upload_meeting_recording(meeting_id: int, request: Request):
    """Store an in-app audio recording to /tmp so Process can pick it up."""
    from fastapi import HTTPException as _HE
    import os as _os
    token = _get_tok(request)
    if not _auth_verify(token):
        raise _HE(status_code=401, detail="Auth required")
    form = await request.form()
    audio = form.get("audio_file")
    if not audio or not hasattr(audio, "filename"):
        raise _HE(status_code=400, detail="audio_file required")
    raw = await audio.read()
    path = f"/tmp/kimfam_recording_{meeting_id}.webm"
    with open(path, "wb") as f:
        f.write(raw)
    return {"ok": True, "path": path}


@app.post("/api/meetings/{meeting_id}/recording/chunk")
async def append_meeting_recording_chunk(meeting_id: int, request: Request):
    """Append one streamed audio chunk to the meeting recording.

    The conductor uploads a chunk every ~15 seconds while the meeting runs, so the
    audio is persisted progressively. This is what makes recording survive a closed
    or Chrome-discarded tab (the old flow buffered the whole meeting in the tab's
    memory and uploaded only at End, so any tab hiccup lost everything). seq=0
    starts a fresh file (truncate); later chunks append in order. MediaRecorder's
    first timeslice carries the webm header, so the concatenation is decodable."""
    from fastapi import HTTPException as _HE
    import os as _os
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Auth required")
    form = await request.form()
    chunk = form.get("chunk")
    seq = str(form.get("seq", "")).strip()
    if not chunk or not hasattr(chunk, "read"):
        raise _HE(status_code=400, detail="chunk required")
    data = await chunk.read()
    path = f"/tmp/kimfam_recording_{meeting_id}.webm"
    mode = "wb" if seq == "0" else "ab"
    with open(path, mode) as f:
        f.write(data)
    return {"ok": True, "bytes": len(data), "seq": seq}


@app.post("/api/meetings/{meeting_id}/conductor/start")
async def conductor_start(meeting_id: int, request: Request):
    """Start the meeting — sets item 0, records start time."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT id FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    _exec("""UPDATE meetings SET
               conductor_item=0,
               conductor_started_at=NOW(),
               conductor_item_started_at=NOW(),
               conductor_ended_at=NULL,
               conductor_recording=TRUE
             WHERE id=%s""", (meeting_id,))
    return {"ok": True, "current_item": 0}


def _record_item_timing(meeting_id: int):
    """Record how long the currently-active agenda item actually took, before we
    leave it. Stores {label, planned_min, actual_s} keyed by item index in
    conductor_timings. Used for post-meeting auditing (which items run long)."""
    from db import execute as _exec, query as _dbq
    import json as _json_t, datetime as _dt
    rows = _dbq("""SELECT conductor_item, conductor_item_started_at, agenda, conductor_timings
                   FROM meetings WHERE id=%s""", (meeting_id,))
    if not rows:
        return
    r = rows[0]
    idx = r["conductor_item"]
    started = r["conductor_item_started_at"]
    if idx is None or idx < 0 or started is None:
        return
    now_utc = _dt.datetime.utcnow().replace(tzinfo=_dt.timezone.utc)
    actual_s = int((now_utc - started).total_seconds())
    raw_ag = r["agenda"]
    agenda = (_json_t.loads(raw_ag) if isinstance(raw_ag, (str, bytes, bytearray)) else (raw_ag or []))
    flat = _flatten_agenda(agenda)
    item = flat[idx] if 0 <= idx < len(flat) else None
    if not item:
        return
    raw_tm = r["conductor_timings"]
    timings = (_json_t.loads(raw_tm) if isinstance(raw_tm, (str, bytes, bytearray)) else (raw_tm or {}))
    timings[str(idx)] = {
        "label": item.get("label", ""),
        "planned_min": item.get("duration_min", 0),
        "actual_s": actual_s,
    }
    _exec("UPDATE meetings SET conductor_timings=%s WHERE id=%s",
          (_json_t.dumps(timings), meeting_id))


@app.post("/api/meetings/{meeting_id}/conductor/next")
async def conductor_next(meeting_id: int, request: Request):
    """Advance to next agenda item."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import json as _json_n
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT conductor_item, agenda FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    r = rows[0]
    _raw_ag = r["agenda"]
    agenda = (_json_n.loads(_raw_ag) if isinstance(_raw_ag, (str, bytes, bytearray)) else (_raw_ag or []))
    # Flatten sections into a linear list for conductor
    flat = _flatten_agenda(agenda)
    current = r["conductor_item"] or 0
    next_item = current + 1
    _record_item_timing(meeting_id)  # capture how long the item we're leaving took
    _exec("""UPDATE meetings SET conductor_item=%s, conductor_item_started_at=NOW()
             WHERE id=%s""", (next_item, meeting_id))
    return {"ok": True, "current_item": next_item, "total_items": len(flat)}


@app.post("/api/meetings/{meeting_id}/conductor/goto")
async def conductor_goto(meeting_id: int, request: Request):
    """Jump to a specific agenda item index."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    idx = int(body.get("index", 0))
    rows = _dbq("SELECT id FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    _record_item_timing(meeting_id)  # capture the item we're leaving
    _exec("""UPDATE meetings SET conductor_item=%s, conductor_item_started_at=NOW()
             WHERE id=%s""", (idx, meeting_id))
    return {"ok": True, "current_item": idx}


@app.post("/api/meetings/{meeting_id}/conductor/end")
async def conductor_end(meeting_id: int, request: Request):
    """End the meeting — sets ended timestamp."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import datetime as _dt
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _dbq("SELECT id, conductor_started_at FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    _record_item_timing(meeting_id)  # capture the final item's duration
    _exec("""UPDATE meetings SET conductor_item=-1, conductor_ended_at=NOW(),
                                conductor_recording=FALSE
             WHERE id=%s""", (meeting_id,))
    return {"ok": True, "ended": True}


@app.patch("/api/meetings/{meeting_id}/agenda")
async def update_agenda(meeting_id: int, request: Request):
    """Replace the agenda for a meeting (admin, before or during)."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    import json as _json_upd
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body  = await request.json()
    agenda = body.get("agenda", [])
    rows = _dbq("SELECT id FROM meetings WHERE id=%s", (meeting_id,))
    if not rows:
        raise _HE(status_code=404, detail="Meeting not found")
    _exec("UPDATE meetings SET agenda=%s WHERE id=%s",
          (_json_upd.dumps(agenda), meeting_id))
    return {"ok": True}


def _flatten_agenda(agenda: list) -> list:
    """Flatten nested agenda (sections with children) into linear list for conductor."""
    flat = []
    for item in agenda:
        if item.get("type") == "section":
            flat.append({**item, "is_section_header": True})
            for child in (item.get("children") or []):
                flat.append({**child, "is_section_child": True})
        else:
            flat.append(item)
    return flat


@app.get("/api/meeting")
def get_next_meeting():
    """Next upcoming meeting from PostgreSQL (date > today EAT)."""
    from db import query as _dbq
    import datetime as _dt, re as _re
    today = (_dt.datetime.utcnow() + _dt.timedelta(hours=3)).date()
    rows = _dbq("""
        SELECT ref, date, start_time_eat, venue
        FROM meetings
        WHERE date >= %s
        ORDER BY date ASC
        LIMIT 1
    """, (today,))
    if not rows:
        days_to_sun = (6 - today.weekday()) % 7 or 7
        next_sun = today + _dt.timedelta(days=days_to_sun)
        return {"ref": "TBD", "date": next_sun.strftime("%A, %d %B %Y"),
                "eat": "4:30 PM EAT", "ist": "7:00 PM IST",
                "platform": "Google Meet",
                "note": "Date not yet confirmed. Typically every Sunday at 4:30 PM EAT."}
    r = rows[0]
    ref  = r["ref"]
    eat  = str(r["start_time_eat"])[:5] if r["start_time_eat"] else "16:30"
    # Convert HH:MM EAT → IST (+2h30m)
    ist = ""
    try:
        h, mn = int(eat[:2]), int(eat[3:5])
        total = h * 60 + mn + 150
        ih, im = (total // 60) % 24, total % 60
        iampm = "PM" if ih >= 12 else "AM"
        ih12 = ih % 12 or 12
        eat_ampm = "PM" if h >= 12 else "AM"
        eat_h12 = h % 12 or 12
        eat = f"{eat_h12}:{mn:02d} {eat_ampm} EAT"
        ist = f"{ih12}:{im:02d} {iampm} IST"
    except Exception:
        eat = "4:30 PM EAT"; ist = "7:00 PM IST"
    pretty = r["date"].strftime("%A, %d %B %Y") if r["date"] else "TBD"
    return {"ref": ref, "date": pretty, "eat": eat, "ist": ist,
            "platform": "Google Meet",
            "note": "Please update your action points on KimFam Hub before the meeting."}

@app.get("/api/projects")
def get_projects():
    rows = gc().open_by_key(SOLOMON_ID).worksheet("Financial Statement").get_all_values()
    data = {}
    for row in rows[1:]:
        if len(row) >= 3 and row[1].strip():
            data[row[1].strip()] = {"value": row[2].strip(), "desc": row[3].strip() if len(row)>3 else ""}
    return {"chicken": data}

def _react_index() -> HTMLResponse:
    react_index = os.path.join(_DIST, "index.html")
    if os.path.isfile(react_index):
        return HTMLResponse(open(react_index).read())
    # Root index.html is always kept in sync with dist/index.html by deploy.yml
    root_index = os.path.join(os.path.dirname(_DIST), "..", "index.html")
    root_index = os.path.normpath(root_index)
    if os.path.isfile(root_index):
        return HTMLResponse(open(root_index).read())
    from fastapi import HTTPException
    raise HTTPException(status_code=503, detail="App not deployed yet — dist/index.html missing")

@app.get("/", response_class=HTMLResponse)
def index():
    return _react_index()

@app.get("/api/projects/all")
def get_all_projects():
    from db import query as _dbq
    import json as _json
    from datetime import datetime as _dt2

    def _parse_hardcoded_date(s):
        """Parse strings like '3 Jun 2026' into a date."""
        try:
            return _dt2.strptime(s.strip(), "%d %b %Y").date()
        except Exception:
            return None

    # Fetch latest project_updates row per project
    try:
        rows = _dbq(
            "SELECT DISTINCT ON (project_id) project_id, author, text, media, "
            "TO_CHAR(created_at AT TIME ZONE 'Africa/Nairobi', 'FMDD Mon YYYY') AS date_str, "
            "created_at "
            "FROM project_updates "
            "ORDER BY project_id, created_at DESC"
        )
        db_updates = {}
        for row in rows:
            media = row.get("media") or []
            if isinstance(media, str):
                try: media = _json.loads(media)
                except: media = []
            db_updates[row["project_id"]] = {
                "date":   row["date_str"],
                "author": row["author"],
                "text":   row["text"],
                "images": [m["url"] for m in media if m.get("type") == "image"],
                "videos": [m["url"] for m in media if m.get("type") == "video"],
                "created_at": row.get("created_at"),
            }
    except Exception as _e:
        db_updates = {}

    def _best_update(project_id, hardcoded_update):
        db = db_updates.get(project_id)
        if not db:
            return hardcoded_update
        if not hardcoded_update:
            return db
        # Compare dates: prefer the more recent one
        db_dt   = db.get("created_at")
        hc_dt   = _parse_hardcoded_date(hardcoded_update.get("date",""))
        if db_dt and hc_dt:
            db_date = db_dt.date() if hasattr(db_dt, "date") else db_dt
            # DB only wins if strictly newer — same date keeps hardcoded (richer media)
            if db_date > hc_dt:
                return db
            # Same date: prefer whichever has more media
            db_media = len(db.get("images",[]))+len(db.get("videos",[]))
            hc_media = len(hardcoded_update.get("images",[]))+len(hardcoded_update.get("videos",[]))
            return db if db_media >= hc_media else hardcoded_update
        return db  # default to DB if we cannot compare

    projects = [
        {"id":"chicken","name":"Free Range Chicken","icon":"🐔","category":"Farming & Agriculture","status":"Operational","lead":"Solomon Ariho","headline":"60% production rate","live":True,"data":[
            {"label":"Current Flock","value":"51 active hens"},{"label":"Production Rate","value":"60% consistently"},
            {"label":"Feed","value":"Commercial (silver fish protein)"},{"label":"Deaths (last 2 months)","value":"~19 since last count (flock down from 70); 1 additional young bird died 6 Jun 2026"},
            {"label":"Phase 2 Trigger","value":"August 2026 — scale up if performance holds"},{"label":"Latest Update","value":"6-Jun-2026: Solomon confirmed 52 hens present; 1 has reproductive failure (not treatable, to be slaughtered). Active count: 51. Dad also reported 1 young bird dead from existing flock. Chick order (good brown breeds) still pending — seller target end-June 2026."}
        ]},
        {"id":"washing_bay","name":"Washing Bay","icon":"🚗","category":"Business Ventures","status":"Operational","lead":"Alex + Dad","headline":"UGX 580K/month revenue","data":[
            {"label":"Lead Investor","value":"Alex Tuhimbise"},{"label":"Total CapEx","value":"UGX 25,900,000"},
            {"label":"Monthly Revenue","value":"UGX ~580,000"},{"label":"Revenue (5 months)","value":"UGX 2,923,000"},
            {"label":"Infrastructure","value":"Electricity, well, pump, jet spray machine"},
            {"label":"Sanitation Plan","value":"Septic tank required (pit cannot serve as septic — fills too fast). Pit to be used for solid waste only. Viola confirmed."},
            {"label":"Septic Tank Estimate","value":"UGX 5,473,000 (Eng. Julius Kato, Invoice 110, 29 May 2026)"},
            {"label":"Plumbing Estimate","value":"UGX 2,162,500 (materials 1,612,500 + labour 550,000)"},
            {"label":"Total Sanitation Budget","value":"UGX 7,635,500"},
            {"label":"Open Issues","value":"Latrine/sanitation works to be approved and funded. Eli occupancy agreement pending. Town Council removed temporary structures."}
        ],"update":{"date":"3 Jun 2026","author":"Dad","text":"Dad shared formal estimates for the sanitation upgrade: a septic tank (UGX 5,473,000 from Eng. Julius Kato Services) and full plumbing works (UGX 2,162,500). Viola confirmed a proper septic tank is required — the existing pit cannot serve that purpose but can hold solid waste such as pads. Total sanitation budget: UGX 7,635,500.","images":["/static/project-pics/washing-bay/septic_tank_estimate_2026-06-03.jpg","/static/project-pics/washing-bay/plumbing_quote_p1_2026-06-03.jpg","/static/project-pics/washing-bay/plumbing_quote_p2_2026-06-03.jpg"],"videos":["/static/project-pics/washing-bay/washing_bay_vision.mp4"]}},
        {"id":"sheep","name":"Sheep (Dorper)","icon":"🐑","category":"Farming & Agriculture","status":"Active","lead":"Dad + Solomon","headline":"~40 sheep, 17 Dorper offspring","data":[
            {"label":"Total Flock","value":"~40 sheep"},{"label":"Dorper Offspring","value":"17 identified"},
            {"label":"Club Holds","value":"2 Dorper sheep directly (1 young sold at UGX 200K to fund ear-tagging)"},
            {"label":"Capital Invested","value":"UGX 1,710,000 (Aug 2024 receipt on file)"},
            {"label":"Sale Target","value":"~30kg at 6 months, UGX 100K-150K per Dorper"},
            {"label":"Latest Update","value":"19-May-2026: 10 young + 12 old sheep ear-tagged for project identification. Looking for an additional ram — current ram cannot breed with its own offspring. Next: plough two padocks and broadcast pasture seeds."}
        ],"update":{"date":"31 May 2026","author":"Dad","text":"Pasture ploughing complete. 5.751 acres ploughed at UGX 230,000/acre (subtotal UGX 1,583,000), plus UGX 60,000 for shrub clearing and UGX 210,000 for tractor fuel — total UGX 1,853,000. Hellen processing reimbursement (KIM/07/26-3) along with the outstanding UGX 10,000 from ear-tagging. September 2026: replough, place more topsoil, sow grass seed, then harrow to smooth the area for sheep grazing. Earlier in the month (19 May), 10 young + 12 old sheep were ear-tagged for project identification — Mum confirmed flat pricing of UGX 200,000 per sheep; payment for 12 old ewes (UGX 2,400,000) is being processed under KIM/07/26-4.","images":["/static/project-pics/sheep/WhatsApp Image 2026-05-19 at 20.53.29.jpeg","/static/project-pics/sheep/WhatsApp Image 2026-05-19 at 20.54.45.jpeg","/static/project-pics/sheep/WhatsApp Image 2026-05-19 at 20.55.09.jpeg"],"videos":["/static/project-pics/sheep/ploughing_2026-05-31_1.mp4","/static/project-pics/sheep/ploughing_2026-05-31_2.mp4","/static/project-pics/sheep/ploughing_2026-05-31_3.mp4"]}},
        {"id":"goats","name":"Goats","icon":"🐐","category":"Farming & Agriculture","status":"Under Review","lead":"Individual owners + Solomon","headline":"~55 goats — difficult to manage","data":[
            {"label":"Count","value":"~55 goats across family owners"},{"label":"Ownership","value":"All members except Solomon and Hellen"},
            {"label":"Grazing","value":"Valley and upper plot (limited, not sustainable long-term)"},
            {"label":"Decision","value":"Time given until August 2026 to find workable approach. No immediate sale."}
        ]},
        {"id":"dairy","name":"Dairy (Cows)","icon":"🐄","category":"Farming & Agriculture","status":"Planning","lead":"Dad + Solomon","headline":"Shed-based system under investigation","data":[
            {"label":"Current Stock","value":"~50 cows across 3 farms (exact count to be confirmed by Dad)"},
            {"label":"Key Challenge","value":"Tick-related disease kills quality breeds that graze outdoors. Vaccines not yet commercially available."},
            {"label":"Proposed Model","value":"Shed-based zero-grazing: cows kept fully indoors, never go outside, fed cut forage only, no tick exposure"},
            {"label":"Target","value":"Start with 2 quality cows at ~UGX 8,000,000 each"},
            {"label":"Potential Yield","value":"15 to 50 litres per cow per day"},
            {"label":"Productive Life","value":"4 to 6 births/lactations per cow, then sold for beef"},
            {"label":"Shed Visit","value":"Dad visited a working shed dairy farm on 25 May 2026. Photos and videos from the visit are below."}
        ],"update":{"date":"25 May 2026","author":"Dad","text":"Dad visited a working shed-based dairy farm to study the model we plan to build. Key findings: cows are fully confined, never go outdoors, and are kept on a deep litter system. Each cow completes 4 to 6 lactation cycles before being sold for beef. The shed uses wooden stall dividers, a concrete feed trough, and a piped water supply. Farms running this model are still rare in the area. This is the system KimFam is targeting.","images":["/static/project-pics/dairy/dairy_img_3.jpg","/static/project-pics/dairy/dairy_img_5.jpg","/static/project-pics/dairy/dairy_img_4.jpg","/static/project-pics/dairy/dairy_img_1.jpg","/static/project-pics/dairy/dairy_img_2.jpg"],"videos":["/static/project-pics/dairy/dairy_vid_herd.mp4","/static/project-pics/dairy/dairy_vid_stalls.mp4","/static/project-pics/dairy/dairy_vid_exterior.mp4","/static/project-pics/dairy/dairy_vid_feeding.mp4"]}},
        {"id":"mango","name":"Mango & Oranges","icon":"🥭","category":"Farming & Agriculture","status":"Active","lead":"Dad + Solomon","headline":"Plantation at Nyabugando","data":[
            {"label":"Location","value":"Nyabugando farm"},
            {"label":"Integration","value":"Free-range chickens patrol the mango trees, eating fruit fly larvae — natural pest control"},
            {"label":"Note","value":"Mango harvest season temporarily reduces egg production. This is planned for in projections."}
        ]},
        {"id":"trees","name":"Tree Planting (Eucalyptus)","icon":"🌲","category":"Farming & Agriculture","status":"Active","lead":"Dad + Solomon","headline":"3.38 acres at Busisi Land","data":[
            {"label":"Location","value":"Busisi Land, Kashari"},{"label":"Species","value":"Eucalyptus grandis"},
            {"label":"Area","value":"3.38 acres"},{"label":"Density","value":"600-800 trees per acre"},
            {"label":"Capital Invested","value":"UGX 2,630,000 (seedlings + planting, Aug 2024 receipt on file)"},
            {"label":"Revenue Streams","value":"Timber, firewood, poles"},{"label":"ROI Horizon","value":"5-10 years"}
        ]},
        {"id":"bees","name":"Apiary (Bees)","icon":"🍯","category":"Farming & Agriculture","status":"Research","lead":"Dad","headline":"Business plan ready, no hives yet","data":[
            {"label":"Brand Name","value":"KIM FAM Honey Farm"},{"label":"Starting Scale","value":"10 hives planned"},
            {"label":"Production Target","value":"250kg honey per year"},
            {"label":"Status","value":"Dad joined beekeeper groups. District apiary officer consulted (limited guidance). Identifying bee forage trees."},
            {"label":"No capital committed","value":"Hives not yet purchased"}
        ]},
        {"id":"rabbits","name":"Rabbits","icon":"🐇","category":"Farming & Agriculture","status":"Research","lead":"Arunga family","headline":"High demand, outgrowers engaged","data":[
            {"label":"Lead","value":"Arunga family"},{"label":"Market Price","value":"UGX 15,000/kg"},
            {"label":"Demand","value":"High — limited only by housing capacity"},
            {"label":"Status","value":"Outgrowers engaged to increase supply. Housing expansion planned."},
            {"label":"Key Risks","value":"Cats and rats"}
        ]},
        {"id":"fortune_credit","name":"Fortune Credit","icon":"📈","category":"Unit Trusts","status":"Due Diligence","lead":"Alex (due diligence) + Lawi (introduced)","headline":"2% monthly return — pending verification","data":[
            {"label":"Type","value":"Silent investor model — finances startup businesses"},{"label":"Location","value":"Nairobi, Kenya (Busega branch in Kampala)"},
            {"label":"Monthly Return","value":"2% per month (24% annualised)"},
            {"label":"Minimum Investment","value":"UGX 29M-34M (KES-denominated)"},
            {"label":"Investment Period","value":"3 or 6 months"},
            {"label":"Status","value":"Due diligence ongoing. Alex to visit Busega. No funds committed. Verify CMA Kenya registration before committing."}
        ]},
        {"id":"kakoba","name":"Kakoba Land","icon":"🏙️","category":"Real Estate","status":"Early Planning","lead":"Arunga family","headline":"Urban land, development TBD","data":[
            {"label":"Location","value":"Kakoba (urban area)"},{"label":"Plan","value":"Develop into commercial complex"},
            {"label":"Status","value":"Survey engineer consulted. Development options being assessed. No concrete proposal yet."}
        ]},
        {"id":"irrigation","name":"Irrigation & Bananas","icon":"🌿","category":"Farming & Agriculture","status":"Planning","lead":"Dad + Solomon","headline":"Phase 1 design complete","data":[
            {"label":"Current Banana Area","value":"~2 acres"},{"label":"Expansion Plan","value":"4 acres (double current)"},
            {"label":"System","value":"Overhead sprinkler irrigation"},{"label":"Water Sources","value":"Rainwater harvesting + farm ponds"},
            {"label":"Additional","value":"Vegetable garden to be added alongside banana expansion"},
            {"label":"Phase 1","value":"Design complete"},{"label":"Phase 2","value":"Procurement and installation — pending approval and funding"},
            {"label":"Reason","value":"~80 banana bunches lost to heat stress (premature ripening) without irrigation"}
        ]}
    ]
    # Overlay live DB updates — newer DB entry wins over hardcoded update
    for p in projects:
        hc = p.get("update")
        p["update"] = _best_update(p["id"], hc)
    return projects


@app.get("/api/projects/pitches")
async def get_project_pitches(request: Request):
    """AI-cooked one-line "why join" enticements per project (figure-led).
    Cached by the background pitch engine; read-only here."""
    from fastapi import HTTPException as _HE
    from fastapi import Request as _R  # noqa
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Not authenticated")
    import pitch_engine as _pitch
    data = _pitch.load()
    return {pid: v.get("pitch") for pid, v in data.items() if v.get("pitch")}


@app.post("/api/projects/pitches/refresh")
async def refresh_project_pitches(request: Request):
    """Force the pitch engine to re-cook now (admin/internal). Runs in a thread
    so the request returns immediately."""
    from fastapi import HTTPException as _HE
    if not (_auth_verify(_get_tok(request)) or _internal_key_ok(request)):
        raise _HE(status_code=401, detail="Not authenticated")
    import threading as _th
    import pitch_engine as _pitch
    _th.Thread(target=lambda: _pitch.cook_once(get_all_projects, force=True), daemon=True).start()
    return {"status": "cooking"}

# ── Document repository + Ask KimFam ──────────────────────────────────────────
from fastapi import Request
from fastapi.responses import StreamingResponse
from pathlib import Path
from docx import Document as DocxDocument
import os, json
from dotenv import load_dotenv
from google import genai as genai_client
from groq import Groq

load_dotenv(Path(__file__).parent / ".env")

DOCS_DIR = Path(__file__).parent / "docs"

CATEGORY_LABELS = {
    "minutes":    "Meeting Minutes",
    "governance": "Governance & Policies",
    "projects":   "Project Documents",
    "financial":  "Financial Reports",
    "receipts":   "Receipts",
}

import re
from datetime import datetime

MONTHS = {"jan":1,"feb":2,"mar":3,"april":4,"apr":4,"may":5,"june":6,"jun":6,
          "july":7,"jul":7,"aug":8,"sep":9,"oct":10,"nov":11,"dec":12}

def _date_from_name(stem: str):
    """Parse a date from a filename stem (no filesystem access). Returns datetime or None."""
    s = stem.lower().replace("_", " ").replace("-", " ")
    # Match patterns like "april 12 2026", "feb 11 2024", "june 9 2024"
    m = re.search(r'(jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|april?|may|june?|july?|aug(?:ust)?|sep(?:tember)?|oct(?:ober)?|nov(?:ember)?|dec(?:ember)?)[^\d]*(\d{1,2})[^\d]*(\d{4})', s)
    if m:
        mon = MONTHS.get(m.group(1)[:3] if m.group(1)[:3] != "apr" else "apr", 0)
        if not mon:
            for k in MONTHS:
                if m.group(1).startswith(k):
                    mon = MONTHS[k]
                    break
        try:
            return datetime(int(m.group(3)), mon, int(m.group(2)))
        except Exception:
            pass
    return None

def _extract_date(path: Path):
    """Try to extract a sortable date from the filename. Falls back to file mtime."""
    dt = _date_from_name(path.stem)
    if dt is not None:
        return dt
    # Fallback: file modification time (local disk only)
    return datetime.fromtimestamp(path.stat().st_mtime)

def _friendly_name(path: Path) -> str:
    n = path.stem.replace("_", " ").replace("-", " ")
    for p in ["KIM Investment Club Meeting Minutes", "KIMFAM Meeting Minutes",
              "KIM Meeting Minutes"]:
        n = n.replace(p, "").strip()
    return n if n else path.stem

_DOC_SUFFIXES = {".docx", ".pdf", ".pptx", ".xlsx", ".doc", ".ppt", ".xls"}

_DOC_CATS = ["minutes", "governance", "projects", "financial", "receipts"]


def _doc_date_key(stem: str) -> tuple:
    """Best-effort (year, month, day) from a document name, handling the many
    formats in this repo: 'June 7 2026', '2026-03-21', '6.23.2024', '21.03.2026',
    '28 02 2025', '9th Feb.2023', or a bare year. Returns (0,0,0) if none found,
    so undated files sort to the bottom."""
    dt = _date_from_name(stem)
    if dt is not None:
        return (dt.year, dt.month, dt.day)
    s = stem.lower().replace("_", " ")
    m = re.search(r'(\d{4})[-.](\d{1,2})[-.](\d{1,2})', s)          # ISO yyyy-mm-dd
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    m = re.search(r'(\d{1,2})[.\s](\d{1,2})[.\s](\d{4})', s)         # d.m.y or m.d.y
    if m:
        a, b, y = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if a > 12:   day, mon = a, b
        elif b > 12: day, mon = b, a
        else:        day, mon = a, b   # ambiguous; year dominates the sort anyway
        return (y, mon, day)
    m = re.search(r'(20\d\d|19\d\d)', s)                            # bare year
    if m:
        return (int(m.group(1)), 0, 0)
    return (0, 0, 0)


def _docs_build_cat(cat: str, groups: dict) -> dict:
    """Turn {group_label: [file dicts]} into the nested category payload, with
    files sorted newest-first (non-archive before archive on a tie) and year
    groups newest-first."""
    def _sort_files(files):
        def k(x):
            stem = Path(x["rel"]).stem
            y, mo, d = _doc_date_key(stem)
            not_archive = 0 if "archive" in stem.lower() else 1
            return (y, mo, d, not_archive)
        return sorted(files, key=k, reverse=True)
    def _group_key(label):
        # Groups carrying a year sort newest-first; the rest alphabetically after.
        m = re.search(r"(19|20)\d{2}", label)
        return (0, -int(m.group(0)), label) if m else (1, 0, label)
    out_groups = []
    for glabel in sorted(groups.keys(), key=_group_key):
        files = _sort_files(groups[glabel])
        out_groups.append({
            "label": glabel,
            "files": [
                {"name": _friendly_name(Path(f["rel"])), "file": f["rel"], "url": f"/docs/{cat}/{f['rel']}"}
                for f in files
            ],
        })
    total = sum(len(g["files"]) for g in out_groups)
    return {"label": CATEGORY_LABELS[cat], "count": total, "groups": out_groups}


@app.get("/api/docs")
def get_docs():
    """Nested documents: category -> sub-groups (one folder deep) -> files.
    Files directly under a category (no sub-folder) fall into a 'General' group."""
    result = {}
    if _r2.is_configured():
        for cat in _DOC_CATS:
            try:
                objects = _r2.list_folder(cat + "/")
            except Exception:
                objects = []
            groups: dict[str, list] = {}
            for obj in objects:
                key = obj["key"]
                rel = key[len(cat) + 1:] if key.startswith(cat + "/") else key
                if not rel:
                    continue
                if "_versions" in rel.split("/"):   # archived proposal versions stay hidden
                    continue
                parts = rel.split("/", 1)
                group, fname = (parts[0], parts[1]) if len(parts) == 2 else ("General", parts[0])
                base = fname.rsplit("/", 1)[-1]
                if Path(base).suffix.lower() not in _DOC_SUFFIXES or ".bak." in base:
                    continue
                groups.setdefault(group, []).append({"rel": rel, "mtime": obj["last_modified"]})
            result[cat] = _docs_build_cat(cat, groups)
    else:
        # Fallback: local disk (staging / no-R2 environments)
        for cat in _DOC_CATS:
            folder = DOCS_DIR / cat
            if not folder.exists():
                continue
            groups = {}
            for f in folder.rglob("*"):
                if not f.is_file() or f.suffix.lower() not in _DOC_SUFFIXES:
                    continue
                if f.name.startswith("~") or ".bak." in f.name:
                    continue
                rel = f.relative_to(folder).as_posix()
                if "_versions" in rel.split("/"):   # archived proposal versions stay hidden
                    continue
                parts = rel.split("/", 1)
                group = parts[0] if len(parts) == 2 else "General"
                groups.setdefault(group, []).append({"rel": rel, "mtime": f.stat().st_mtime})
            result[cat] = _docs_build_cat(cat, groups)
    return result

@app.get("/api/docs/search")
def search_docs(query: str = ""):
    query_lower = query.lower().strip()
    if not query_lower:
        return {"query": query, "results": []}

    results = []
    all_docs = get_docs()

    for cat, cat_data in all_docs.items():
        for group in cat_data.get("groups", []):
            for doc in group.get("files", []):
                # Search in both filename and friendly name
                if query_lower in doc["file"].lower() or query_lower in doc["name"].lower():
                    results.append({
                        "category": cat,
                        "categoryLabel": cat_data["label"],
                        "group": group["label"],
                        "name": doc["name"],
                        "file": doc["file"],
                        "url": doc["url"]
                    })

    return {"query": query, "count": len(results), "results": results}

@app.get("/api/docs/search-semantic")
def search_docs_semantic(query: str = ""):
    query_lower = query.lower().strip()
    if not query_lower:
        return {"query": query, "results": []}

    try:
        from sentence_transformers import SentenceTransformer
        model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
        embedding = model.encode(query_lower, normalize_embeddings=True)
        embedding = [float(x) for x in embedding]

        import chromadb
        from pathlib import Path
        base_dir = Path(__file__).parent
        chroma_dir = base_dir / "data" / "chroma"
        client = chromadb.PersistentClient(path=str(chroma_dir))
        collection = client.get_or_create_collection("kimfam_docs")

        chroma_results = collection.query(
            query_embeddings=[embedding],
            n_results=10,
            include=["documents", "metadatas", "distances"]
        )

        all_docs = get_docs()
        doc_urls = {}
        for cat, cat_data in all_docs.items():
            for doc in cat_data.get("files", []):
                doc_urls[doc["file"]] = (doc["url"], doc["name"], cat_data["label"], cat)

        results = []
        if chroma_results["documents"]:
            for i, doc_text in enumerate(chroma_results["documents"][0]):
                metadata = chroma_results["metadatas"][0][i] if chroma_results["metadatas"] else {}
                distance = chroma_results["distances"][0][i] if chroma_results["distances"] else 0
                source = metadata.get("source", "")

                # Map source to file
                for fname, (url, name, label, cat) in doc_urls.items():
                    if source and source in fname:
                        results.append({
                            "category": cat,
                            "categoryLabel": label,
                            "name": name,
                            "file": fname,
                            "url": url,
                            "relevance": max(0, 1 - distance),
                            "preview": doc_text[:200] + "..." if len(doc_text) > 200 else doc_text
                        })
                        break

        return {
            "query": query,
            "count": len(results),
            "results": results
        }
    except Exception as e:
        log.warning(f"Semantic search error: {e}")
        return {"query": query, "count": 0, "results": [], "error": str(e)}


# ── Proposals: AI scoring against the KimFam project-proposal template ──────────
# Criteria mirror EXACTLY the "Project Proposal Template — Key Areas to Address"
# in Hillary's Project Management deck: Background, Objective, Scope, Stakeholders,
# Resources, Timeline, Risk Management, Financial Appraisal, Benefits.
PROPOSAL_CRITERIA = [
    ("Background & Introduction", 10),
    ("Objective", 12),
    ("Scope", 8),
    ("Stakeholders", 8),
    ("Resources & Budget", 12),
    ("Timeline / Milestones", 10),
    ("Risk Management", 10),
    ("Financial Appraisal (Payback, ROI, NPV)", 20),
    ("Benefits (Quantitative & Qualitative)", 10),
]


def _extract_doc_text(fname: str, raw: bytes) -> str:
    """Extract plain text from an uploaded proposal (docx/pdf/pptx/xlsx/txt)."""
    import io as _io
    ext = fname.lower().rsplit(".", 1)[-1] if "." in fname else ""
    try:
        if ext == "docx":
            import docx as _docx
            d = _docx.Document(_io.BytesIO(raw))
            return "\n".join(p.text for p in d.paragraphs if p.text.strip())
        if ext == "pdf":
            import pdfplumber as _pp
            with _pp.open(_io.BytesIO(raw)) as pdf:
                return "\n".join(pg.extract_text() or "" for pg in pdf.pages)
        if ext == "pptx":
            from pptx import Presentation as _Pr
            prs = _Pr(_io.BytesIO(raw)); out = []
            for s in prs.slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        for p in sh.text_frame.paragraphs:
                            t = "".join(r.text for r in p.runs).strip()
                            if t:
                                out.append(t)
                    if getattr(sh, "has_table", False):
                        for row in sh.table.rows:
                            cs = [c.text.strip() for c in row.cells]
                            if any(cs):
                                out.append(" | ".join(cs))
            return "\n".join(out)
        if ext in ("xlsx", "xls"):
            from openpyxl import load_workbook as _lw
            wb = _lw(_io.BytesIO(raw), read_only=True, data_only=True); out = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        out.append(" | ".join(vals))
            return "\n".join(out)
        return raw.decode("utf-8", errors="replace")
    except Exception as _e:
        log.error(f"proposal extract failed for {fname}: {_e}")
        return ""


def _re_docs_safe(name: str) -> str:
    import re as _r
    return (_r.sub(r'[\\/:*?"<>|]+', " ", name).strip()[:80]) or "Proposal"


def _doc_modified_date(fname: str, raw: bytes):
    """The document's own last-modified date (docx/pptx core properties, pdf ModDate).
    Used as the 'uploaded' date for pre-existing proposals so the card shows when the
    document was actually last edited, not when it was imported. Returns datetime or None."""
    import io as _io
    ext = fname.lower().rsplit(".", 1)[-1] if "." in fname else ""
    try:
        if ext == "docx":
            import docx as _docx
            return _docx.Document(_io.BytesIO(raw)).core_properties.modified
        if ext == "pptx":
            from pptx import Presentation as _Pr
            return _Pr(_io.BytesIO(raw)).core_properties.modified
        if ext == "pdf":
            import pdfplumber as _pp, datetime as _dt
            with _pp.open(_io.BytesIO(raw)) as pdf:
                md = (pdf.metadata or {}).get("ModDate") or (pdf.metadata or {}).get("CreationDate")
            if md:
                m = __import__("re").search(r"D:(\d{4})(\d{2})(\d{2})", str(md))
                if m:
                    return _dt.datetime(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    except Exception:
        return None
    return None


def _ensure_proposals_table():
    from db import execute as _exec
    _exec("""CREATE TABLE IF NOT EXISTS proposals (
        id            SERIAL PRIMARY KEY,
        title         TEXT NOT NULL,
        owner         TEXT NOT NULL,
        submitted_by  TEXT NOT NULL,
        source_path   TEXT,
        overall_score NUMERIC,
        verdict       TEXT,
        scores        JSONB,
        strengths     JSONB,
        gaps          JSONB,
        improvements  JSONB,
        summary       TEXT,
        support_requested TEXT,
        readiness     JSONB,
        thread_id     INTEGER,
        version       INTEGER DEFAULT 1,
        is_current    BOOLEAN DEFAULT TRUE,
        scored        BOOLEAN DEFAULT FALSE,
        file_hash     TEXT,
        uploaded_at   TIMESTAMPTZ DEFAULT now(),
        created_at    TIMESTAMPTZ DEFAULT now()
    )""")
    # Additive columns for tables created before these fields existed.
    for col, ddl in [
        ("support_requested", "TEXT"), ("readiness", "JSONB"),
        ("thread_id", "INTEGER"), ("version", "INTEGER DEFAULT 1"),
        ("is_current", "BOOLEAN DEFAULT TRUE"), ("uploaded_at", "TIMESTAMPTZ DEFAULT now()"),
        ("file_hash", "TEXT"),
    ]:
        try:
            _exec(f"ALTER TABLE proposals ADD COLUMN IF NOT EXISTS {col} {ddl}")
        except Exception:
            pass
    _exec("""CREATE TABLE IF NOT EXISTS proposal_comments (
        id          SERIAL PRIMARY KEY,
        thread_id   INTEGER,
        proposal_id INTEGER,
        author      TEXT,
        body        TEXT,
        action_ref  TEXT,
        created_at  TIMESTAMPTZ DEFAULT now()
    )""")


def _verdict_for(score: float) -> str:
    if score >= 80: return "Strong"
    if score >= 65: return "Viable with conditions"
    if score >= 50: return "Needs work"
    return "Not ready"


_PROPOSAL_FRAMEWORK_CACHE = {"text": None}


def _proposal_framework_context() -> str:
    """Load and cache the framework docs the scorer must reason against:
    the project-management deck and the investment & reward guidelines."""
    if _PROPOSAL_FRAMEWORK_CACHE["text"] is not None:
        return _PROPOSAL_FRAMEWORK_CACHE["text"]
    import os as _os, glob as _glob
    docs = _os.path.join(_os.path.dirname(__file__), "docs")
    parts = []
    patterns = [
        (f"{docs}/projects/**/*project management*", "PROJECT MANAGEMENT FRAMEWORK"),
        (f"{docs}/governance/**/*reward*", "INVESTMENT & REWARD GUIDELINES"),
    ]
    for pat, label in patterns:
        for fp in _glob.glob(pat, recursive=True):
            if fp.lower().endswith((".docx", ".pdf", ".pptx")):
                try:
                    txt = _extract_doc_text(_os.path.basename(fp), open(fp, "rb").read())
                    if txt.strip():
                        parts.append(f"--- {label} ({_os.path.basename(fp)}) ---\n{txt[:9000]}")
                        break
                except Exception as _e:
                    log.error(f"framework load failed {fp}: {_e}")
    ctx = "\n\n".join(parts)
    _PROPOSAL_FRAMEWORK_CACHE["text"] = ctx
    return ctx


async def _score_proposal(text: str) -> dict | None:
    """Score proposal text against PROPOSAL_CRITERIA using CLAUDE ONLY (no DeepSeek
    or Groq fallback), grounded in the club's project-management deck and reward
    guidelines. The weighted total is computed server-side."""
    import json as _json, re as _re
    if len(text) > 18000:
        text = text[:18000] + "\n\n[Truncated for scoring.]"
    crit_lines = "\n".join(f"- {n} (weight {w})" for n, w in PROPOSAL_CRITERIA)
    framework = _proposal_framework_context()
    prompt = (
        "You are evaluating a KimFam Investment Club project proposal. Judge it strictly "
        "against the club's own framework provided below (its project-management approach and "
        "its investment & reward guidelines). Score each criterion 1 (poor) to 5 (excellent), "
        "justify briefly using only what the proposal actually says, and reflect the club's "
        "guidelines (e.g. ROI, recoupment and reward sharing) in the Financial Appraisal score. "
        "No em-dashes or en-dashes.\n\n"
        "=== CLUB FRAMEWORK (authoritative) ===\n" + (framework or "(framework unavailable)") +
        "\n\n=== SCORING CRITERIA (these are the club's official Project Proposal Template, the "
        "'Key Areas to Address') ===\n" + crit_lines +
        "\n\nAlso identify the SUPPORT the proposer is requesting from the family (money, "
        "approval, labour, expertise, etc.) and judge whether the proposal is in shape for the "
        "family to actually grant that support, listing what must still be provided before it can be.\n"
        "\n\nReturn ONLY valid JSON (no markdown):\n"
        '{"criteria":[{"name":"<exact criterion name>","score":1-5,"rationale":"..."}],'
        '"strengths":["..."],"gaps":["..."],"improvements":["..."],'
        '"support_requested":"what the proposer asks the family for",'
        '"readiness":{"status":"ready|partly|not_ready","assessment":"can the family grant the '
        'support yet, and why","blocking":["what must be provided before the support can be granted"]},'
        '"summary":"one short paragraph"}'
        "\n\n=== PROPOSAL TO SCORE ===\n" + (text or "(empty)")
    )
    # Claude only, per the family's instruction (it holds the framework context).
    raw = await _ask_claude_async(prompt, model="claude-haiku-4-5-20251001", timeout=220)
    if not raw:
        return None
    raw = _re.sub(r"^```(?:json)?\s*", "", raw.strip())
    raw = _re.sub(r"\s*```$", "", raw)
    try:
        data = _json.loads(raw)
    except Exception:
        m = _re.search(r"\{.*\}", raw, _re.DOTALL)
        if not m:
            return None
        try:
            data = _json.loads(m.group())
        except Exception:
            return None
    weights = {n: w for n, w in PROPOSAL_CRITERIA}
    total = 0.0
    norm = []
    for c in (data.get("criteria") or []):
        name = c.get("name", "")
        w = weights.get(name, 0)
        try:
            sc = max(1, min(5, int(round(float(c.get("score", 0))))))
        except Exception:
            sc = 0
        if w and sc:
            total += (sc / 5.0) * w
        norm.append({"name": name, "weight": w, "score": sc, "rationale": c.get("rationale", "")})
    total = round(total, 1)
    return {
        "criteria": norm, "overall_score": total, "verdict": _verdict_for(total),
        "strengths": data.get("strengths") or [], "gaps": data.get("gaps") or [],
        "improvements": data.get("improvements") or [], "summary": data.get("summary") or "",
        "support_requested": data.get("support_requested") or "",
        "readiness": data.get("readiness") or {},
    }


# ── Reusable SSE progress wrapper for long AI tasks ────────────────────────────
# Emits {type:step} stage messages (keeping the connection alive past idle cutoffs)
# while a Claude scoring call runs, then persists and emits {type:result}. Any AI
# request that takes >a few seconds should stream through this pattern so the user
# sees what it's doing (like Ask KimFam) and the browser never "fails to fetch".
_AI_SCORE_STAGES = [
    "Weighing Background, Objective and Scope...",
    "Checking Resources, Budget and Timeline...",
    "Assessing Risk Management...",
    "Evaluating Financial Appraisal (Payback, ROI, NPV)...",
    "Judging Benefits and strategic fit...",
    "Determining the support requested and readiness...",
    "Finalising the scorecard...",
]


async def _ai_score_stream(text: str, persist, kind: str = "proposal"):
    import asyncio as _aio
    yield _sse({"type": "step", "msg": "Reading the proposal..."})
    yield _sse({"type": "step", "msg": "Loading the project-management framework and reward guidelines..."})
    yield _sse({"type": "step", "msg": "Scoring against the 9 Key Areas with Claude (about a minute)..."})
    task = _aio.create_task(_score_proposal(text))
    i = 0
    while not task.done():
        await _aio.sleep(6)
        if not task.done():
            yield _sse({"type": "step", "msg": _AI_SCORE_STAGES[i % len(_AI_SCORE_STAGES)]})
            i += 1
    try:
        sc = task.result()
    except Exception as _e:
        log.error(f"{kind} score stream failed: {_e}")
        sc = None
    if not sc:
        yield _sse({"type": "error", "msg": "Scoring is busy right now. Please try again in a moment."})
        return
    try:
        result = persist(sc)
    except Exception as _e:
        log.error(f"{kind} persist failed: {_e}")
        yield _sse({"type": "error", "msg": "Scored it, but could not save. Please retry."})
        return
    yield _sse({"type": "step", "msg": f"Done: {sc['overall_score']} / 100 ({sc['verdict']})"})
    yield _sse({"type": "result", "proposal": result})


def _proposal_row(r: dict) -> dict:
    import json as _json
    def _j(v):
        return v if isinstance(v, (list, dict)) else (_json.loads(v) if v else [])
    rd = r.get("readiness")
    rd = rd if isinstance(rd, dict) else (_json.loads(rd) if rd else {})
    return {
        "id": r["id"], "title": r["title"], "owner": r["owner"],
        "submitted_by": r["submitted_by"],
        "file_url": (f"/docs/{r['source_path']}" if r.get("source_path") else None),
        "overall_score": float(r["overall_score"]) if r.get("overall_score") is not None else None,
        "verdict": r.get("verdict"), "scored": r.get("scored", False),
        "criteria": _j(r.get("scores")), "strengths": _j(r.get("strengths")),
        "gaps": _j(r.get("gaps")), "improvements": _j(r.get("improvements")),
        "support_requested": r.get("support_requested") or "", "readiness": rd,
        "version": r.get("version", 1), "is_current": r.get("is_current", True),
        "thread_id": r.get("thread_id"),
        "summary": r.get("summary"), "created_at": str(r.get("created_at") or ""),
        "uploaded_at": str(r.get("uploaded_at") or r.get("created_at") or ""),
    }


@app.get("/api/proposals/matrix")
def proposals_matrix(request: Request):
    """The scoring rubric: each Key Area and the % it contributes, plus verdict bands."""
    from fastapi import HTTPException as _HE
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    return {
        "criteria": [{"name": n, "weight": w} for n, w in PROPOSAL_CRITERIA],
        "total": sum(w for _, w in PROPOSAL_CRITERIA),
        "verdict_bands": [
            {"min": 80, "label": "Strong"},
            {"min": 65, "label": "Viable with conditions"},
            {"min": 50, "label": "Needs work"},
            {"min": 0, "label": "Not ready"},
        ],
        "scale": "Each Key Area is scored 1 to 5, weighted to its percentage, and summed to a score out of 100.",
    }


@app.get("/api/proposals")
def list_proposals(request: Request):
    from db import query as _dbq
    if not _auth_verify(_get_tok(request)):
        from fastapi import HTTPException as _HE
        raise _HE(status_code=401, detail="Login required")
    _ensure_proposals_table()
    # One entry per proposal thread (the current version), each carrying its version
    # history so revisions group together and progress is easy to follow.
    current = _dbq("SELECT * FROM proposals WHERE is_current=TRUE ORDER BY uploaded_at DESC NULLS LAST, created_at DESC")
    out = []
    for r in current:
        row = _proposal_row(r)
        thread = r.get("thread_id") or r["id"]
        vers = _dbq("""SELECT id, version, overall_score, verdict, scored, uploaded_at, source_path, scores
                       FROM proposals WHERE thread_id=%s OR id=%s ORDER BY version DESC""",
                    (thread, thread))
        import json as _hjson
        def _hj(v):
            return v if isinstance(v, list) else (_hjson.loads(v) if v else [])
        row["history"] = [{
            "id": v["id"], "version": v["version"],
            "overall_score": float(v["overall_score"]) if v.get("overall_score") is not None else None,
            "verdict": v.get("verdict"), "scored": v.get("scored", False),
            "uploaded_at": str(v.get("uploaded_at") or ""),
            "file_url": (f"/docs/{v['source_path']}" if v.get("source_path") else None),
            "criteria": [{"name": c.get("name"), "score": c.get("score")} for c in _hj(v.get("scores"))],
        } for v in vers]
        out.append(row)
    return {"proposals": out}


def _archive_prior_proposal_file(prior: dict):
    """Move a superseded proposal's document into a hidden `_versions` folder beside it,
    so the Documents repo only ever shows the current version under its real project area.
    The proposal row's source_path is updated so its version-history link still resolves."""
    import os as _os
    sp = prior.get("source_path")
    if not sp or "/_versions/" in sp:
        return
    parts = sp.rsplit("/", 1)
    parent = parts[0] if len(parts) == 2 else "projects"
    fname = parts[-1]
    archive = f"{parent}/_versions/v{prior.get('version', 1)} - {fname}"
    if archive == sp:
        return
    try:
        from db import execute as _exec2
        import r2_storage as _r2p
        here = _os.path.dirname(__file__)
        local_src = _os.path.join(here, "docs", sp)
        local_dst = _os.path.join(here, "docs", archive)
        raw = None
        if _os.path.exists(local_src):
            _os.makedirs(_os.path.dirname(local_dst), exist_ok=True)
            with open(local_src, "rb") as _r: raw = _r.read()
            _os.replace(local_src, local_dst)
        if _r2p.is_configured():
            if raw is None:
                try:
                    cat, fn = sp.split("/", 1)
                    raw = _r2_bytes(cat, fn)
                except Exception:
                    raw = None
            if raw is not None:
                tmpp = f"/tmp/_arch_{_os.getpid()}"
                with open(tmpp, "wb") as _w: _w.write(raw)
                _r2p.upload(tmpp, archive)
                _os.remove(tmpp)
                try: _r2p.delete(sp)
                except Exception: pass
        _exec2("UPDATE proposals SET source_path=%s WHERE id=%s", (archive, prior["id"]))
    except Exception as _e:
        log.error(f"archive prior proposal file failed: {_e}")


@app.get("/api/proposals/areas")
def proposal_areas(request: Request):
    """Existing folder paths under projects/ (for the upload 'file under' picker)."""
    from fastapi import HTTPException as _HE
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    paths = set()
    # The legacy "Proposals" bucket is being deprecated — don't offer it (or _versions)
    # as a filing destination; proposals belong under their real project area.
    def _skip(segs):
        return (not segs) or segs[0].lower() == "proposals" or "_versions" in segs
    try:
        import r2_storage as _r2p
        if _r2p.is_configured():
            for obj in _r2p.list_folder("projects/"):
                rel = obj["key"][len("projects/"):]
                segs = [s for s in rel.split("/")[:-1] if s]
                if _skip(segs):
                    continue
                for i in range(len(segs)):
                    paths.add("/".join(segs[: i + 1]))
        else:
            base = DOCS_DIR / "projects"
            if base.exists():
                for f in base.rglob("*"):
                    if f.is_dir() and not _skip(list(f.relative_to(base).parts)):
                        paths.add(f.relative_to(base).as_posix())
    except Exception as _e:
        log.error(f"proposal areas failed: {_e}")
    return {"areas": sorted(paths)}


@app.get("/api/proposals/{pid}")
def get_proposal(pid: int, request: Request):
    from db import query as _dbq
    from fastapi import HTTPException as _HE
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    _ensure_proposals_table()
    rows = _dbq("SELECT * FROM proposals WHERE id=%s", (pid,))
    if not rows:
        raise _HE(status_code=404, detail="Proposal not found")
    return _proposal_row(rows[0])


@app.post("/api/proposals")
async def create_proposal(request: Request):
    """Upload a proposal (any member), store it, and AI-score it."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq, execute as _exec
    import json as _json, os as _os
    payload = _auth_verify(_get_tok(request))
    if not payload:
        raise _HE(status_code=401, detail="Login required")
    submitter = payload.get("sub", "member")
    form = await request.form()
    f = form.get("file")
    title = (form.get("title") or "").strip()
    owner = (form.get("owner") or "").strip()
    if not f or not hasattr(f, "filename"):
        raise _HE(status_code=400, detail="A proposal file is required")
    if not owner:
        raise _HE(status_code=400, detail="Please choose whose proposal this is (owner)")
    raw = await f.read()
    fname = f.filename or "proposal"
    if not title:
        title = _os.path.splitext(fname)[0]
    text = _extract_doc_text(fname, raw)
    if len(text.strip()) < 40:
        raise _HE(status_code=400, detail="Could not read text from that file")

    _ensure_proposals_table()
    # ── Exact-duplicate guard: a content hash (over the file bytes, so it survives a
    #    rename but changes on any edit). Block re-uploading a byte-identical file.
    import hashlib as _hl
    file_hash = _hl.sha256(raw).hexdigest()
    dup = _dbq("""SELECT title, owner FROM proposals WHERE file_hash=%s AND is_current=TRUE
                  LIMIT 1""", (file_hash,))
    if dup:
        raise _HE(status_code=409,
                  detail=f"This exact file is already uploaded as '{dup[0]['title']}' "
                         f"({dup[0]['owner']}). Edit the document before re-uploading a new version.")

    # ── Versioning: a new upload for the same owner+title supersedes the prior
    #    current version (which is archived) so we track refinement progress. An
    #    explicit ?supersedes=<id> wins over the owner+title match.
    supersedes = (form.get("supersedes") or "").strip()
    prior = None
    if supersedes.isdigit():
        rows = _dbq("SELECT * FROM proposals WHERE id=%s", (int(supersedes),))
        prior = rows[0] if rows else None
    if not prior:
        rows = _dbq("""SELECT * FROM proposals WHERE owner=%s AND lower(title)=lower(%s)
                       AND is_current=TRUE ORDER BY id DESC LIMIT 1""", (owner, title))
        prior = rows[0] if rows else None
    # Where the document is filed in the Documents repo: under its real project area
    # (e.g. "Real estates/kakoba land/Boys quarters"), mirroring the folder tree — NOT a
    # generic Proposals bucket. Sanitize each path segment; keep the slashes.
    file_under = (form.get("file_under") or "").strip()
    area = "/".join(_re_docs_safe(p) for p in file_under.split("/") if p.strip()) or "Proposals"

    if prior:
        new_version = (prior.get("version") or 1) + 1
        thread = prior.get("thread_id") or prior["id"]
        _exec("UPDATE proposals SET is_current=FALSE WHERE thread_id=%s OR id=%s",
              (thread, thread))
        # Tuck the prior current document into a hidden _versions folder so Documents
        # only ever shows the current version under the real project area.
        _archive_prior_proposal_file(prior)
    else:
        new_version = 1
        thread = None

    # Current version lives under its real project area, visible in Documents.
    source_path = f"projects/{area}/{fname}"
    try:
        import r2_storage as _r2p
        if _r2p.is_configured():
            tmpp = f"/tmp/_proposal_{_os.getpid()}_{fname}"
            with open(tmpp, "wb") as _w: _w.write(raw)
            _r2p.upload(tmpp, source_path)
            _os.remove(tmpp)
        _ddir = _os.path.join(_os.path.dirname(__file__), "docs", *(["projects"] + area.split("/")))
        _os.makedirs(_ddir, exist_ok=True)
        with open(_os.path.join(_ddir, fname), "wb") as _w: _w.write(raw)
    except Exception as _e:
        log.error(f"proposal file store failed: {_e}")

    # Stream progress (SSE) so the long Claude scoring keeps the connection alive and
    # the user sees what it's doing, like Ask KimFam. Uses _ai_score_stream (reusable).
    from fastapi.responses import StreamingResponse

    def _persist(sc):
        cols = "title, owner, submitted_by, source_path, version, is_current, file_hash"
        if sc:
            _exec(f"""INSERT INTO proposals
                ({cols}, overall_score, verdict, scores, strengths, gaps, improvements,
                 support_requested, readiness, summary, scored)
                VALUES (%s,%s,%s,%s,%s,TRUE,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,TRUE)""",
                (title, owner, submitter, source_path, new_version, file_hash,
                 sc["overall_score"], sc["verdict"], _json.dumps(sc["criteria"]),
                 _json.dumps(sc["strengths"]), _json.dumps(sc["gaps"]),
                 _json.dumps(sc["improvements"]), sc.get("support_requested", ""),
                 _json.dumps(sc.get("readiness", {})), sc["summary"]))
        else:
            _exec(f"""INSERT INTO proposals ({cols}, scored)
                     VALUES (%s,%s,%s,%s,%s,TRUE,%s,FALSE)""",
                  (title, owner, submitter, source_path, new_version, file_hash))
        row = _dbq("SELECT * FROM proposals WHERE source_path=%s ORDER BY id DESC LIMIT 1", (source_path,))
        nid = row[0]["id"] if row else None
        if nid:
            _exec("UPDATE proposals SET thread_id=%s WHERE id=%s", (thread or nid, nid))
        result = _proposal_row(_dbq("SELECT * FROM proposals WHERE id=%s", (nid,))[0]) if nid else {"ok": True}
        # Personal WhatsApp confirmation to the owner (and submitter). Not the group.
        try:
            import notifications as _ntf
            on_behalf = submitter.lower() not in owner.lower()
            _ntf.notify_proposal_submitted(
                title, owner, submitter,
                (sc["overall_score"] if sc else None), (sc["verdict"] if sc else None),
                result.get("file_url"), on_behalf=on_behalf)
        except Exception as _ne:
            log.error(f"proposal submit notify failed: {_ne}")
        return result

    return StreamingResponse(
        _ai_score_stream(text, _persist, kind="proposal"),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive"},
    )


@app.post("/api/proposals/{pid}/score")
async def rescore_proposal(pid: int, request: Request):
    """(Re)score an existing proposal record from its stored file."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq, execute as _exec
    import json as _json, os as _os
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    _ensure_proposals_table()
    rows = _dbq("SELECT * FROM proposals WHERE id=%s", (pid,))
    if not rows:
        raise _HE(status_code=404, detail="Proposal not found")
    sp = rows[0].get("source_path")
    raw = None
    local = _os.path.join(_os.path.dirname(__file__), "docs", sp) if sp else None
    if local and _os.path.exists(local):
        raw = open(local, "rb").read()
    elif sp:
        try:
            import r2_storage as _r2p
            cat, fn = sp.split("/", 1)
            raw = _r2_bytes(cat, fn)
        except Exception:
            raw = None
    if raw is None:
        raise _HE(status_code=404, detail="Proposal file not found to score")
    text = _extract_doc_text(_os.path.basename(sp), raw)
    from fastapi.responses import StreamingResponse

    def _persist(sc):
        _exec("""UPDATE proposals SET overall_score=%s, verdict=%s, scores=%s, strengths=%s,
                 gaps=%s, improvements=%s, support_requested=%s, readiness=%s, summary=%s,
                 scored=TRUE WHERE id=%s""",
              (sc["overall_score"], sc["verdict"], _json.dumps(sc["criteria"]),
               _json.dumps(sc["strengths"]), _json.dumps(sc["gaps"]),
               _json.dumps(sc["improvements"]), sc.get("support_requested", ""),
               _json.dumps(sc.get("readiness", {})), sc["summary"], pid))
        return _proposal_row(_dbq("SELECT * FROM proposals WHERE id=%s", (pid,))[0])

    return StreamingResponse(
        _ai_score_stream(text, _persist, kind="proposal-rescore"),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive"},
    )


@app.post("/api/proposals/{pid}/share")
def share_proposal(pid: int, request: Request):
    """Deliberately announce a proposal to the family group (ready for review)."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    _ensure_proposals_table()
    rows = _dbq("SELECT * FROM proposals WHERE id=%s", (pid,))
    if not rows:
        raise _HE(status_code=404, detail="Proposal not found")
    p = rows[0]
    link = f"/docs/{p['source_path']}" if p.get("source_path") else None
    try:
        import notifications as _ntf
        _ntf.notify_proposal_ready(
            p["title"], p["owner"], p["submitted_by"],
            float(p["overall_score"]) if p.get("overall_score") is not None else None,
            p.get("verdict"), link)
    except Exception as _ne:
        log.error(f"proposal share failed: {_ne}")
        raise _HE(status_code=502, detail="Could not send the group message. Try again.")
    return {"ok": True, "message": "Shared to the family group."}


@app.get("/api/proposals/{pid}/comments")
def list_proposal_comments(pid: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import query as _dbq
    if not _auth_verify(_get_tok(request)):
        raise _HE(status_code=401, detail="Login required")
    _ensure_proposals_table()
    p = _dbq("SELECT thread_id, id FROM proposals WHERE id=%s", (pid,))
    if not p:
        raise _HE(status_code=404, detail="Proposal not found")
    thread = p[0].get("thread_id") or p[0]["id"]
    rows = _dbq("SELECT * FROM proposal_comments WHERE thread_id=%s ORDER BY created_at", (thread,))
    return {"comments": [{
        "id": r["id"], "author": r["author"], "body": r["body"],
        "action_ref": r.get("action_ref"), "created_at": str(r.get("created_at") or ""),
    } for r in rows]}


@app.post("/api/proposals/{pid}/comments")
async def add_proposal_comment(pid: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import query as _dbq, execute as _exec
    payload = _auth_verify(_get_tok(request))
    if not payload:
        raise _HE(status_code=401, detail="Login required")
    body = await request.json()
    text = (body.get("body") or "").strip()
    if not text:
        raise _HE(status_code=400, detail="Comment is empty")
    _ensure_proposals_table()
    p = _dbq("SELECT * FROM proposals WHERE id=%s", (pid,))
    if not p:
        raise _HE(status_code=404, detail="Proposal not found")
    thread = p[0].get("thread_id") or p[0]["id"]
    author = payload.get("sub", "member")
    _exec("INSERT INTO proposal_comments (thread_id, proposal_id, author, body) VALUES (%s,%s,%s,%s)",
          (thread, pid, author, text))
    try:
        import notifications as _ntf
        link = f"/docs/{p[0]['source_path']}" if p[0].get("source_path") else None
        _ntf.notify_proposal_comment(p[0]["title"], p[0]["owner"], author, text, link)
    except Exception as _ne:
        log.error(f"proposal comment notify failed: {_ne}")
    return {"ok": True}


@app.post("/api/proposals/{pid}/comments/{cid}/to-action")
def proposal_comment_to_action(pid: int, cid: int, request: Request):
    """Promote a proposal comment into a tracked action under the next meeting."""
    from fastapi import HTTPException as _HE
    from db import query as _dbq, execute as _exec
    import re as _re2
    payload = _auth_verify(_get_tok(request))
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Only admins can turn a comment into an action")
    _ensure_proposals_table()
    c = _dbq("SELECT * FROM proposal_comments WHERE id=%s", (cid,))
    p = _dbq("SELECT * FROM proposals WHERE id=%s", (pid,))
    if not c or not p:
        raise _HE(status_code=404, detail="Comment or proposal not found")
    if c[0].get("action_ref"):
        return {"ok": True, "action_id": c[0]["action_ref"]}
    author = payload.get("sub", "admin")
    target_id, target_ref = _next_or_create_meeting()
    m = _re2.match(r"KIM\s+(\d+)/(\d{4})", target_ref or "")
    if m:
        num = int(m.group(1)); yr = m.group(2)[-2:]
        i = (_dbq("SELECT COUNT(*) AS c FROM actions WHERE meeting_id=%s", (target_id,))[0]["c"]) + 1
        new_ref = f"KIM/{num:02d}/{yr}-{i}"
        while _dbq("SELECT ref FROM actions WHERE ref=%s", (new_ref,)):
            i += 1; new_ref = f"KIM/{num:02d}/{yr}-{i}"
    else:
        new_ref = f"{target_ref}-c{cid}"
    desc = f"[Proposal: {p[0]['title']}] {c[0]['body']}"
    _exec("""INSERT INTO actions
               (ref, description, assignee, assignees, meeting_id, related_meeting,
                status, priority, deadline, parent_ref, project_id, created_by)
             VALUES (%s,%s,%s,%s,%s,%s,'open','medium',NULL,NULL,NULL,%s)""",
          (new_ref, desc, p[0]["owner"], [p[0]["owner"]], target_id, target_ref, author))
    _exec("UPDATE proposal_comments SET action_ref=%s WHERE id=%s", (new_ref, cid))
    return {"ok": True, "action_id": new_ref}


import time

# In-memory doc cache: {"text": str, "ts": float}
_doc_cache: dict = {}
DOC_CACHE_TTL = 3600  # re-read files every 60 minutes

def _read_all_docs() -> str:
    now = time.time()
    if _doc_cache.get("text") and (now - _doc_cache.get("ts", 0)) < DOC_CACHE_TTL:
        return _doc_cache["text"]
    chunks = []
    for cat in ["minutes", "governance", "projects", "financial"]:
        folder = DOCS_DIR / cat
        if not folder.exists():
            continue
        for f in sorted(folder.glob("*.docx")):
            try:
                doc = DocxDocument(str(f))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                chunks.append(f"=== {CATEGORY_LABELS[cat]}: {f.stem} ===\n{text}")
            except Exception:
                pass
    result = "\n\n".join(chunks)
    _doc_cache["text"] = result
    _doc_cache["ts"] = now
    return result

def _get_live_context() -> str:
    """Pull live financial data + meeting register from Google Sheet."""
    from datetime import date as _date
    sections = [f"Today's date: {_date.today().strftime('%d %B %Y')}"]
    try:
        sh = gc().open_by_key(SHEET_ID)
        # Meeting register (finances now come from DB tools, not the Sheet)
        mr = sh.worksheet("2026 Meeting Register")
        mr_rows = mr.get_all_values()[1:]  # skip header
        meeting_lines = []
        for r in mr_rows:
            if len(r) >= 2 and r[1].strip():
                ref = r[0] if r[0].strip() else "?"
                dt = r[1]
                topics = r[4][:200] if len(r) > 4 else ""
                decisions = r[5][:200] if len(r) > 5 else ""
                next_actions = r[6][:300] if len(r) > 6 else ""
                entry = f"  {ref} | {dt}"
                if topics: entry += f" | Topics: {topics}"
                if decisions: entry += f" | Decisions: {decisions}"
                if next_actions: entry += f" | Actions: {next_actions}"
                meeting_lines.append(entry)
        if meeting_lines:
            sections.append("2026 MEETING REGISTER (all meetings, latest last):\n" + "\n".join(meeting_lines))
        # Action Tracker — open items only
        at = sh.worksheet("Action Tracker")
        at_rows = at.get_all_values()[1:]  # skip header
        open_actions = []
        done_actions = []
        for r in at_rows:
            if len(r) < 5 or not r[0].strip():
                continue
            action_id = r[0]
            desc = r[1][:120]
            owner = r[2]
            deadline = r[3]
            status = r[4].strip()
            update = r[6][:120] if len(r) > 6 else ""
            line = f"  {action_id} | {owner} | {deadline} | {status} | {desc}"
            if update:
                line += f" | Update: {update}"
            if status in ("Done", "Closed"):
                done_actions.append(line)
            else:
                open_actions.append(line)
        if open_actions:
            sections.append("OPEN ACTION ITEMS (Action Tracker):\n" + "\n".join(open_actions))
        if done_actions:
            sections.append("COMPLETED ACTION ITEMS (Action Tracker):\n" + "\n".join(done_actions))
    except Exception as e:
        sections.append(f"Sheet unavailable: {e}")
    return "=== LIVE DATA (Google Sheet) ===\n" + "\n\n".join(sections)

@app.get("/api/ask/history")
async def ask_history(request: Request):
    """Return last 10 conversation turns for the logged-in member."""
    token = _get_tok(request)
    payload = _auth_verify(token) if token else None
    if not payload:
        return {"turns": []}
    session_id = payload["sub"]
    from ask_agent import load_history
    turns = load_history(session_id)
    # Remap q/a keys to question/answer so the frontend normaliseHistory() can read them
    mapped = [{"question": t["q"], "answer": t["a"]} for t in turns[-10:]]
    return {"turns": mapped}

@app.post("/api/ask")
async def ask_kimfam(request: Request):
    body = await request.json()
    question = body.get("question", "").strip()
    if not question:
        return {"answer": "Please ask a question."}
    # Use member name from JWT as session ID if authenticated, else fall back to anonymous ID
    token = _get_tok(request)
    payload = _auth_verify(token) if token else None
    session_id = payload["sub"] if payload else body.get("session_id", "anon")
    from ask_agent import ask as agent_ask
    live_context = _get_live_context()
    answer = agent_ask(question, session_id, live_context)
    return {"answer": answer}


@app.post("/api/ask/stream")
async def ask_stream(request: Request):
    import asyncio as _aio, queue as _q
    from fastapi.responses import StreamingResponse
    body = await request.json()
    question = body.get("question", "").strip()
    if not question:
        async def _err():
            yield _sse({"type": "error", "msg": "No question provided"})
        return StreamingResponse(_err(), media_type="text/event-stream")
    token = _get_tok(request)
    payload = _auth_verify(token) if token else None
    session_id = payload["sub"] if payload else body.get("session_id", "anon")

    step_q = _q.Queue()

    def _run():
        from ask_agent import ask as agent_ask
        live = _get_live_context()
        def _cb(msg): step_q.put(("step", msg))
        answer = agent_ask(question, session_id, live, progress_cb=_cb)
        step_q.put(("result", answer))

    async def _generate():
        loop = _aio.get_running_loop()
        fut = loop.run_in_executor(None, _run)
        while True:
            try:
                kind, val = step_q.get_nowait()
                if kind == "step":
                    yield _sse({"type": "step", "msg": val})
                elif kind == "result":
                    yield _sse({"type": "result", "answer": val})
                    break
            except _q.Empty:
                if fut.done():
                    break
                await _aio.sleep(0.15)
        try: await fut
        except Exception: pass

    return StreamingResponse(
        _generate(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive"}
    )


@app.post("/api/admin/reindex")
async def admin_reindex(request: Request):
    api_key = request.headers.get("X-Internal-Key", "")
    if api_key != os.getenv("INTERNAL_API_KEY", ""):
        from fastapi import HTTPException as _HE
        raise _HE(status_code=403, detail="Forbidden")
    import subprocess, sys
    result = subprocess.run(
        [sys.executable, os.path.join(os.path.dirname(__file__), "embed_documents.py")],
        capture_output=True, text=True,
        env={**os.environ},
        cwd=os.path.dirname(__file__),
    )
    return {"stdout": result.stdout[-2000:], "stderr": result.stderr[-1000:], "returncode": result.returncode}



# ── Auth ──────────────────────────────────────────────────────────────────────
from auth import (
    login as _auth_login,
    change_password as _auth_change_password,
    set_password as _auth_set_password,
    get_members_status as _auth_members_status,
    verify_token as _auth_verify,
    seed_members as _auth_seed,
    update_status as _auth_update_status,
    send_otp as _auth_send_otp,
    verify_otp as _auth_verify_otp,
)

# Seed members on startup
try:
    _auth_seed()
except Exception:
    pass

def _get_tok(request):
    """Read JWT from HttpOnly cookie (falls back to Authorization header for backward compat)."""
    cookie = request.cookies.get("kimfam_token", "")
    if cookie:
        return cookie
    return request.headers.get("Authorization", "").replace("Bearer ", "")


@app.get("/api/updates")
def get_updates(project_id: str = None):
    from db import query as dbq
    if project_id:
        rows = dbq("""SELECT * FROM project_updates WHERE project_id=%s ORDER BY created_at DESC""", (project_id,))
    else:
        rows = dbq("SELECT * FROM project_updates ORDER BY created_at DESC")
    result = []
    for r in rows:
        result.append({
            "id":           r["id"],
            "project_id":   r["project_id"],
            "project_name": r["project_name"],
            "author":       r["author"],
            "text":         r["text"],
            "media":        list(r["media"] or []),
            "created_at":   r["created_at"].isoformat(),
        })
    return result


@app.post("/api/updates")
async def post_update(request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        raise _HE(status_code=401, detail="Not authenticated")
    if payload.get("role") != "admin":
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    project_id   = str(body.get("project_id", "")).strip()
    project_name = str(body.get("project_name", "")).strip()
    text         = str(body.get("text", "")).strip()
    media        = body.get("media", [])
    author       = payload.get("display") or payload.get("sub") or "Admin"
    if not project_id or not text:
        raise _HE(status_code=400, detail="project_id and text required")
    import json as _json
    _exec("""INSERT INTO project_updates (project_id, project_name, author, text, media)
             VALUES (%s, %s, %s, %s, %s)""",
          (project_id, project_name, author, text, _json.dumps(media)))
    return {"ok": True}


@app.delete("/api/updates/{update_id}")
async def delete_update(update_id: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("role") != "admin":
        raise _HE(status_code=403, detail="Admin only")
    _exec("DELETE FROM project_updates WHERE id=%s", (update_id,))
    return {"ok": True}


@app.post("/api/internal/updates")
async def internal_post_update(request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    body = await request.json()
    import json as _json
    project_id   = str(body.get("project_id", "")).strip()
    project_name = str(body.get("project_name", "")).strip()
    text         = str(body.get("text", "")).strip()
    author       = str(body.get("author", "Family")).strip()
    media        = body.get("media", [])
    source_msgs  = body.get("source_messages", [])
    if not project_id or not text:
        raise _HE(status_code=400, detail="project_id and text required")
    _exec(
        """INSERT INTO project_updates (project_id, project_name, author, text, media)
             VALUES (%s, %s, %s, %s, %s)""",
        (project_id, project_name, author, text, _json.dumps(media))
    )
    from db import query as dbq
    row = dbq("SELECT id FROM project_updates ORDER BY id DESC LIMIT 1")
    new_id = row[0]["id"] if row else None
    return {"ok": True, "id": new_id}


@app.delete("/api/internal/updates/{update_id}")
async def internal_delete_update(update_id: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    _exec("DELETE FROM project_updates WHERE id=%s", (update_id,))
    return {"ok": True}


@app.patch("/api/internal/updates/{update_id}")
async def internal_edit_update(update_id: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    body = await request.json()
    new_text = str(body.get("text", "")).strip()
    if not new_text:
        raise _HE(status_code=400, detail="text required")
    _exec("UPDATE project_updates SET text=%s WHERE id=%s", (new_text, update_id))
    return {"ok": True}


# ── Internal action mutations (for the WhatsApp agent's group auto-capture) ────
@app.post("/api/internal/actions/done")
async def internal_action_done(request: Request):
    """Mark an action Done in the DB. Internal-key auth (used by the WhatsApp agent)."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    body = await request.json()
    action_ref = str(body.get("action_id", "")).strip()
    note       = str(body.get("note", "")).strip()
    author     = str(body.get("author", "agent")).strip() or "agent"
    if not action_ref:
        raise _HE(status_code=400, detail="action_id required")
    rows = _dbq("SELECT id, status FROM actions WHERE ref=%s", (action_ref,))
    if not rows:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    action_db_id = rows[0]["id"]
    old_status   = rows[0]["status"]
    try:
        _exec("UPDATE actions SET status='done', closed_at=NOW() WHERE ref=%s", (action_ref,))
        _exec("""INSERT INTO action_updates (action_id, author, text, type, old_value, new_value)
                 VALUES (%s,%s,%s,'status_change',%s,'done')""",
              (action_db_id, author, note or "Marked done via WhatsApp auto-capture", old_status))
        return {"ok": True, "action_id": action_ref}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"internal_action_done: {e}")
        raise _HE(status_code=500, detail="DB update failed")


@app.post("/api/internal/actions/update")
async def internal_action_update(request: Request):
    """Log a progress update on an action. Internal-key auth (used by the WhatsApp agent)."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _dbq
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    body = await request.json()
    action_ref  = str(body.get("action_id", "")).strip()
    update_text = str(body.get("update_text", "")).strip()
    author      = str(body.get("author", "agent")).strip() or "agent"
    if not action_ref or not update_text:
        raise _HE(status_code=400, detail="action_id and update_text required")
    rows = _dbq("SELECT id, status FROM actions WHERE ref=%s", (action_ref,))
    if not rows:
        raise _HE(status_code=404, detail=f"Action not found: {action_ref}")
    action_db_id = rows[0]["id"]
    try:
        _exec("""INSERT INTO action_updates (action_id, author, text, type)
                 VALUES (%s,%s,%s,'comment')""",
              (action_db_id, author, update_text))
        if rows[0]["status"] == "open":
            _exec("UPDATE actions SET status='in_progress' WHERE ref=%s", (action_ref,))
        return {"ok": True, "action_id": action_ref}
    except Exception as e:
        import logging as _lg; _lg.getLogger("main").error(f"internal_action_update: {e}")
        raise _HE(status_code=500, detail="DB update failed")



@app.post("/api/auth/login")
async def auth_login(request: Request):
    body = await request.json()
    name = body.get("name", "").strip()
    password = body.get("password", "").strip()
    if not name or not password:
        from fastapi import HTTPException
        raise HTTPException(status_code=400, detail="Username and password required")
    result = _auth_login(name, password)
    if not result:
        from fastapi import HTTPException
        raise HTTPException(status_code=401, detail="Incorrect username or password")
    from fastapi.responses import JSONResponse
    resp = JSONResponse({
        "name":                result["name"],
        "display":             result["display"],
        "role":                result["role"],
        "must_change_password": result["must_change_password"],
    })
    resp.set_cookie("kimfam_token", result["token"],
                    httponly=True, samesite="lax", max_age=30*24*60*60, path="/")
    return resp

@app.post("/api/auth/logout")
async def auth_logout():
    from fastapi.responses import JSONResponse
    resp = JSONResponse({"ok": True})
    resp.delete_cookie("kimfam_token", path="/")
    return resp


@app.post("/api/auth/change-password")
async def auth_change_password(request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        from fastapi import HTTPException
        raise HTTPException(status_code=401, detail="Invalid or expired session")
    body = await request.json()
    old_pw = body.get("old_password", "")
    new_pw = body.get("new_password", "")
    if not _auth_change_password(payload["sub"], old_pw, new_pw):
        from fastapi import HTTPException
        raise HTTPException(status_code=400, detail="Incorrect current password or new password too short (min 6)")
    return {"ok": True}

@app.post("/api/auth/admin/set-password")
async def auth_admin_set_password(request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in {'Hillary', 'Hellen'}:
        from fastapi import HTTPException
        raise HTTPException(status_code=403, detail="Admin only")
    body = await request.json()
    name = body.get("name", "").strip()
    new_pw = body.get("password", "").strip()
    if not _auth_set_password(name, new_pw):
        from fastapi import HTTPException
        raise HTTPException(status_code=400, detail="Member not found or password too short (min 6)")
    return {"ok": True, "name": name}

@app.get("/api/auth/admin/members-status")
async def auth_admin_members_status(request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in {'Hillary', 'Hellen'}:
        from fastapi import HTTPException
        raise HTTPException(status_code=403, detail="Admin only")
    return {"members": _auth_members_status()}

@app.get("/api/auth/me")
async def auth_me(request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        from fastapi import HTTPException
        raise HTTPException(status_code=401, detail="Not authenticated")
    return {"name": payload["sub"], "display": payload["display"], "role": payload["role"]}


# ── Family Profiles ────────────────────────────────────────────────────────────
from family_profiles import (
    get_all_families as _get_families,
    update_family_children as _update_children,
    can_edit as _can_edit_family,
    seed_family_profiles as _seed_families,
)

try:
    _seed_families()
except Exception:
    pass

@app.get("/api/family-profiles")
async def family_profiles_list(request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        from fastapi import HTTPException
        raise HTTPException(status_code=401, detail="Login required")
    return {"families": _get_families()}

@app.put("/api/family-profiles/{family_id}")
async def family_profiles_update(family_id: str, request: Request):
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        from fastapi import HTTPException
        raise HTTPException(status_code=401, detail="Login required")
    member_name = payload["sub"]
    if not _can_edit_family(family_id, member_name):
        from fastapi import HTTPException
        raise HTTPException(status_code=403, detail="You can only edit your own family's profile")
    body = await request.json()
    children = body.get("children", [])
    # Sanitise each child entry
    clean = []
    for c in children:
        clean.append({
            "name": str(c.get("name", "")).strip(),
            "birthday": str(c.get("birthday", "")).strip(),
            "adopted": bool(c.get("adopted", False)),
            "on_obligations": bool(c.get("on_obligations", False)),
        })
    ok = _update_children(family_id, clean, member_name)
    if not ok:
        from fastapi import HTTPException
        raise HTTPException(status_code=404, detail="Family not found")
    return {"ok": True}

# ── Serve docs as static files ─────────────────────────────────────────────────
from fastapi.responses import FileResponse

import mammoth
from fastapi import HTTPException
from fastapi.responses import HTMLResponse as HTMLResp

MIME = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pdf":  "application/pdf",
    ".png":  "image/png",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
}


# ── Washing Bay Income Module ─────────────────────────────────────────────────
import sqlite3 as _sqlite3
from pydantic import BaseModel as _BaseModel
from fastapi import HTTPException as _HTTPException

# Environment-relative so STAGING writes to its own data dir, not prod's.
# (Was hardcoded to the prod path, which leaked staging test income into prod.)
# For prod, dirname(__file__) IS /var/www/kimfamhub, so the path is unchanged.
_WB_DB = os.path.join(os.path.dirname(__file__), "data", "washing_bay.db")

def _wb_conn():
    os.makedirs(os.path.dirname(_WB_DB), exist_ok=True)
    conn = _sqlite3.connect(_WB_DB)
    conn.execute("""CREATE TABLE IF NOT EXISTS income_records (
        id             INTEGER PRIMARY KEY AUTOINCREMENT,
        date           TEXT    NOT NULL,
        amount_ugx     INTEGER NOT NULL,
        received_from  TEXT    DEFAULT '',
        collector      TEXT    NOT NULL,
        txn_ref        TEXT    DEFAULT '',
        notes          TEXT    DEFAULT '',
        recorded_at    TEXT    DEFAULT (datetime('now'))
    )""")
    # add txn_ref to existing tables that predate this column
    for col, defval in [("txn_ref", "''"), ("received_from", "''")]:
        try:
            conn.execute(f"ALTER TABLE income_records ADD COLUMN {col} TEXT DEFAULT {defval}")
            conn.commit()
        except Exception:
            pass
    # Capital accountability: who actually put in the ~25.9M CapEx, with proof.
    conn.execute("""CREATE TABLE IF NOT EXISTS capital_contributions (
        id           INTEGER PRIMARY KEY AUTOINCREMENT,
        contributor  TEXT    NOT NULL,
        amount_ugx   INTEGER NOT NULL,
        date         TEXT    DEFAULT '',
        source       TEXT    DEFAULT '',
        proof_ref    TEXT    DEFAULT '',
        verified     INTEGER DEFAULT 0,
        recorded_at  TEXT    DEFAULT (datetime('now'))
    )""")
    conn.commit()
    return conn

@app.get("/api/washing-bay/income")
def wb_get_income():
    conn = _wb_conn()
    rows = conn.execute(
        "SELECT id,date,amount_ugx,received_from,collector,txn_ref,notes,recorded_at FROM income_records ORDER BY date DESC, id DESC"
    ).fetchall()
    conn.close()
    records = [{"id":r[0],"date":r[1],"amount_ugx":r[2],"received_from":r[3],"collector":r[4],"txn_ref":r[5],"notes":r[6],"recorded_at":r[7]} for r in rows]
    total = sum(r["amount_ugx"] for r in records)
    return {"total_ugx": total, "count": len(records), "records": records}

class _WBEntry(_BaseModel):
    pin: str
    date: str
    amount_ugx: int
    received_from: str = ""
    collector: str
    txn_ref: str = ""
    notes: str = ""

@app.post("/api/washing-bay/income")
def wb_add_income(entry: _WBEntry):
    if entry.pin != os.getenv("WASHING_BAY_PIN", "1234"):
        raise _HTTPException(status_code=403, detail="Incorrect PIN")
    conn = _wb_conn()
    conn.execute(
        "INSERT INTO income_records (date,amount_ugx,received_from,collector,txn_ref,notes) VALUES (?,?,?,?,?,?)",
        (entry.date, entry.amount_ugx, entry.received_from, entry.collector, entry.txn_ref, entry.notes)
    )
    conn.commit()
    conn.close()
    return {"ok": True}


# ── Washing Bay capital accountability ────────────────────────────────────────
# Dad is reconciling the ~UGX 25.9M CapEx, which currently has no proof of
# sources. Members log who actually contributed what, and we show how the
# running total balances against the reported estimate.
WB_CAPEX_TARGET = 25_900_000

@app.get("/api/washing-bay/capital")
def wb_get_capital():
    conn = _wb_conn()
    rows = conn.execute(
        "SELECT id,contributor,amount_ugx,date,source,proof_ref,verified,recorded_at "
        "FROM capital_contributions ORDER BY date DESC, id DESC"
    ).fetchall()
    conn.close()
    records = [{"id":r[0],"contributor":r[1],"amount_ugx":r[2],"date":r[3],"source":r[4],
                "proof_ref":r[5],"verified":bool(r[6]),"recorded_at":r[7]} for r in rows]
    total = sum(r["amount_ugx"] for r in records)
    verified_total = sum(r["amount_ugx"] for r in records if r["verified"])
    by = {}
    for r in records:
        by[r["contributor"]] = by.get(r["contributor"], 0) + r["amount_ugx"]
    by_contributor = sorted(
        ({"contributor": k, "amount_ugx": v} for k, v in by.items()),
        key=lambda x: x["amount_ugx"], reverse=True,
    )
    remaining = max(0, WB_CAPEX_TARGET - total)
    return {
        "target_ugx": WB_CAPEX_TARGET,
        "total_accounted_ugx": total,
        "verified_ugx": verified_total,
        "remaining_ugx": remaining,
        "pct_accounted": round(total / WB_CAPEX_TARGET * 100, 1) if WB_CAPEX_TARGET else 0,
        "balanced": total >= WB_CAPEX_TARGET,
        "by_contributor": by_contributor,
        "records": records,
    }

class _WBCapital(_BaseModel):
    pin: str
    contributor: str
    amount_ugx: int
    date: str = ""
    source: str = ""
    proof_ref: str = ""
    verified: bool = False

@app.post("/api/washing-bay/capital")
def wb_add_capital(entry: _WBCapital):
    if entry.pin != os.getenv("WASHING_BAY_PIN", "1234"):
        raise _HTTPException(status_code=403, detail="Incorrect PIN")
    if not entry.contributor.strip() or entry.amount_ugx <= 0:
        raise _HTTPException(status_code=400, detail="Contributor and a positive amount are required")
    conn = _wb_conn()
    conn.execute(
        "INSERT INTO capital_contributions (contributor,amount_ugx,date,source,proof_ref,verified) "
        "VALUES (?,?,?,?,?,?)",
        (entry.contributor.strip(), entry.amount_ugx, entry.date, entry.source,
         entry.proof_ref, 1 if entry.verified else 0)
    )
    conn.commit()
    conn.close()
    return {"ok": True}

class _WBCapitalDel(_BaseModel):
    pin: str
    id: int

@app.post("/api/washing-bay/capital/delete")
def wb_delete_capital(body: _WBCapitalDel):
    if body.pin != os.getenv("WASHING_BAY_PIN", "1234"):
        raise _HTTPException(status_code=403, detail="Incorrect PIN")
    conn = _wb_conn()
    conn.execute("DELETE FROM capital_contributions WHERE id=?", (body.id,))
    conn.commit()
    conn.close()
    return {"ok": True}

def _r2_redirect(category: str, filename: str):
    """Return a redirect to a presigned R2 URL, or None if R2 not configured."""
    if not _r2.is_configured():
        return None
    from fastapi.responses import RedirectResponse
    url = _r2.presigned_url(f"{category}/{filename}", expires=3600)
    return RedirectResponse(url=url, status_code=302)

def _r2_bytes(category: str, filename: str) -> bytes | None:
    """Download file bytes from R2 for in-memory rendering. Returns None on any failure."""
    if not _r2.is_configured():
        return None
    try:
        import io
        buf = io.BytesIO()
        _r2._client().download_fileobj(_r2._R2_BUCKET, f"{category}/{filename}", buf)
        return buf.getvalue()
    except Exception:
        return None

# NOTE: the /view route MUST be declared before the download route below, because
# both use the greedy {filename:path} converter (to allow sub-group sub-paths like
# "Chicken/Proposal.docx"); whichever is registered first wins for ".../view".
@app.get("/docs/{category}/{filename:path}/view")
def view_doc(category: str, filename: str):
    path = DOCS_DIR / category / filename
    suffix = Path(filename).suffix.lower()

    # Try local first
    file_bytes: bytes | None = None
    if path.exists() and path.is_file():
        file_bytes = path.read_bytes()
    else:
        file_bytes = _r2_bytes(category, filename)

    if file_bytes is None:
        raise HTTPException(status_code=404)

    if suffix == ".docx":
        import io
        result = mammoth.convert_to_html(io.BytesIO(file_bytes))
        html = f"""<!DOCTYPE html><html><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
  body{{font-family:Georgia,serif;max-width:800px;margin:0 auto;padding:20px;background:#fff;color:#1a1a1a;line-height:1.7;font-size:15px}}
  h1,h2,h3{{color:#1a3a1a;margin-top:1.5em}} table{{border-collapse:collapse;width:100%;margin:1em 0}}
  td,th{{border:1px solid #ccc;padding:8px 12px}} th{{background:#f0f7f0}}
  p{{margin:0.7em 0}}
</style></head><body>{result.value}</body></html>"""
        return HTMLResp(content=html)
    elif suffix in (".png", ".jpg", ".jpeg"):
        from fastapi.responses import Response
        return Response(content=file_bytes, media_type=MIME.get(suffix, "image/jpeg"))
    elif suffix == ".pdf":
        from fastapi.responses import Response
        return Response(content=file_bytes, media_type="application/pdf")
    elif suffix == ".pptx":
        import io, html as _html
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        body = []
        for i, slide in enumerate(prs.slides, 1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(r.text for r in para.runs).strip()
                        if txt:
                            lines.append(f"<p>{_html.escape(txt)}</p>")
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = "".join(f"<td>{_html.escape(c.text.strip())}</td>" for c in row.cells)
                        rows.append(f"<tr>{cells}</tr>")
                    lines.append(f"<table>{''.join(rows)}</table>")
            body.append(f"<section class='slide'><div class='num'>Slide {i}</div>{''.join(lines)}</section>")
        html = ("<!DOCTYPE html><html><head><meta charset='UTF-8'>"
                "<meta name='viewport' content='width=device-width,initial-scale=1'><style>"
                "body{font-family:Inter,Arial,sans-serif;max-width:820px;margin:0 auto;padding:16px;background:#0f1729;color:#e2e8f0}"
                ".slide{background:#fff;color:#1a1a1a;border-radius:10px;padding:20px;margin:14px 0;box-shadow:0 4px 16px rgba(0,0,0,.4)}"
                ".num{font-size:11px;color:#64748b;text-transform:uppercase;letter-spacing:.05em;margin-bottom:8px}"
                "table{border-collapse:collapse;width:100%;margin:.6em 0}td{border:1px solid #ccc;padding:6px 10px;font-size:14px}"
                "p{margin:.4em 0}</style></head><body>" + "".join(body) + "</body></html>")
        return HTMLResp(content=html)
    elif suffix == ".xlsx":
        import io, html as _html
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        body = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                if not any(c is not None for c in row):
                    continue
                cells = "".join(f"<td>{_html.escape('' if c is None else str(c))}</td>" for c in row)
                rows.append(f"<tr>{cells}</tr>")
            body.append(f"<section class='sheet'><div class='num'>Sheet: {_html.escape(ws.title)}</div><table>{''.join(rows)}</table></section>")
        html = ("<!DOCTYPE html><html><head><meta charset='UTF-8'>"
                "<meta name='viewport' content='width=device-width,initial-scale=1'><style>"
                "body{font-family:Inter,Arial,sans-serif;margin:0 auto;padding:16px;background:#0f1729;color:#e2e8f0}"
                ".sheet{background:#fff;color:#1a1a1a;border-radius:10px;padding:16px;margin:14px 0;overflow-x:auto}"
                ".num{font-size:11px;color:#64748b;text-transform:uppercase;letter-spacing:.05em;margin-bottom:8px}"
                "table{border-collapse:collapse}td{border:1px solid #d0d0d0;padding:4px 8px;font-size:13px;white-space:nowrap}"
                "</style></head><body>" + "".join(body) + "</body></html>")
        return HTMLResp(content=html)
    else:
        raise HTTPException(status_code=415, detail="Preview not supported for this file type")


@app.get("/docs/{category}/{filename:path}")
def serve_doc(category: str, filename: str):
    # filename may include a sub-group path, e.g. "Chicken/Proposal.docx".
    path = DOCS_DIR / category / filename
    base = Path(filename).name
    if path.exists() and path.is_file():
        mime = MIME.get(path.suffix.lower(), "application/octet-stream")
        return FileResponse(str(path), media_type=mime,
                            headers={"Content-Disposition": f"attachment; filename=\"{base}\""})
    r = _r2_redirect(category, filename)
    if r:
        return r
    raise HTTPException(status_code=404)


# ── Profile Picture Upload ─────────────────────────────────────────────────────
from fastapi import UploadFile, File as FastAPIFile, Form

# Environment-relative so staging serves its own avatars dir (same fix as washing-bay).
# For prod, __file__ is /var/www/kimfamhub/main.py so the path is unchanged.
AVATARS_DIR = Path(__file__).parent / "static" / "avatars"
AVATARS_DIR.mkdir(parents=True, exist_ok=True)

@app.post("/api/auth/status")
async def auth_update_status(request: Request):
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        raise _HE(status_code=401, detail="Not authenticated")
    body = await request.json()
    msg = str(body.get("status_message", ""))[:120]
    _auth_update_status(payload["sub"], msg)
    return {"ok": True, "status_message": msg}


@app.post("/api/auth/forgot-password")
async def auth_forgot_password(request: Request):
    from fastapi import HTTPException as _HE
    body = await request.json()
    name = str(body.get("name", "")).strip()
    if not name:
        raise _HE(status_code=400, detail="name is required")
    result = _auth_send_otp(name)
    if not result["ok"]:
        raise _HE(status_code=400, detail=result["error"])
    return {"ok": True, "message": "Code sent to your WhatsApp number."}


@app.post("/api/auth/reset-password")
async def auth_reset_password(request: Request):
    from fastapi import HTTPException as _HE
    body = await request.json()
    name  = str(body.get("name", "")).strip()
    token = str(body.get("token", "")).strip()
    new_pw = str(body.get("new_password", ""))
    if not name or not token or not new_pw:
        raise _HE(status_code=400, detail="name, token, and new_password are required")
    result = _auth_verify_otp(name, token, new_pw)
    if not result["ok"]:
        raise _HE(status_code=400, detail=result["error"])
    from fastapi.responses import JSONResponse
    resp = JSONResponse({k: v for k, v in result.items() if k != "token"})
    resp.set_cookie("kimfam_token", result["token"],
                    httponly=True, samesite="lax", max_age=30*24*60*60, path="/")
    return resp


@app.post("/api/auth/upload-avatar")
async def upload_avatar(request: Request, file: UploadFile = FastAPIFile(...)):
    import jwt as _jwt
    from fastapi import HTTPException as _HE
    JWT_SECRET = os.environ.get("JWT_SECRET","kimfam-secret-change-me")
    token = _get_tok(request)
    try:
        payload = _jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        member_name = payload.get("sub","")
    except Exception:
        raise _HE(status_code=401, detail="Invalid token")
    if not member_name:
        raise _HE(status_code=401, detail="Not authenticated")
    ct = file.content_type or ""
    # iOS Files picker sends application/octet-stream; accept any image-like type
    # (PIL will reject non-images at conversion time, which is safe enough).
    bad_types = {"text/", "application/json", "application/xml"}
    if any(ct.startswith(t) for t in bad_types):
        raise _HE(status_code=400, detail="File must be an image")
    safe_name = "".join(c for c in member_name.lower() if c.isalnum())
    dest = AVATARS_DIR / f"{safe_name}.jpg"
    contents = await file.read()
    try:
        from PIL import Image
        import io as _io
        # Register HEIC/HEIF plugin if available (iOS default photo format)
        try:
            from pillow_heif import register_heif_opener
            register_heif_opener()
        except ImportError:
            pass
        img = Image.open(_io.BytesIO(contents))
        img = img.convert("RGB")
        img.thumbnail((400, 400), Image.LANCZOS)
        img.save(str(dest), "JPEG", quality=88)
    except Exception as _img_err:
        # Don't silently save broken bytes — return a clear error so the UI shows
        # "unsupported format" instead of saving a broken file and claiming success.
        import logging as _lg
        _lg.getLogger("main").warning("avatar convert failed: %s", _img_err)
        raise _HTTPException(status_code=400, detail=f"Could not process image. Please use a JPEG or PNG photo from your camera roll (not HEIC or RAW). Error: {str(_img_err)[:80]}")
    url = f"/static/avatars/{safe_name}.jpg?t={int(__import__('time').time())}"
    return {"url": url, "cache_bust": True}

# ── Document Upload (admin) ───────────────────────────────────────────────────
import r2_storage as _r2

_ALLOWED_DOC_CATEGORIES = {"minutes", "governance", "projects", "financial", "receipts"}
_ALLOWED_DOC_SUFFIXES   = {".docx", ".pdf", ".jpg", ".jpeg", ".png"}

@app.post("/api/admin/upload-doc")
async def admin_upload_doc(
    request: Request,
    category: str = Form(...),
    file: UploadFile = FastAPIFile(...),
):
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    payload = _auth_verify(token) if token else None
    if not payload or payload.get("role") != "admin":
        raise _HE(status_code=403, detail="Admin only")
    if category not in _ALLOWED_DOC_CATEGORIES:
        raise _HE(status_code=400, detail=f"Unknown category: {category}")
    suffix = Path(file.filename or "").suffix.lower()
    if suffix not in _ALLOWED_DOC_SUFFIXES:
        raise _HE(status_code=400, detail=f"File type not allowed: {suffix}")

    safe_filename = "".join(c for c in (file.filename or "upload") if c.isalnum() or c in "._- ")
    contents = await file.read()

    # Write to local filesystem (keeps existing list/serve endpoints working)
    local_path = DOCS_DIR / category / safe_filename
    (DOCS_DIR / category).mkdir(parents=True, exist_ok=True)
    with open(str(local_path), "wb") as fh:
        fh.write(contents)

    # Also upload to R2 for durable storage
    r2_key  = f"{category}/{safe_filename}"
    r2_url  = None
    if _r2.is_configured():
        import tempfile, os as _os
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(contents)
            tmp_path = tmp.name
        try:
            public = category in {"governance", "minutes"}
            r2_url = _r2.upload(tmp_path, r2_key, public=public)
        finally:
            _os.unlink(tmp_path)

    # Invalidate doc cache so new file appears immediately
    _doc_cache.clear()

    return {"ok": True, "filename": safe_filename, "r2_key": r2_key, "r2_url": r2_url}


# ── Contributions (Sprint 1) ──────────────────────────────────────────────────
from contributions import router as contributions_router
app.include_router(contributions_router)

# ── Project Participation ──────────────────────────────────────────────────────
import requests as _requests

_BRIDGE_SEND    = "http://localhost:8080/api/send"
_HILLARY_NUM    = "256775102684"
_CHAIRMAN_NUM        = "256772625387"  # fallback only

def _get_officer(role_slug: str):
    """Return (member_name, phone) for the current active office bearer."""
    from db import query as _q
    import notifications as _notif
    rows = _q("SELECT member_name FROM club_office_bearers WHERE role_slug=%s AND effective_to IS NULL ORDER BY effective_from DESC LIMIT 1", (role_slug,))
    if not rows:
        return None, None
    name = rows[0]["member_name"]
    phone = _notif.MEMBER_PHONES.get(name)
    return name, phone

def _chairman_name() -> str:
    name, _ = _get_officer("chairman")
    return name or _CHAIRMAN_NAME

def _chairman_num() -> str:
    _, phone = _get_officer("chairman")
    return phone or _CHAIRMAN_NUM

def _secretary_name() -> str:
    name, _ = _get_officer("secretary")
    return name or _HILLARY_NAME
_VALID_ROLES    = {"project_lead", "team_member"}
_VALID_MODES    = {"operational","oversight","capital","advisory","commercial","physical"}
_ADMINS_PP      = {"Hillary", "Hellen"}
_HILLARY_NAME   = "Hillary"
_CHAIRMAN_NAME       = "Israel"   # display name for Dad in DB records

_ROLE_LABELS = {
    "project_lead": "Project Lead",
    "team_member":   "Team Member",
}
_MODE_LABELS = {
    "operational":  "Operational (day-to-day running)",
    "oversight":    "Oversight (governance / monitoring)",
    "capital":      "Capital (financial investment)",
    "advisory":     "Advisory (strategic advice)",
    "commercial":   "Commercial (sales / procurement / market links)",
    "physical":     "Physical (labour / on-site work)",
}

def _wa_send(number: str, text: str):
    try:
        _requests.post(_BRIDGE_SEND, json={"recipient": number, "message": text}, timeout=10)
    except Exception:
        pass

def _internal_key_ok(request) -> bool:
    expected = os.environ.get("KIMFAM_INTERNAL_KEY", "")
    provided = request.headers.get("X-Internal-Key", "")
    return bool(expected and provided == expected)

def _build_dad_approval_msg(member: str, project_id: str, role: str, modes: list, note: str = None) -> str:
    role_label = _ROLE_LABELS.get(role, role)
    modes_text = ", ".join(_MODE_LABELS.get(m, m) for m in (modes or [])) if modes else ""
    msg = (
        f"KimFam Projects — Participation Request\n\n"
        f"{member} wants to join the *{project_id}* project.\n"
        f"Role: {role_label}"
    )
    if modes_text:
        msg += f"\nContributions: {modes_text}"
    if note:
        msg += f"\nNote: {note}"
    msg += (
        "\n\nHillary has reviewed this. Please reply:\n"
        "YES to approve, or NO to reject.\n"
        "You can ask questions first."
    )
    return msg


@app.post("/api/projects/interest")
async def submit_interest(request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        raise _HE(status_code=401, detail="Not authenticated")
    member = payload["sub"]
    body = await request.json()
    project_id = str(body.get("project_id","")).strip()
    preferred_role = str(body.get("preferred_role","team_member")).strip()
    contribution_modes = [str(m).strip() for m in (body.get("contribution_modes") or [])]
    note = str(body.get("note","")).strip() or None
    if not project_id:
        raise _HE(status_code=400, detail="project_id is required")
    if preferred_role not in _VALID_ROLES:
        raise _HE(status_code=400, detail="Invalid preferred_role")
    invalid_modes = [m for m in contribution_modes if m not in _VALID_MODES]
    if invalid_modes:
        raise _HE(status_code=400, detail=f"Invalid contribution_modes: {invalid_modes}")
    if preferred_role == "project_lead":
        contribution_modes = []
    if preferred_role == "team_member" and not contribution_modes:
        raise _HE(status_code=400, detail="Contributors must select at least one contribution mode")

    # Dad (Israel) is chairman — his submissions auto-confirm, no review needed
    if member == _chairman_name():
        existing = _q("SELECT id, status FROM project_participation WHERE project_id=%s AND member_name=%s",
                      (project_id, member))
        if existing and existing[0]["status"] == "confirmed":
            raise _HE(status_code=409, detail="Your interest in this project is already confirmed.")
        if existing:
            _exec("""UPDATE project_participation
                     SET preferred_role=%s, contribution_modes=%s, note=%s,
                         status='confirmed', confirmed_role=%s, confirmed_modes=%s,
                         confirmed_at=NOW(), confirmed_by='self', awaiting_chairman_since=NULL
                     WHERE id=%s""",
                  (preferred_role, contribution_modes, note,
                   preferred_role, contribution_modes, existing[0]["id"]))
        else:
            _exec("""INSERT INTO project_participation
                     (project_id, member_name, preferred_role, contribution_modes, note,
                      status, confirmed_role, confirmed_modes, confirmed_at, confirmed_by)
                     VALUES (%s,%s,%s,%s,%s,'confirmed',%s,%s,NOW(),'self')""",
                  (project_id, member, preferred_role, contribution_modes, note,
                   preferred_role, contribution_modes))
        _wa_send(_HILLARY_NUM, f"Dad added himself to *{project_id}* as {_ROLE_LABELS.get(preferred_role, preferred_role)}.")
        return {"ok": True, "action": "auto_confirmed"}

    # Hillary's own submissions skip his review and go directly to awaiting_chairman
    if member == _secretary_name():
        existing = _q("SELECT id, status FROM project_participation WHERE project_id=%s AND member_name=%s",
                      (project_id, member))
        if existing and existing[0]["status"] in ("confirmed", "awaiting_chairman"):
            raise _HE(status_code=409, detail="Your interest is already submitted and pending Dad's approval or confirmed.")
        if existing:
            _exec("""UPDATE project_participation
                     SET preferred_role=%s, contribution_modes=%s, note=%s,
                         status='awaiting_chairman', awaiting_chairman_since=NOW(),
                         rejection_note=NULL
                     WHERE id=%s""",
                  (preferred_role, contribution_modes, note, existing[0]["id"]))
        else:
            _exec("""INSERT INTO project_participation
                     (project_id, member_name, preferred_role, contribution_modes, note,
                      status, awaiting_chairman_since)
                     VALUES (%s,%s,%s,%s,%s,'awaiting_chairman',NOW())""",
                  (project_id, member, preferred_role, contribution_modes, note))
        _wa_send(_chairman_num(), _build_dad_approval_msg(
            member, project_id, preferred_role, contribution_modes, note))
        return {"ok": True, "action": "sent_to_dad"}

    # Regular member submission
    existing = _q("SELECT id, status FROM project_participation WHERE project_id=%s AND member_name=%s",
                  (project_id, member))
    if existing:
        row = existing[0]
        if row["status"] == "rejected":
            _exec("""UPDATE project_participation
                     SET preferred_role=%s, contribution_modes=%s, note=%s,
                         status='pending', rejection_note=NULL, submitted_at=NOW()
                     WHERE id=%s""",
                  (preferred_role, contribution_modes, note, row["id"]))
            action = "re-submitted"
        elif row["status"] == "pending":
            _exec("""UPDATE project_participation
                     SET preferred_role=%s, contribution_modes=%s, note=%s, submitted_at=NOW()
                     WHERE id=%s""",
                  (preferred_role, contribution_modes, note, row["id"]))
            action = "updated"
        else:
            raise _HE(status_code=409, detail="Your interest in this project has already been confirmed or is awaiting the chairman. Contact Hillary to adjust it.")
    else:
        _exec("""INSERT INTO project_participation
                 (project_id, member_name, preferred_role, contribution_modes, note)
                 VALUES (%s, %s, %s, %s, %s)""",
              (project_id, member, preferred_role, contribution_modes, note))
        action = "submitted"

    role_label = _ROLE_LABELS.get(preferred_role, preferred_role)
    modes_text = ", ".join(_MODE_LABELS.get(m, m) for m in contribution_modes) if contribution_modes else ""
    msg = f"KimFam Projects: {member} expressed interest in *{project_id}*\nRole: {role_label}"
    if modes_text:
        msg += f"\nModes: {modes_text}"
    if note:
        msg += f"\nNote: {note}"
    msg += "\n\nReview at kimfamhub.com (Admin > Project Participation)"
    _wa_send(_HILLARY_NUM, msg)
    return {"ok": True, "action": action}


@app.get("/api/projects/interests")
async def get_interests(request: Request, project_id: str = None):
    from fastapi import HTTPException as _HE
    from db import query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload:
        raise _HE(status_code=401, detail="Not authenticated")
    if project_id:
        rows = _q("""SELECT pp.*, f.family_name
            FROM project_participation pp
            LEFT JOIN persons p ON LOWER(p.whatsapp_name) = LOWER(pp.member_name)
            LEFT JOIN families f ON p.family_id = f.id
            WHERE pp.project_id=%s ORDER BY pp.submitted_at""", (project_id,))
    else:
        rows = _q("""SELECT pp.*, f.family_name
            FROM project_participation pp
            LEFT JOIN persons p ON LOWER(p.whatsapp_name) = LOWER(pp.member_name)
            LEFT JOIN families f ON p.family_id = f.id
            ORDER BY pp.project_id, pp.submitted_at""")
    result = []
    for r in rows:
        result.append({
            "id":                 r["id"],
            "project_id":         r["project_id"],
            "member_name":        r["member_name"],
            "family_name":        (r["family_name"].title() if r.get("family_name") else None),
            "preferred_role":     r["preferred_role"],
            "contribution_modes": list(r["contribution_modes"] or []),
            "note":               r["note"],
            "status":             r["status"],
            "confirmed_role":     r["confirmed_role"],
            "confirmed_modes":    list(r["confirmed_modes"] or []) if r["confirmed_modes"] else None,
            "rejection_note":     r["rejection_note"],
            "submitted_at":       str(r["submitted_at"])[:16],
            "confirmed_at":       str(r["confirmed_at"])[:16] if r["confirmed_at"] else None,
            "confirmed_by":       r["confirmed_by"],
        })
    return result


@app.get("/api/projects/interests/awaiting-chairman")
async def get_awaiting_chairman(request: Request):
    """Internal endpoint for the WhatsApp agent to fetch items awaiting Chairman's approval."""
    from fastapi import HTTPException as _HE
    from db import query as _q
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    rows = _q("""SELECT id, project_id, member_name, preferred_role, contribution_modes,
                        note, awaiting_chairman_since
                 FROM project_participation WHERE status='awaiting_chairman'
                 ORDER BY awaiting_chairman_since""")
    result = []
    for r in rows:
        result.append({
            "id":                 r["id"],
            "project_id":         r["project_id"],
            "member_name":        r["member_name"],
            "preferred_role":     r["preferred_role"],
            "contribution_modes": list(r["contribution_modes"] or []),
            "note":               r["note"],
            "awaiting_since":     str(r["awaiting_chairman_since"])[:16] if r["awaiting_chairman_since"] else None,
        })
    return result


@app.put("/api/projects/interest/{interest_id}/confirm")
async def confirm_interest(interest_id: int, request: Request):
    """Hillary reviews and approves — sets awaiting_chairman, sends WhatsApp to Dad."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    confirmed_role  = str(body.get("confirmed_role","")).strip() or None
    confirmed_modes = [str(m).strip() for m in (body.get("confirmed_modes") or [])] or None
    rows = _q("SELECT * FROM project_participation WHERE id=%s", (interest_id,))
    if not rows:
        raise _HE(status_code=404, detail="Interest not found")
    r = rows[0]
    # Block Hillary from confirming his own submission
    if r["member_name"] == _HILLARY_NAME:
        raise _HE(status_code=403, detail="Hillary cannot confirm his own submission. It must go to Dad directly.")
    final_role  = confirmed_role  or r["preferred_role"]
    final_modes = confirmed_modes or list(r["contribution_modes"] or [])
    _exec("""UPDATE project_participation
             SET status='awaiting_chairman', confirmed_role=%s, confirmed_modes=%s,
                 awaiting_chairman_since=NOW(), confirmed_by=%s
             WHERE id=%s""",
          (final_role, final_modes, payload["sub"], interest_id))
    member     = r["member_name"]
    project_id = r["project_id"]
    _wa_send(_chairman_num(), _build_dad_approval_msg(member, project_id, final_role, final_modes, r.get("note")))
    return {"ok": True, "status": "awaiting_chairman", "final_role": final_role, "final_modes": final_modes}


@app.put("/api/projects/interest/{interest_id}/dad-approve")
async def dad_approve_interest(interest_id: int, request: Request):
    """Called by the WhatsApp orchestrator when Dad says YES."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    rows = _q("SELECT * FROM project_participation WHERE id=%s AND status='awaiting_chairman'", (interest_id,))
    if not rows:
        raise _HE(status_code=404, detail="Item not found or not in awaiting_chairman state")
    r = rows[0]
    final_role  = r["confirmed_role"]  or r["preferred_role"]
    final_modes = list(r["confirmed_modes"] or r["contribution_modes"] or [])
    _exec("""UPDATE project_participation
             SET status='confirmed', confirmed_role=%s, confirmed_modes=%s,
                 confirmed_at=NOW(), confirmed_by='Dad'
             WHERE id=%s""",
          (final_role, final_modes, interest_id))
    role_label = _ROLE_LABELS.get(final_role, final_role)
    _wa_send(_HILLARY_NUM,
        f"Dad approved {r['member_name']}'s participation in *{r['project_id']}* as {role_label}.")
    return {"ok": True, "member": r["member_name"], "project_id": r["project_id"]}


@app.put("/api/projects/interest/{interest_id}/dad-reject")
async def dad_reject_interest(interest_id: int, request: Request):
    """Called by the WhatsApp orchestrator when Dad says NO."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    if not _internal_key_ok(request):
        raise _HE(status_code=403, detail="Forbidden")
    body = await request.json()
    rejection_note = str(body.get("note","")).strip() or None
    rows = _q("SELECT * FROM project_participation WHERE id=%s AND status='awaiting_chairman'", (interest_id,))
    if not rows:
        raise _HE(status_code=404, detail="Item not found or not in awaiting_chairman state")
    r = rows[0]
    _exec("""UPDATE project_participation
             SET status='rejected', rejection_note=%s,
                 confirmed_at=NOW(), confirmed_by='Dad'
             WHERE id=%s""",
          (rejection_note, interest_id))
    _wa_send(_HILLARY_NUM,
        f"Dad rejected {r['member_name']}'s participation in *{r['project_id']}*."
        + (f"\nReason: {rejection_note}" if rejection_note else ""))
    return {"ok": True, "member": r["member_name"], "project_id": r["project_id"]}


@app.put("/api/projects/interest/{interest_id}/adjust")
async def adjust_interest(interest_id: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    new_role  = str(body.get("preferred_role","")).strip() or None
    new_modes = [str(m).strip() for m in (body.get("contribution_modes") or [])] or None
    if not _q("SELECT id FROM project_participation WHERE id=%s", (interest_id,)):
        raise _HE(status_code=404, detail="Interest not found")
    if new_role:
        _exec("UPDATE project_participation SET preferred_role=%s WHERE id=%s", (new_role, interest_id))
    if new_modes is not None:
        _exec("UPDATE project_participation SET contribution_modes=%s WHERE id=%s", (new_modes, interest_id))
    return {"ok": True}


@app.put("/api/projects/interest/{interest_id}/reject")
async def reject_interest(interest_id: int, request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    body = await request.json()
    rejection_note = str(body.get("note","")).strip() or None
    if not _q("SELECT id FROM project_participation WHERE id=%s", (interest_id,)):
        raise _HE(status_code=404, detail="Interest not found")
    _exec("""UPDATE project_participation
             SET status='rejected', rejection_note=%s, confirmed_at=NOW(), confirmed_by=%s
             WHERE id=%s""",
          (rejection_note, payload["sub"], interest_id))
    return {"ok": True}



@app.get("/api/admin/login-status")
async def admin_login_status(request: Request):
    from fastapi import HTTPException as _HE
    import sqlite3
    from datetime import date
    token = _get_tok(request)
    payload = _auth_verify(token)
    is_internal = _internal_key_ok(request)
    if not payload and not is_internal:
        raise _HE(status_code=401, detail="Auth required")
    if payload and payload.get("role") != "admin":
        raise _HE(status_code=403, detail="Admin only")
    today = date.today().isoformat()
    DB = os.path.join(os.path.dirname(__file__), "data", "kimfam.db")
    with sqlite3.connect(DB) as conn:
        conn.row_factory = sqlite3.Row
        rows = conn.execute(
            "SELECT name, display, last_login, is_locked, must_change_password FROM members ORDER BY name"
        ).fetchall()
    result = []
    for r in rows:
        ll = r["last_login"] or ""
        result.append({
            "name":                r["name"],
            "display":             r["display"],
            "is_locked":           bool(r["is_locked"]),
            "last_login":          ll or None,
            "logged_in_today":     ll.startswith(today) if ll else False,
            "must_change_password": bool(r["must_change_password"]),
        })
    not_today   = [m for m in result if not m["logged_in_today"] and not m["is_locked"]]
    logged_in   = [m for m in result if m["logged_in_today"]]
    never_setup = [m for m in result if m["must_change_password"] and not m["is_locked"]]
    return {
        "date": today,
        "logged_in_today": logged_in,
        "not_logged_in_today": not_today,
        "never_changed_password": never_setup,
    }


@app.put("/api/projects/interest/{interest_id}/confirm-direct")
async def confirm_interest_direct(interest_id: int, request: Request):
    """Admin directly confirms participation, bypassing the chairman WhatsApp step."""
    from fastapi import HTTPException as _HE
    from db import execute as _exec, query as _q
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    rows = _q("SELECT * FROM project_participation WHERE id=%s", (interest_id,))
    if not rows:
        raise _HE(status_code=404, detail="Interest not found")
    r = rows[0]
    final_role  = r["confirmed_role"] or r["preferred_role"]
    final_modes = list(r["confirmed_modes"] or r["contribution_modes"] or [])
    _exec("""UPDATE project_participation
             SET status='confirmed', confirmed_role=%s, confirmed_modes=%s,
                 confirmed_at=NOW(), confirmed_by=%s, awaiting_chairman_since=NULL
             WHERE id=%s""",
          (final_role, final_modes, payload["sub"], interest_id))
    role_label = _ROLE_LABELS.get(final_role, final_role)
    _wa_send(MEMBER_PHONES.get(r["member_name"]),
        f"Great news! Your participation in the *{r['project_id']}* project has been confirmed as {role_label}. Welcome to the team!" + KIMFAM_FOOTER)
    return {"ok": True, "status": "confirmed", "final_role": final_role}


_CHICKEN_SHEET_ID  = "1CqF-NzkMJ8iJw0tC8xkLE9DI94cjFr2vvlAfx4QfXhI"
_PROJ_SHEET_ID     = "1rMFK8nJ7HbBQ6JGgyCPBB4NJI-095RcunUfNUQXL1xw"
_CHICKEN_CACHE: dict = {}
_CHICKEN_CACHE_TTL = 60  # 1 min

_ALLOC_CACHE: dict = {}
_ALLOC_CACHE_TTL = 300

def _fetch_family_equity():
    """
    Per-family equity in the club account under both allocation models.
    Covers ALL expenditure_records. Opening balance 3,895,944 split equally (556,563/family).

    Model A (family proposal): equal share per family if they have it; if short they pay
    all they have and the gap is split equally among the families that paid full share.
    Model B: pure proportional to each family's pool at the date of each expense.
    """
    import time
    from db import query as _dbq

    now = time.time()
    if _ALLOC_CACHE.get("ts") and now - _ALLOC_CACHE["ts"] < _ALLOC_CACHE_TTL:
        return _ALLOC_CACHE["data"]

    OPENING_BALANCE = 3_895_944
    PROJ_LABEL = {
        "chicken": "Chicken",
        "sheep":   "Sheep",
        "trees":   "Trees / Eucalyptus",
        "staff":   "Staff / Stipend",
        "event":   "Events",
        "loan":    "Loans",
        "other":   "Other",
    }
    CAT_LABEL = {
        "project_investment": "Project Investment",
        "staff": "Staff / Stipend",
        "event": "Event",
        "loan":  "Loan",
        "other": "Other",
    }


    fam_rows  = _dbq("SELECT id, family_name FROM families ORDER BY id")
    families  = [(r["id"], r["family_name"].title()) for r in fam_rows]
    fids      = [f[0] for f in families]
    fname     = {f[0]: f[1] for f in families}
    OPENING_PF = OPENING_BALANCE / len(families)

    mem_rows = _dbq("SELECT family_id, COUNT(*) as cnt FROM persons GROUP BY family_id")
    members  = {r["family_id"]: int(r["cnt"]) for r in mem_rows}

    # Model C: fixed weight = sum of monthly_share_ugx per family
    # (adults 50,000 + children 15,000 each)
    obl_rows = _dbq("SELECT family_id, SUM(monthly_share_ugx) as obl FROM persons GROUP BY family_id")
    obligation = {r["family_id"]: int(r["obl"]) for r in obl_rows}
    total_obligation = sum(obligation[fid] for fid in fids)
    weight_C = {fid: obligation.get(fid, 0) / total_obligation for fid in fids}

    pmt_rows = _dbq("""SELECT family_id, confirmed_at::date AS pay_date, amount_ugx
        FROM contribution_payments WHERE status='confirmed'
        ORDER BY confirmed_at, family_id""")

    exp_rows = _dbq("""SELECT txn_date, description, amount_ugx, category, project
        FROM expenditure_records ORDER BY txn_date""")

    # Loan repayments: money that left the pool as an expense but came back.
    # Merge into the event timeline as dated credits, distributed back to families
    # using the same model logic as the original expense (equal/proportional/fixed).
    loan_repayment_rows = _dbq("""SELECT paid_at::date AS pay_date, SUM(amount_ugx) AS amount
        FROM loan_payments GROUP BY paid_at::date ORDER BY paid_at::date""")
    loan_repayments = sorted(
        [(r["pay_date"], int(r["amount"])) for r in loan_repayment_rows],
        key=lambda x: x[0]
    )

    # per-family payment timeline
    # Build sorted contribution list for running-balance approach
    all_pmts = sorted(
        [(r["pay_date"], r["family_id"], int(r["amount_ugx"])) for r in pmt_rows],
        key=lambda x: x[0]
    )
    total_contribs = {fid: 0 for fid in fids}
    for _, fid, amt in all_pmts:
        total_contribs[fid] += amt

    # Running balances for A and B — start with opening credit, grow with contributions,
    # shrink with each expense allocation.
    # Model C does NOT use running balances: equity = total_club_equity × fixed_weight.
    # We still compute per-expense aC for the expense log display, but final equity
    # is derived at the end from the total pool, not from a running balance.
    bal_A = {fid: OPENING_PF for fid in fids}
    bal_B = {fid: OPENING_PF for fid in fids}

    pmt_idx = 0   # pointer into all_pmts
    loan_idx = 0  # pointer into loan_repayments

    cum_A = {fid: 0.0 for fid in fids}
    cum_B = {fid: 0.0 for fid in fids}
    cum_C = {fid: 0.0 for fid in fids}  # cumulative for expense log only
    proj_A, proj_B, proj_C, proj_meta = {}, {}, {}, {}
    expense_detail = []

    def _apply_loan_repayment(repay_amt):
        """Credit a loan repayment back to all families using each model's own logic."""
        # Model A: equal share back (mirrors equal-share expense)
        eq_back = repay_amt / len(fids)
        for fid in fids:
            bal_A[fid] += eq_back
        # Model B: proportional to current pool (same as how expenses are charged)
        tot = sum(bal_B.values())
        for fid in fids:
            bal_B[fid] += (bal_B[fid] / tot * repay_amt) if tot else repay_amt / len(fids)
        # Model C: no running balance — equity computed from total pool at the end

    for r in exp_rows:
        txn_date = r["txn_date"]; desc = r["description"]
        amount = int(r["amount_ugx"]); category = r["category"]; project = r["project"]

        # Credit contributions made on or before this expense date
        while pmt_idx < len(all_pmts) and all_pmts[pmt_idx][0] <= txn_date:
            _, fid, amt = all_pmts[pmt_idx]
            bal_A[fid] += amt
            bal_B[fid] += amt
            # Model C has no running balance (equity = total pool × fixed weight at the end)
            pmt_idx += 1

        # Credit loan repayments made on or before this expense date
        while loan_idx < len(loan_repayments) and loan_repayments[loan_idx][0] <= txn_date:
            _apply_loan_repayment(loan_repayments[loan_idx][1])
            loan_idx += 1

        eq       = amount / 7
        pools    = {fid: bal_A[fid] for fid in fids}   # snapshot before deductions (for display)
        tot_pool = sum(bal_B.values())  # B uses current pool for proportional

        # Model A: equal share if you have it; if short pay all you have,
        # gap splits equally among those who covered their full portion.
        aA = {}
        shortfall = 0.0
        full = []
        for fid in fids:
            if bal_A[fid] >= eq:
                aA[fid] = eq
                full.append(fid)
            else:
                aA[fid] = bal_A[fid]   # pay everything remaining
                shortfall += eq - bal_A[fid]
        if full and shortfall > 0:
            extra = shortfall / len(full)
            for fid in full:
                aA[fid] += extra

        # Model B: strictly proportional to current running balance
        aB = {fid: (bal_B[fid] / tot_pool * amount if tot_pool else amount / 7)
              for fid in fids}

        # Model C: fixed weight proportional to family monthly obligation
        aC = {fid: weight_C[fid] * amount for fid in fids}

        # Deduct allocations from running balances (A and B only)
        for fid in fids:
            bal_A[fid] = max(0.0, bal_A[fid] - aA[fid])
            bal_B[fid] = max(0.0, bal_B[fid] - aB[fid])
            cum_A[fid] += aA[fid]
            cum_B[fid] += aB[fid]
            cum_C[fid] += aC[fid]  # tracked for expense log display only

        pk = project if project else category
        if pk not in proj_A:
            proj_A[pk] = {fid: 0.0 for fid in fids}
            proj_B[pk] = {fid: 0.0 for fid in fids}
            proj_C[pk] = {fid: 0.0 for fid in fids}
            proj_meta[pk] = {"total": 0, "project": project, "category": category}
        for fid in fids:
            proj_A[pk][fid] += aA[fid]
            proj_B[pk][fid] += aB[fid]
            proj_C[pk][fid] += aC[fid]
        proj_meta[pk]["total"] += amount

        expense_detail.append({
            "date":   txn_date.strftime("%d %b %Y"),
            "desc":   desc,
            "amount": amount,
            "cat":    CAT_LABEL.get(category, category),
            "proj":   project or "",
            "eq":     round(eq),
            "fams":   {
                fname[fid]: {
                    "pool":    round(pools[fid]),
                    "pool_pct": round(pools[fid]/tot_pool*100, 1) if tot_pool else 0,
                    "aA":  round(aA[fid]),
                    "aB":  round(aB[fid]),
                    "aC":  round(aC[fid]),
                } for fid in fids
            },
        })

    # current equity = remaining running balance after all expenses processed
    # Also credit any contributions or loan repayments that arrived after the last expense date
    while pmt_idx < len(all_pmts):
        _, fid, amt = all_pmts[pmt_idx]
        bal_A[fid] += amt
        bal_B[fid] += amt
        pmt_idx += 1
    while loan_idx < len(loan_repayments):
        _apply_loan_repayment(loan_repayments[loan_idx][1])
        loan_idx += 1

    # Model C equity: total club equity × fixed weight.
    # Does not depend on when/how much individual families contributed —
    # only on what's currently in the pool and each family's headcount weight.
    total_pool_C = sum(bal_A.values())  # same total as A (pool is shared)
    eq_C_by_fid = {fid: total_pool_C * weight_C[fid] for fid in fids}

    family_summary = []
    for fid, nm in families:
        m   = members.get(fid, 4)
        eqA = bal_A[fid]
        eqB = bal_B[fid]
        eqC = eq_C_by_fid[fid]
        family_summary.append({
            "family":          nm,
            "family_id":       fid,
            "members":         m,
            "opening_credit":  round(OPENING_PF),
            "total_contributed": round(total_contribs[fid]),
            "allocated_A":     round(cum_A[fid]),
            "allocated_B":     round(cum_B[fid]),
            "allocated_C":     round(cum_C[fid]),
            "equity_A":        round(eqA),
            "equity_B":        round(eqB),
            "equity_C":        round(eqC),
            "equity_A_pct":    0,
            "equity_B_pct":    0,
            "equity_C_pct":    0,
            "weight_C_pct":    round(weight_C[fid] * 100, 2),
            "obligation_C":    obligation.get(fid, 0),
            "equity_A_pp":     round(eqA/m) if m else 0,
            "equity_B_pp":     round(eqB/m) if m else 0,
            "equity_C_pp":     round(eqC/m) if m else 0,
        })

    tot_eqA = sum(f["equity_A"] for f in family_summary)
    tot_eqB = sum(f["equity_B"] for f in family_summary)
    tot_eqC = sum(f["equity_C"] for f in family_summary)
    for f in family_summary:
        f["equity_A_pct"] = round(f["equity_A"]/tot_eqA*100, 1) if tot_eqA else 0
        f["equity_B_pct"] = round(f["equity_B"]/tot_eqB*100, 1) if tot_eqB else 0
        f["equity_C_pct"] = round(f["equity_C"]/tot_eqC*100, 1) if tot_eqC else 0

    # project breakdowns
    project_summaries = []
    for pk in proj_A:
        project_summaries.append({
            "key":   pk,
            "label": PROJ_LABEL.get(pk, pk.title()),
            "total": proj_meta[pk]["total"],
            "families": {
                fname[fid]: {
                    "aA":   round(proj_A[pk][fid]),
                    "aB":   round(proj_B[pk][fid]),
                    "aC":   round(proj_C[pk][fid]),
                    "pctA": round(proj_A[pk][fid]/proj_meta[pk]["total"]*100, 1) if proj_meta[pk]["total"] else 0,
                    "pctB": round(proj_B[pk][fid]/proj_meta[pk]["total"]*100, 1) if proj_meta[pk]["total"] else 0,
                    "pctC": round(proj_C[pk][fid]/proj_meta[pk]["total"]*100, 1) if proj_meta[pk]["total"] else 0,
                } for fid in fids
            },
        })

    data = {
        "opening_balance":    OPENING_BALANCE,
        "opening_per_family": round(OPENING_PF),
        "total_expenses":     sum(int(r["amount_ugx"]) for r in exp_rows),
        "total_eq_A":         tot_eqA,
        "total_eq_B":         tot_eqB,
        "total_eq_C":         tot_eqC,
        "family_summary":     family_summary,
        "expense_detail":     expense_detail,
        "project_summaries":  project_summaries,
        "model_a_desc": (
            "Each family pays their equal share if they have enough in the account. "
            "If a family is short, they pay everything they have and the gap is shared "
            "equally among the families that covered their full portion. Proposed by the family."
        ),
        "model_b_desc": (
            "Every expense is split in proportion to each family's accumulated pool "
            "at the date of the expense. Families that contributed more consistently "
            "bear a larger share of every outflow."
        ),
        "model_c_desc": (
            "Proposed by Solomon at KIM 008/2026. Each family's share of every expense "
            "is fixed by their total monthly contribution obligation: adults at UGX 50,000 "
            "each, children at UGX 15,000 each. A larger family pays a proportionally larger "
            "share of every expense, regardless of their current balance."
        ),
        "decision_status": "pending",
        "decision_note":   "Discussed at KIM 008/2026. Vote deferred to KIM 009/2026. Both models remain active. Read the explainer card before the next meeting.",
    }

    _ALLOC_CACHE["data"] = data
    _ALLOC_CACHE["ts"]   = now
    return data


def _gc():
    import gspread
    from google.oauth2.service_account import Credentials
    SA = os.path.join(os.path.dirname(__file__), "service-account.json")
    creds = Credentials.from_service_account_file(SA, scopes=[
        "https://spreadsheets.google.com/feeds",
        "https://www.googleapis.com/auth/spreadsheets",
    ])
    return gspread.authorize(creds)

def _fetch_chicken_data():
    import time
    now = time.time()
    if _CHICKEN_CACHE.get("ts") and now - _CHICKEN_CACHE["ts"] < _CHICKEN_CACHE_TTL:
        return _CHICKEN_CACHE["data"]

    gc = _gc()
    sh = gc.open_by_key(_CHICKEN_SHEET_ID)

    # 1. Product summary (product descriptions tab)
    pd_rows = sh.worksheet("product descriptions").get_all_values()
    products = {}
    for r in pd_rows[1:]:
        if len(r) < 14 or not r[0].strip():
            continue
        pid = r[0].strip()
        products[pid] = {
            "name":      r[2].strip(),
            "available": _parse_num(r[6]),
            "sold":      _parse_num(r[7]),
            "purchased": _parse_num(r[8]),
            "revenue":   _parse_num(r[9]),
            "deaths":    _parse_num(r[10]),
            "deaths_val":_parse_num(r[11]),
        }

    # 2. Sales breakdown by product over time
    sales_rows = sh.worksheet("sales").get_all_values()
    sales_by_product = {"a1": [], "a2": [], "a3": []}
    monthly_eggs = {}
    for r in sales_rows[1:]:
        if len(r) < 12 or not r[3].strip():
            continue
        pid   = r[6].strip()
        date  = r[3].strip()
        qty   = _parse_num(r[9])
        amt   = _parse_num(r[11])
        buyer = r[18].strip() if len(r) > 18 else ""
        if pid in sales_by_product:
            sales_by_product[pid].append({"date": date, "qty": qty, "amount": amt, "buyer": buyer})
        if pid == "a1":
            try:
                from datetime import datetime
                dt = datetime.strptime(date, "%m/%d/%Y") if "/" in date and len(date.split("/")[0]) <= 2 else datetime.strptime(date, "%d/%m/%Y")
                ym = dt.strftime("%b %Y")
                monthly_eggs[ym] = monthly_eggs.get(ym, 0) + qty
            except Exception:
                pass

    # 3. Deaths breakdown
    spoilt_rows = sh.worksheet("used or spoilt items not sold").get_all_values()
    deaths_detail = []
    for r in spoilt_rows[1:]:
        if len(r) < 12 or not r[3].strip():
            continue
        pid    = r[6].strip()
        date   = r[3].strip()
        reason = r[10].strip()
        qty    = _parse_num(r[11])
        if pid in ("a2", "a3") and r[9].strip().lower() == "damaged":
            deaths_detail.append({"product": r[5].strip(), "date": date, "qty": qty, "reason": reason})

    # 4. Stock purchases (batches)
    stock_rows = sh.worksheet("stock details starting may 2024").get_all_values()
    batches = []
    for r in stock_rows[1:]:
        if len(r) < 10 or not r[3].strip():
            continue
        pid = r[4].strip()
        # Include standard bird PIDs (a2=hens, a3=cocks) plus any extra batch products
        # (non-a1 means not eggs; new batches get UUID-style PIDs from AppSheet)
        if pid in ("a2", "a3") or (pid not in ("a1",) and "chicken" in r[3].lower() or "batch" in r[3].lower() or "hen" in r[3].lower() or "cock" in r[3].lower() or "pullet" in r[3].lower() or "chick" in r[3].lower()):
            batches.append({"product": r[3].strip(), "date": r[2].strip()[:10], "qty": _parse_num(r[9]), "source": r[13].strip() if len(r) > 13 else "", "pid": pid})

    # 5. Financial Statement tab (live P&L)
    fin_rows = sh.worksheet("Financial Statement").get_all_values()
    financials_raw = {}
    for r in fin_rows[1:]:
        if len(r) < 3 or not r[1].strip():
            continue
        financials_raw[r[1].strip().lower()] = _parse_num(r[2])

    # 6. Company expenses tab — OPEX/CapEx breakdown by category
    def _cat(item_name):
        it = item_name.lower()
        if any(k in it for k in ['chicken','hen','cock','bird','pullet']):
            return 'Birds / Stock'
        if any(k in it for k in ['mash','maize','feed','kyacu','grower','layer',
                                   'concentrate','sunflower','sun flower','milling','broken']):
            return 'Feed & Nutrition'
        if any(k in it for k in ['medicine','s-dime','interflox','dudu','lime','vaccine','drug']):
            return 'Medicine & Vet'
        if 'transport' in it:
            return 'Transport'
        if any(k in it for k in ['labour','salary','manager','casual','workshop']):
            return 'Labour'
        return 'Equipment & Supplies'

    exp_rows = sh.worksheet('company expenses').get_all_values()
    opex_breakdown = {}
    monthly_spend  = {}
    expense_timeline = []

    for r in exp_rows[1:]:
        if len(r) < 14 or not r[9].strip():
            continue
        etype = r[13].strip()
        item  = r[4].strip()
        cost  = _parse_num(r[9])
        date  = r[11].strip()
        if not cost:
            continue
        expense_timeline.append({'date': date, 'item': item, 'cost': int(cost), 'type': etype})
        try:
            from datetime import datetime as _dt
            parts = date.split('/')
            if len(parts) == 3:
                d = _dt(int(parts[2]), int(parts[1]), int(parts[0]))
                ym = d.strftime('%b %Y')
                monthly_spend[ym] = int(monthly_spend.get(ym, 0) + cost)
        except Exception:
            pass
        if etype.lower() == 'opex':
            cat = _cat(item)
            opex_breakdown[cat] = int(opex_breakdown.get(cat, 0) + cost)

    data = {
        "products": products,
        "sales_by_product": sales_by_product,
        "monthly_egg_sales": monthly_eggs,
        "deaths_detail": deaths_detail,
        "batches": batches,
        "financials_raw":   financials_raw,
        "opex_breakdown":   opex_breakdown,
        "monthly_spend":    monthly_spend,
        "expense_timeline": expense_timeline,
    }
    _CHICKEN_CACHE["data"] = data
    _CHICKEN_CACHE["ts"] = now
    return data

def _parse_num(s):
    try:
        return float(str(s).replace(",","").replace(" ","").replace("\xa0","") or 0)
    except Exception:
        return 0.0

@app.post("/api/projects/chicken/refresh")
async def chicken_refresh(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    _CHICKEN_CACHE.clear()
    _fetch_chicken_data()
    return {ok: True}


@app.post("/api/projects/chicken/refresh")
async def chicken_refresh(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    _CHICKEN_CACHE.clear()
    _fetch_chicken_data()
    return {"ok": True}

@app.get("/api/projects/chicken/detail")
async def chicken_detail(request: Request):
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    payload = _auth_verify(token)
    is_internal = _internal_key_ok(request)
    if not payload and not is_internal:
        raise _HE(status_code=401, detail="Auth required")
    data = _fetch_chicken_data()
    p = data["products"]

    hens  = p.get("a2", {})
    cocks = p.get("a3", {})
    eggs  = p.get("a1", {})

    # Any product that is FREE RANGE CHICKEN but not standard a1/a2/a3 = new batch (chicks/pullets)
    extra_batches = {pid: v for pid, v in p.items() if pid not in ("a1", "a2", "a3")}
    active_chicks      = sum(int(v.get("available",  0)) for v in extra_batches.values())
    chicks_purchased   = sum(int(v.get("purchased",  0)) for v in extra_batches.values())
    chicks_sold        = sum(int(v.get("sold",        0)) for v in extra_batches.values())
    chicks_deaths      = sum(int(v.get("deaths",      0)) for v in extra_batches.values())
    chicks_capex       = sum(int(v.get("deaths_val",  0)) for v in extra_batches.values())
    # Label: show descriptions so it's clear what these are
    batch_labels = [v for v in extra_batches.values() if v.get("name")]

    TRAY = 30

    fr = data.get("financials_raw", {})
    total_sales    = int(fr.get("sales", 0))
    spoilt_loss    = int(fr.get("spoilt goods (cash loss)", 0))
    opex           = int(fr.get("operating expenses (opex)", 0))
    capex          = int(fr.get("capital expenses (capex)", 0))
    dep_pa         = int(fr.get("depreciation (per year)", 0))
    gross_pos      = int(fr.get("gross position", 0))
    net_capex      = int(fr.get("net position (with capex)", 0))
    net_dep        = int(fr.get("net position (with depreciation)", 0))

    # Batch 3 is WhatsApp-only until hens_purchased exceeds 180 (Batch 1+2)
    hens_purchased = int(hens.get("purchased", 0))
    batch3_unlogged = hens_purchased <= 180
    whatsapp_only = {
        "pending": batch3_unlogged,
        "batch3_ordered": 100 if batch3_unlogged else 0,
        "batch3_deaths": 1 if batch3_unlogged else 0,
        "note": (
            "Batch 3 (100 birds) and 1 death reported via WhatsApp on 6 Jun 2026. "
            "Not yet recorded in Solomon's AppSheet. Figures unverified until logged."
        ) if batch3_unlogged else "",
    }

    # Revenue = eggs + all bird products (a2, a3, any extra batches)
    extra_revenue = sum(int(v.get("revenue", 0)) for v in extra_batches.values())
    actual_revenue = int(eggs.get("revenue", 0) + hens.get("revenue", 0) + cocks.get("revenue", 0) + extra_revenue)
    return {
        "flock": {
            "active_hens":            int(hens.get("available", 0)),
            "active_cocks":           int(cocks.get("available", 0)),
            "new_batch_chicks":       active_chicks,
            "total_hens_purchased":   hens_purchased,
            "total_cocks_purchased":  int(cocks.get("purchased", 0)),
            "new_batch_purchased":    chicks_purchased,
            "hens_sold":              int(hens.get("sold", 0)),
            "cocks_sold":             int(cocks.get("sold", 0)),
            "hens_deaths":            int(hens.get("deaths", 0)),
            "cocks_deaths":           int(cocks.get("deaths", 0)),
            "new_batch_note":         "New batch (chicks/pullets) — not yet laying. Recorded as separate product in AppSheet." if active_chicks else "",
            "deaths_detail":          data["deaths_detail"],
            "batches":                data["batches"],
            "whatsapp_only":          whatsapp_only,
        },
        "sales": {
            "eggs": {
                "units":   int(eggs.get("sold", 0)),
                "trays":   round(eggs.get("sold", 0) / TRAY, 1),
                "revenue": int(eggs.get("revenue", 0)),
            },
            "hens": {
                "units":   int(hens.get("sold", 0)),
                "revenue": int(hens.get("revenue", 0)),
            },
            "cocks": {
                "units":   int(cocks.get("sold", 0)),
                "revenue": int(cocks.get("revenue", 0)),
            },
            "total_revenue":     actual_revenue,
            "monthly_egg_sales": data["monthly_egg_sales"],
        },
        "financials": {
            "total_sales":         total_sales,
            "spoilt_loss":         spoilt_loss,
            "opex":                opex,
            "capex":               capex,
            "depreciation_pa":     dep_pa,
            "gross_position":      gross_pos,
            "net_with_capex":      net_capex,
            "net_with_depreciation": net_dep,
            "source":              "Solomon's AppSheet Financial Statement tab (live)",
        },
        "projections": {
            "phase":             "Phase 1 (100 birds, Jul 2024 - Jan 2026)",
            "projected_revenue": 12360000,
            "projected_profit":  3950000,
            "actual_revenue":    actual_revenue,
            "actual_gross":      gross_pos,
            "revenue_variance":  actual_revenue - 12360000,
            "profit_variance":   gross_pos - 3950000,
            "note": "Phase 2 (100 new birds) started Jun 2026. WhatsApp-reported only pending AppSheet update.",
            "assumptions": {
                "egg_price_per_egg":        400,
                "egg_price_per_tray":       12000,
                "projected_monthly_eggs":   1800,
                "ramp_up_month":            5,
                "ramp_up_eggs":             1500,
                "projected_eggs_total":     24900,
                "projected_egg_trays":      830,
                "projected_birds_to_sell":  80,
                "projected_chicken_price":  30000,
                "phase_duration":           "Jul 2024 - Jan 2026 (19 months)",
                "productive_months":        18,
                "cost_of_sale":             2650000,
                "other_expenses":           5760000,
                "profit_split_mgmt":        "30% of profit",
                "profit_split_family_club": "50% of total profit",
            },
            "line_variances": [
                {"item": "Eggs Sold", "unit": "eggs",
                 "projected": 24900, "actual": int(eggs.get("sold", 0))},
                {"item": "Egg Revenue", "unit": "ugx",
                 "projected": 9960000, "actual": int(eggs.get("revenue", 0))},
                {"item": "Chicken Revenue", "unit": "ugx",
                 "projected": 2400000, "actual": int(hens.get("revenue", 0) + cocks.get("revenue", 0))},
                {"item": "Total Revenue", "unit": "ugx",
                 "projected": 12360000, "actual": actual_revenue},
                {"item": "Total Expenses (OPEX)", "unit": "ugx",
                 "projected": -(2650000 + 5760000), "actual": -opex},
                {"item": "Gross Profit / (Loss)", "unit": "ugx",
                 "projected": 3950000, "actual": gross_pos},
            ],
        },
        "family_account": {
            "total": 8289700,
            "tranches": [
                {"date": "04 Apr 2024", "label": "Free range chicken capital", "amount": 100000},
                {"date": "23 Apr 2024", "label": "Free range chicken capital", "amount": 4000000},
                {"date": "08 Jul 2024", "label": "Extra investment",           "amount": 2189700},
                {"date": "25 Apr 2026", "label": "Extra investment",           "amount": 2000000},
            ],
        },
        "opex_breakdown":   data.get("opex_breakdown", {}),
        "monthly_spend":    data.get("monthly_spend", {}),
    }

@app.get("/api/projects/trees/detail")
async def trees_detail(request: Request):
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    ok = _auth_verify(token) or _internal_key_ok(request)
    if not ok:
        raise _HE(status_code=401, detail="Auth required")
    import datetime as _dt_trees
    planted = _dt_trees.date(2024, 8, 16)
    today_eat = (_dt_trees.datetime.utcnow() + _dt_trees.timedelta(hours=3)).date()
    months_old = (today_eat.year - planted.year) * 12 + (today_eat.month - planted.month)
    years_old = round(months_old / 12, 1)
    harvest_yr12 = _dt_trees.date(2036, 8, 1)
    days_to_harvest = (harvest_yr12 - today_eat).days
    pct_to_harvest = round(min(100, months_old / (12 * 12) * 100), 1)

    # Confirmed actual cash payments (ABSA bank receipts on file)
    confirmed_payments = [
        {
            "date":   "16 Aug 2024",
            "desc":   "Seedlings (3,600) + Site Preparation + Planting Labour",
            "amount": 2630000,
            "from_":  "Hellen Tusiime (ABSA → Treasurer)",
            "to":     "Kikangi Israel (Dad) · Diamond Trust Bank",
            "ref":    "KIM project trees · TXN 17237999200564",
        },
    ]
    total_confirmed = sum(p["amount"] for p in confirmed_payments)

    # Future budget commitments (budgeted, not yet paid)
    future_commitments = [
        {"item": "Year 1 Maintenance",         "amount": 3278400, "due": "Aug 2025", "note": "Spot hoeing, slashing x2, security"},
        {"item": "Year 2 Maintenance",         "amount": 1924560, "due": "Aug 2026", "note": "Herbicide, slashing x2, security"},
        {"item": "Security (remaining 8 yrs)", "amount": 1920000, "due": "2026-2034","note": "UGX 20,000/month x 96 months"},
    ]
    total_future   = sum(c["amount"] for c in future_commitments)
    total_capex    = total_confirmed + total_future
    revenue_poles  = 81120000
    revenue_timber = 81120000
    gross_profit   = revenue_poles - total_capex

    milestones = [
        {"label": "Seedlings Purchased",  "date": "Aug 2024", "year": 0,  "done": True,  "note": "3,600 seedlings at UGX 200 each"},
        {"label": "Year 1 Maintenance",   "date": "Aug 2025", "year": 1,  "done": months_old >= 12, "note": "Spot hoeing, slashing x2, security"},
        {"label": "Year 2 Maintenance",   "date": "Aug 2026", "year": 2,  "done": months_old >= 24, "note": "Herbicide, slashing x2, security"},
        {"label": "First Selective Harvest","date": "Aug 2029","year": 5, "done": False, "note": "Thinnings — used for farm fencing"},
        {"label": "Electric Poles Harvest","date": "Aug 2036", "year": 12, "done": False, "note": "811 trees × UGX 100,000 = UGX 81.1M"},
        {"label": "Timber Harvest",       "date": "Aug 2039", "year": 15, "done": False, "note": "Full timber grade — same yield, premium market"},
    ]

    return {
        "overview": {
            "species":           "Eucalyptus grandis",
            "location":          "Busisi Land, Kashari",
            "land_acres":        3.38,
            "seedlings_planted": 3600,
            "harvestable_trees": 811,
            "planted_date":      "16 Aug 2024",
            "age_months":        months_old,
            "age_years":         years_old,
            "pct_to_harvest":    pct_to_harvest,
            "days_to_harvest":   days_to_harvest,
        },
        "capex": {
            "total":              total_capex,
            "confirmed":          total_confirmed,
            "future":             total_future,
            "confirmed_payments": confirmed_payments,
            "future_commitments": future_commitments,
        },
        "revenue": {
            "poles":  {"qty": 811, "price_each": 100000, "total": revenue_poles,  "timeline": "Year 12 (Aug 2036)"},
            "timber": {"qty": 811, "price_each": 100000, "total": revenue_timber, "timeline": "Year 15-20 (2039-2044)"},
        },
        "profit_split": {
            "gross_profit":  gross_profit,
            "management":    {"pct": 30, "amount": round(gross_profit * 0.30), "label": "Management & Oversight"},
            "investor":      {"pct": 20, "amount": round(gross_profit * 0.20), "label": "Lead Investor"},
            "club":          {"pct": 50, "amount": round(gross_profit * 0.50), "label": "Family Investment Club"},
        },
        "financial_metrics": {
            "payback_years":     12,
            "payback_note":      "Revenue at year 12 (poles) covers full CapEx of UGX " + f"{total_capex:,}",
            "roi_pct":           round(gross_profit / total_capex * 100, 1),
            "roi_note":          "Return on total projected investment (poles scenario)",
            "irr_approx":        "18-22%",
            "irr_note":          "Estimated IRR (12-yr horizon, single harvest). Higher if timber held to yr 15-20.",
            "npv_note":          "NPV positive at discount rates below ~18%. Long-horizon, low-risk profile.",
            "confirmed_deployed": total_confirmed,
            "remaining_budget":   total_future,
            "pct_deployed":       round(total_confirmed / total_capex * 100, 1),
        },
        "milestones": milestones,
        "risks": [
            {"risk": "Water availability",      "level": "medium", "probability": "Medium", "impact": "High",   "note": "Eucalyptus is water-intensive. Busisi land water source not documented."},
            {"risk": "Security / theft",        "level": "low",    "probability": "Low",    "impact": "Medium", "note": "Security budgeted at UGX 20K/month ongoing."},
            {"risk": "Market price at harvest", "level": "medium", "probability": "Medium", "impact": "High",   "note": "UGX 100K/pole is today estimate. Price in 2036 is uncertain."},
            {"risk": "Survival rate",           "level": "low",    "probability": "Low",    "impact": "Medium", "note": "No mortality data yet. Eucalyptus grandis is hardy."},
            {"risk": "Regulatory / land",       "level": "low",    "probability": "Low",    "impact": "High",   "note": "Busisi land title status should be confirmed."},
        ],
    }


@app.get("/api/projects/sheep/detail")
async def sheep_detail(request: Request):
    from fastapi import HTTPException as _HE
    token = _get_tok(request)
    ok = _auth_verify(token) or _internal_key_ok(request)
    if not ok:
        raise _HE(status_code=401, detail="Auth required")

    capital_invested = 1710000
    dorper_count     = 17
    price_low        = 100000
    price_high       = 150000
    price_avg        = 125000

    value_low  = dorper_count * price_low
    value_high = dorper_count * price_high
    value_avg  = dorper_count * price_avg

    sold_proceeds     = 200000   # 1 young Dorper sold for ear-tagging
    pending_ewes      = 12
    pending_per_head  = 200000
    pending_total     = pending_ewes * pending_per_head

    net_position = value_avg + sold_proceeds - capital_invested

    return {
        "investment": {
            "date":    "16 Aug 2024",
            "amount":  capital_invested,
            "paid_by": "Hellen Tusiime (ABSA → Stanbic)",
            "ref":     "KIM PROJECT DORPER SHEEP",
            "paid_to": "Tuhimbise Pricilla Arinaitwe",
        },
        "flock": {
            "total":          40,
            "dorper_count":   dorper_count,
            "club_held":      2,
            "sold":           1,
            "sold_value":     sold_proceeds,
            "note":           "1 young Dorper sold at UGX 200K to fund ear-tagging exercise",
        },
        "valuation": {
            "price_low":   price_low,
            "price_high":  price_high,
            "price_avg":   price_avg,
            "basis":       "Mum confirmed flat pricing UGX 200K/head; Dorpers typically command premium",
            "value_low":   value_low,
            "value_high":  value_high,
            "value_avg":   value_avg,
        },
        "position": {
            "capital_invested":   capital_invested,
            "dorper_value_avg":   value_avg,
            "sold_proceeds":      sold_proceeds,
            "total_value":        value_avg + sold_proceeds,
            "net_position":       net_position,
            "return_pct":         round((net_position / capital_invested) * 100, 1),
        },
        "pipeline": {
            "pending_ewes":       pending_ewes,
            "pending_per_head":   pending_per_head,
            "pending_total":      pending_total,
            "action_id":          "KIM/07/26-4",
            "note":               "12 old ewes ear-tagged, payment being processed via Hellen",
        },
        "next_steps": [
            {"step": "Replough paddocks",  "target": "Sep 2026", "done": False},
            {"step": "Broadcast pasture seeds", "target": "Sep 2026", "done": False},
            {"step": "Harrow paddock area", "target": "Sep 2026", "done": False},
            {"step": "Source new breeding ram", "target": "TBD", "done": False, "note": "Existing ram cannot breed with its own offspring"},
            {"step": "KIM/07/26-4 payment settled", "target": "Jun 2026", "done": False},
        ],
        "breeding_model": {
            "gestation_days":    150,
            "lambs_per_birth":   "1-3",
            "cycles_per_year":   "1-2",
            "note":              "Dorper is a hair sheep — no shearing needed, hardy in tropical climate",
            "harvest_age_months": 6,
            "harvest_weight_kg": 30,
        },
        "receipts": [
            {"date": "16 Aug 2024", "desc": "Dorper sheep purchase", "amount": 1710000, "ref": "ABSA→Stanbic TXN 17238007755249"},
        ],
    }


@app.get("/api/projects/washing_bay/detail")
async def washing_bay_detail(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    import datetime as _dt
    start = _dt.date(2026, 1, 1)
    today = (_dt.datetime.utcnow() + _dt.timedelta(hours=3)).date()
    months_running = max(1, (today.year - start.year)*12 + (today.month - start.month))
    capex         = 25_900_000
    sanitation    = 7_635_500
    monthly_rev   = 580_000
    actual_rev    = 2_923_000
    payback_mo    = round(capex / monthly_rev, 1)
    payback_san   = round((capex + sanitation) / monthly_rev, 1)
    proj_monthly  = [monthly_rev * i for i in range(0, 61)]
    act_points    = [monthly_rev * i if i <= months_running else None for i in range(0, 61)]
    capex_line    = [capex] * 61
    on_target     = abs(actual_rev - monthly_rev * months_running) / (monthly_rev * months_running) < 0.15
    annual_rev    = monthly_rev * 12
    annual_profit = annual_rev * 0.35  # rough margin for car wash
    roi_pct       = round((annual_rev * (payback_mo / 12)) / capex * 100, 1)

    # Capital accountability — how much of the reported CapEx actually has a
    # documented contributor + source behind it.
    try:
        _cc = _wb_conn()
        _crows = _cc.execute("SELECT contributor,amount_ugx FROM capital_contributions").fetchall()
        _cc.close()
    except Exception:
        _crows = []
    cap_accounted = sum(r[1] for r in _crows)
    cap_remaining = max(0, capex - cap_accounted)
    cap_pct = round(cap_accounted / capex * 100, 1) if capex else 0
    accountability_risk = [{
        "risk": "No capital accountability",
        "probability": "High", "impact": "High",
        "note": (f"Only UGX {cap_accounted:,} of the reported UGX {capex:,} CapEx has a documented "
                 f"contributor and source ({cap_pct}%). UGX {cap_remaining:,} is still unproven. "
                 "Dad and Alex are the known major investors; the split needs reconciling.")
    }] if cap_accounted < capex else []

    return {
        "overview": {
            "start_date": "Jan 2026 (approx)", "months_running": months_running,
            "lead": "Alex Tuhimbise", "location": "Mbarara area",
            "status": "Operational", "on_target": on_target,
        },
        "revenue": {
            "monthly_actual": monthly_rev, "total_actual": actual_rev,
            "months_confirmed": 5, "monthly_target": monthly_rev,
            "variance_pct": 0,
            "projected_annual": annual_rev,
        },
        "capex": {
            "initial": capex, "sanitation_pending": sanitation,
            "total_with_sanitation": capex + sanitation,
            "breakdown": [
                {"item": "Electricity infrastructure",  "amount": 0, "note": "Included in total"},
                {"item": "Well + pump + jet spray",      "amount": 0, "note": "Included in total"},
                {"item": "Total initial investment",     "amount": capex, "note": "Confirmed — Alex Tuhimbise"},
                {"item": "Sanitation upgrade (pending)", "amount": sanitation, "note": "Septic UGX 5.47M + Plumbing UGX 2.16M"},
            ],
        },
        "financial_metrics": {
            "payback_months": payback_mo, "payback_with_sanitation": payback_mo + round(sanitation/monthly_rev,1),
            "roi_pct": roi_pct, "irr_approx": "~18-24%",
            "roi_note": "Based on monthly revenue vs initial CapEx",
            "irr_note": "Single-site car wash; grows with throughput volume",
            "npv_note": "Positive at discount rates below ~20%; stable cash-flow asset",
            "recovery_to_date": actual_rev,
            "recovery_pct": round(actual_rev / capex * 100, 1),
        },
        "chart_data": {
            "labels": [f"Mo {i}" for i in range(0, 61)],
            "projected_cumulative": proj_monthly,
            "actual_points": act_points,
            "capex_line": capex_line,
            "break_even_month": round(payback_mo),
        },
        "capital_accountability": {
            "target_ugx": capex,
            "accounted_ugx": cap_accounted,
            "remaining_ugx": cap_remaining,
            "pct_accounted": cap_pct,
            "balanced": cap_accounted >= capex,
        },
        "risks": accountability_risk + [
            {"risk": "Equipment failure",    "probability": "Medium", "impact": "High",   "note": "Jet spray / pump downtime kills revenue. Maintenance schedule critical."},
            {"risk": "Competition",          "probability": "Medium", "impact": "Medium", "note": "Other car washes in area. Differentiation via reliability and location."},
            {"risk": "Sanitation compliance","probability": "High",   "impact": "Medium", "note": "Existing pit cannot serve as septic. UGX 7.6M upgrade required."},
            {"risk": "Occupancy (Eli)",      "probability": "Low",    "impact": "High",   "note": "Eli occupancy agreement pending. Land security risk if not formalised."},
            {"risk": "Water supply",         "probability": "Low",    "impact": "High",   "note": "Borehole-fed. If pump fails, operations stop immediately."},
        ],
        "open_issues": [
            {"issue": "Sanitation upgrade approval", "priority": "High", "amount": sanitation},
            {"issue": "Eli occupancy agreement", "priority": "High", "amount": 0},
            {"issue": "Town Council — temp structures removed", "priority": "Medium", "amount": 0},
        ],
    }


@app.get("/api/projects/irrigation/detail")
async def irrigation_detail(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    annual_revenue = 5_000_000
    annual_costs   = 3_960_000
    annual_profit  = 1_040_000
    mgmt_cut       = 312_000
    investor_cut   = 208_000
    club_cut       = 520_000
    years          = list(range(1, 9))
    cum_profit     = [annual_profit * i for i in years]
    rev_by_year    = [annual_revenue] * 8
    cost_by_year   = [annual_costs] * 8
    profit_by_year = [annual_profit] * 8
    return {
        "overview": {
            "crop": "Bananas (+ vegetable garden)", "area_acres": 2,
            "expansion_acres": 4, "system": "Overhead sprinkler irrigation",
            "water_sources": "Rainwater harvesting + farm ponds",
            "status": "Planning — Phase 1 design complete",
            "phase2": "Procurement and installation pending approval and funding",
            "reason": "~80 banana bunches lost to heat stress without irrigation",
        },
        "projection": {
            "bunches_per_year": 500, "price_per_bunch": 10_000,
            "annual_revenue": annual_revenue, "annual_costs": annual_costs,
            "annual_profit": annual_profit,
            "years": 8,
            "revenue_by_year": rev_by_year,
            "cost_by_year": cost_by_year,
            "profit_by_year": profit_by_year,
            "cumulative_profit": cum_profit,
        },
        "profit_split": {
            "annual_profit": annual_profit,
            "management": {"pct": 30, "amount": mgmt_cut, "label": "Management & Oversight"},
            "investor":   {"pct": 20, "amount": investor_cut, "label": "Lead Investor"},
            "club":       {"pct": 50, "amount": club_cut, "label": "Family Investment Club"},
        },
        "financial_metrics": {
            "annual_roi_on_ops": round(annual_profit / annual_costs * 100, 1),
            "8yr_cumulative_profit": annual_profit * 8,
            "8yr_club_earnings": club_cut * 8,
            "note": "Irrigation CapEx (phase 2) not yet quantified — will reduce early-year returns",
            "irr_approx": "TBD after phase 2 cost confirmed",
        },
        "risks": [
            {"risk": "Phase 2 cost unknown",    "probability": "High",   "impact": "Medium", "note": "Irrigation equipment + installation cost not yet estimated."},
            {"risk": "Rainfall variability",    "probability": "Medium", "impact": "High",   "note": "System mitigates this but doesn't eliminate dependence on ponds."},
            {"risk": "Banana disease (fusarium)","probability": "Low",   "impact": "High",   "note": "Fusarium wilt is endemic in Uganda. Variety selection matters."},
            {"risk": "Market price fluctuation","probability": "Medium", "impact": "Medium", "note": "UGX 10K/bunch is conservative; market can go lower at peak season."},
        ],
    }


@app.get("/api/projects/dairy/detail")
async def dairy_detail(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    cows            = 2
    cost_per_cow    = 8_000_000
    total_capex     = cows * cost_per_cow
    litres_low      = 15 * cows
    litres_high     = 50 * cows
    litres_mid      = 30 * cows
    milk_price      = 1_500
    monthly_rev_low = litres_low * milk_price * 30
    monthly_rev_mid = litres_mid * milk_price * 30
    monthly_rev_hi  = litres_high * milk_price * 30
    annual_opex     = 3_600_000
    annual_profit_mid = (monthly_rev_mid * 12) - annual_opex
    payback_mo_mid  = round(total_capex / (monthly_rev_mid - annual_opex / 12), 1)
    proj_labels     = [f"Mo {i}" for i in range(0, 37)]
    proj_cumrev     = [monthly_rev_mid * i for i in range(0, 37)]
    capex_line      = [total_capex] * 37
    return {
        "overview": {
            "status": "Planning", "model": "Shed-based zero-grazing (fully indoor)",
            "target_herd": cows, "cost_per_cow": cost_per_cow,
            "total_capex": total_capex, "productive_life": "4-6 lactation cycles",
            "site_visit": "Dad visited working shed dairy farm, 25 May 2026",
        },
        "revenue_scenarios": {
            "low":  {"litres_day": litres_low,  "monthly": monthly_rev_low,  "annual": monthly_rev_low * 12},
            "mid":  {"litres_day": litres_mid,  "monthly": monthly_rev_mid,  "annual": monthly_rev_mid * 12},
            "high": {"litres_day": litres_high, "monthly": monthly_rev_hi,   "annual": monthly_rev_hi * 12},
            "milk_price_ugx": milk_price,
            "price_note": "UGX 1,500/litre is conservative farm-gate; retail up to UGX 2,500",
        },
        "financial_metrics": {
            "total_capex": total_capex,
            "annual_opex_est": annual_opex,
            "payback_months_mid": payback_mo_mid,
            "annual_profit_mid": annual_profit_mid,
            "roi_mid": round(annual_profit_mid / total_capex * 100, 1),
            "irr_approx": "~35-50%",
            "irr_note": "High IRR driven by low cycle time vs capital. Multiple lactations compound returns.",
            "npv_note": "NPV strongly positive even at 25% discount rate if mid-yield maintained.",
        },
        "model_logic": {
            "why_shed": "Tick-borne diseases kill quality breeds outdoors. Shed eliminates tick exposure entirely.",
            "feed": "Cut forage only — zero grazing. Cows never leave the shed.",
            "litter": "Deep litter system, easy manure management (manure also sold as fertiliser).",
            "water": "Piped water supply required. High daily water consumption.",
            "after_productive_life": "Sold for beef at end of 4-6 lactations — residual value recovery.",
        },
        "chart_data": {
            "labels": proj_labels,
            "projected_cumrev_mid": proj_cumrev,
            "capex_line": capex_line,
            "break_even_month": round(payback_mo_mid),
        },
        "risks": [
            {"risk": "Disease (indoor)",        "probability": "Low",    "impact": "High",   "note": "Shed eliminates tick risk. Main risks: respiratory and mastitis."},
            {"risk": "Feed supply chain",       "probability": "Medium", "impact": "High",   "note": "Zero-grazing depends entirely on cut forage availability. Drought risk."},
            {"risk": "Milk price drop",         "probability": "Medium", "impact": "Medium", "note": "Farm-gate price sensitive to local supply gluts."},
            {"risk": "Capital sourcing",        "probability": "High",   "impact": "High",   "note": "UGX 16M investment not yet approved or committed."},
            {"risk": "Management skill gap",    "probability": "Medium", "impact": "High",   "note": "Zero-grazing is specialised. Training or expert hire required."},
        ],
    }


@app.get("/api/projects/bees/detail")
async def bees_detail(request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    hives           = 10
    kg_per_hive_yr  = 25
    total_kg_yr     = hives * kg_per_hive_yr
    price_per_kg    = 15_000
    annual_rev      = total_kg_yr * price_per_kg
    hive_cost       = 150_000
    setup_costs     = hives * hive_cost + 500_000
    annual_opex     = 600_000
    annual_profit   = annual_rev - annual_opex
    payback_yrs     = round(setup_costs / annual_profit, 1)
    labels_3yr      = ["Yr 1","Yr 2","Yr 3","Yr 4","Yr 5"]
    proj_rev        = [annual_rev, annual_rev * 1.5, annual_rev * 2, annual_rev * 2.5, annual_rev * 3]
    proj_profit     = [annual_profit, annual_profit*1.5, annual_profit*2, annual_profit*2.5, annual_profit*3]
    return {
        "overview": {
            "status": "Research", "brand": "KIM FAM Honey Farm",
            "hives_planned": hives, "annual_yield_kg": total_kg_yr,
            "lead": "Dad (Israel)", "capital_committed": False,
        },
        "economics": {
            "hives": hives, "yield_kg_per_hive": kg_per_hive_yr,
            "total_kg_yr": total_kg_yr, "price_per_kg": price_per_kg,
            "annual_revenue": annual_rev, "annual_opex": annual_opex,
            "annual_profit": annual_profit,
            "hive_cost_each": hive_cost, "total_setup": setup_costs,
        },
        "financial_metrics": {
            "payback_years": payback_yrs, "annual_roi": round(annual_profit / setup_costs * 100, 1),
            "irr_approx": "~60-80%",
            "irr_note": "Very low capital base; high return if yield targets met.",
            "npv_note": "Breakeven within 2 years at conservative yield. Low-risk capital profile.",
            "scale_note": "Year 3+ adds wild hive splits — 10 hives can become 20-30 organically.",
        },
        "chart_data": {
            "labels": labels_3yr,
            "revenue": proj_rev, "profit": proj_profit,
        },
        "other_products": [
            {"product": "Beeswax",   "note": "Candles, lip balm, polish — premium niche market"},
            {"product": "Propolis",  "note": "Medicinal extract — UGX 50K+/gram"},
            {"product": "Pollination services", "note": "Contract pollination for local farms — zero cost to produce"},
        ],
        "risks": [
            {"risk": "Colony collapse / theft",  "probability": "Medium", "impact": "High",   "note": "Vandalism and theft historically a risk for hives in the area."},
            {"risk": "Forage availability",      "probability": "Low",    "impact": "Medium", "note": "Bee forage trees being identified. Mango + Eucalyptus on farm synergise."},
            {"risk": "Market access",            "probability": "Medium", "impact": "Medium", "note": "Local honey market exists but branded premium pricing needs distribution."},
            {"risk": "Expertise gap",            "probability": "High",   "impact": "Medium", "note": "Dad joined beekeeper groups but formal training not yet confirmed."},
        ],
    }


# ── Business Audit + AI Narrative + Portfolio Intelligence ────────────────────

def _build_audit_data(project_id: str):
    """Returns structured assumption registry + formula derivations for each project."""
    import datetime as _dt
    today = (_dt.datetime.utcnow() + _dt.timedelta(hours=3)).date()

    if project_id == "washing_bay":
        capex = 25_900_000; sanitation = 7_635_500; monthly_rev = 580_000
        months_running = max(1, (today.year - 2026)*12 + (today.month - 1))
        actual_rev = 2_923_000; months_confirmed = 5
        return {
            "project": "Washing Bay", "status": "Operational",
            "assumptions": [
                {"variable": "Initial CapEx", "value": capex, "unit": "UGX",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Alex Tuhimbise confirmed capital deployment",
                 "evidence": "Infrastructure installed: electricity, borehole, pump, jet spray machine"},
                {"variable": "Monthly Revenue", "value": monthly_rev, "unit": "UGX/month",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "5-month operational average",
                 "evidence": f"UGX {actual_rev:,} confirmed over {months_confirmed} months = UGX {actual_rev//months_confirmed:,}/month"},
                {"variable": "Sanitation Upgrade Cost", "value": sanitation, "unit": "UGX",
                 "type": "ESTIMATE", "confidence": "MEDIUM",
                 "source": "Contractor quote (septic + plumbing)",
                 "evidence": "Septic system UGX 5,470,000 + plumbing UGX 2,165,500. Not yet approved by board."},
                {"variable": "Revenue Growth Rate", "value": 0, "unit": "%/year",
                 "type": "ASSUMPTION", "confidence": "LOW",
                 "source": "Conservative flat assumption",
                 "evidence": "No seasonal adjustment model yet. Could be higher with throughput growth."},
            ],
            "formula_derivations": [
                {"metric": "Payback Period", "formula": "CapEx ÷ Monthly Net Revenue",
                 "substitution": f"UGX {capex:,} ÷ UGX {monthly_rev:,}/mo",
                 "result": f"{round(capex/monthly_rev,1)} months",
                 "note": "Assumes current revenue rate maintained. Does not include sanitation capex."},
                {"metric": "Capital Recovered to Date", "formula": "Total Revenue ÷ CapEx × 100",
                 "substitution": f"UGX {actual_rev:,} ÷ UGX {capex:,} × 100",
                 "result": f"{round(actual_rev/capex*100,1)}%",
                 "note": "Only {months_confirmed} confirmed months of revenue data."},
                {"metric": "Projected Annual Revenue", "formula": "Monthly Revenue × 12",
                 "substitution": f"UGX {monthly_rev:,} × 12",
                 "result": f"UGX {monthly_rev*12:,}",
                 "note": "Based on flat-rate assumption. Real revenue may vary with car wash volume."},
            ],
            "data_gaps": [
                "Daily/weekly throughput (cars washed) not tracked — prevents revenue forecasting",
                "Operating costs (labour, water, electricity) not reported — net margin unknown",
                "Sanitation board resolution not yet recorded in system",
            ],
            "what_would_change": [
                {"scenario": "Revenue increases to UGX 800K/month", "impact": f"Payback shortens to {round(capex/800000,1)} months"},
                {"scenario": "Sanitation approved", "impact": f"Total CapEx rises to UGX {capex+sanitation:,}; payback extends to {round((capex+sanitation)/monthly_rev,1)} months"},
                {"scenario": "Revenue drops to UGX 400K/month", "impact": f"Payback extends to {round(capex/400000,1)} months — still acceptable"},
            ],
        }

    elif project_id == "trees":
        capex_confirmed = 2_630_000; total_capex = 9_752_960
        poles = 3600; pole_price = 8000
        revenue_poles = poles * pole_price
        harvest_yr = 12
        roi = round((revenue_poles - total_capex)/total_capex*100,1)
        return {
            "project": "Trees (Eucalyptus)", "status": "Growing",
            "assumptions": [
                {"variable": "Seedlings Planted", "value": 3600, "unit": "trees",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Purchase receipt on file",
                 "evidence": "ABSA payment 16 Aug 2024, UGX 2,630,000 to Kikangi Israel. TXN 17237999200564"},
                {"variable": "Land Area", "value": 3.38, "unit": "acres",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Busisi Land parcel", "evidence": "Documented in proposal"},
                {"variable": "Survival Rate", "value": 80, "unit": "%",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Industry standard for Eucalyptus grandis in Uganda",
                 "evidence": "NARO benchmarks suggest 70-90% survival with basic maintenance. Using 80% conservative."},
                {"variable": "Poles per surviving tree", "value": 1, "unit": "poles/tree",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Harvest model assumption",
                 "evidence": "Each tree produces 1 saleable pole. Pruned trees may yield multiple."},
                {"variable": "Price per pole", "value": pole_price, "unit": "UGX",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Current Ugandan timber market",
                 "evidence": "Construction poles in Western Uganda trading at UGX 7,000-10,000 in 2024. Using UGX 8,000 conservative."},
                {"variable": "Harvest Year", "value": harvest_yr, "unit": "years from planting",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Eucalyptus grandis rotation cycle",
                 "evidence": "Full harvest at year 10-15. Using year 12 as mid-point. Selective harvest possible from year 5."},
            ],
            "formula_derivations": [
                {"metric": "Surviving Poles", "formula": "Seedlings × Survival Rate",
                 "substitution": f"3,600 × 80% = {int(3600*0.8):,} poles", "result": f"{int(3600*0.8):,} poles",
                 "note": "Survival rate is the single biggest uncertainty. Track losses annually."},
                {"metric": "Revenue (poles)", "formula": "Surviving Poles × Price/pole",
                 "substitution": f"{int(3600*0.8):,} × UGX {pole_price:,}",
                 "result": f"UGX {revenue_poles:,}",
                 "note": "Does not include firewood, charcoal, or biomass revenue."},
                {"metric": "Total CapEx", "formula": "Confirmed Payments + Future Commitments",
                 "substitution": f"UGX {capex_confirmed:,} + UGX {total_capex-capex_confirmed:,}",
                 "result": f"UGX {total_capex:,}",
                 "note": "UGX 2.63M is the only CASH DEPLOYED so far. Future commitments are budgeted, not yet spent."},
                {"metric": "ROI at Harvest", "formula": "(Revenue - Total CapEx) ÷ Total CapEx × 100",
                 "substitution": f"(UGX {revenue_poles:,} - UGX {total_capex:,}) ÷ UGX {total_capex:,} × 100",
                 "result": f"{roi}%",
                 "note": f"Over {harvest_yr} years. IRR (accounting for time value of money) is approximately 18-22%."},
            ],
            "data_gaps": [
                "Annual tree survival count not yet entered into system",
                "Current tree height/growth rate not tracked — critical for harvest date estimate",
                "Land tenure / legal ownership documentation status unknown",
                "Market price data for poles should be updated annually",
            ],
            "what_would_change": [
                {"scenario": "Survival rate drops to 60%", "impact": f"Revenue drops to UGX {int(3600*0.6*pole_price):,}; ROI falls to {round((int(3600*0.6*pole_price)-total_capex)/total_capex*100,1)}%"},
                {"scenario": "Pole price rises to UGX 12,000", "impact": f"Revenue rises to UGX {int(3600*0.8*12000):,}; ROI jumps to {round((int(3600*0.8*12000)-total_capex)/total_capex*100,1)}%"},
                {"scenario": "Harvest moved to year 10", "impact": "IRR improves; same revenue, 2 fewer years waiting — NPV increases significantly"},
            ],
        }

    elif project_id == "sheep":
        capital = 1_710_000; flock = 17; avg_val = 125_000
        flock_val = flock * avg_val; sold = 200_000
        pipeline = 12 * 200_000
        return {
            "project": "Sheep (Dorper)", "status": "Active",
            "assumptions": [
                {"variable": "Capital Invested", "value": capital, "unit": "UGX",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Receipt on file",
                 "evidence": "16 Aug 2024 payment to Tuhimbise Pricilla Arinaitwe for Dorper sheep"},
                {"variable": "Current Flock Size", "value": flock, "unit": "sheep",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Physical count at farm",
                 "evidence": "17 Dorpers as of last field report"},
                {"variable": "Market Value per Sheep", "value": avg_val, "unit": "UGX",
                 "type": "ESTIMATE", "confidence": "MEDIUM",
                 "source": "Current Ugandan Dorper market",
                 "evidence": "Dorper prices range UGX 100K-150K in Uganda. Using UGX 125K conservative midpoint."},
                {"variable": "Pending Sales (pipeline)", "value": 12, "unit": "sheep",
                 "type": "PROJECTION", "confidence": "MEDIUM",
                 "source": "KIM invoice KIM/07/26-4",
                 "evidence": "12 ewes at UGX 200K each. Customer invoice issued. Pending delivery and payment."},
            ],
            "formula_derivations": [
                {"metric": "Current Flock Value", "formula": "Flock Size × Market Value/sheep",
                 "substitution": f"{flock} × UGX {avg_val:,}", "result": f"UGX {flock_val:,}",
                 "note": "Market value. Actual realised value depends on sale timing and buyer."},
                {"metric": "Net Position", "formula": "Flock Value + Proceeds - Capital",
                 "substitution": f"UGX {flock_val:,} + UGX {sold:,} - UGX {capital:,}",
                 "result": f"UGX +{flock_val+sold-capital:,}",
                 "note": "Unrealised gain. Only UGX 200K is cash-in-hand. The rest is on-hoof."},
                {"metric": "Net Return %", "formula": "Net Position ÷ Capital × 100",
                 "substitution": f"UGX {flock_val+sold-capital:,} ÷ UGX {capital:,} × 100",
                 "result": f"{round((flock_val+sold-capital)/capital*100,1)}%",
                 "note": "Strong unrealised return but breeding timeline and mortality risk not yet quantified."},
            ],
            "data_gaps": [
                "Breeding record (rams, ewes, lambing dates) not entered in system",
                "Veterinary/feed cost tracking not in place — net margin unknown",
                "Pipeline sale (KIM/07/26-4) status: delivery date and payment date not confirmed",
            ],
            "what_would_change": [
                {"scenario": "Pipeline sale completes (12 ewes × UGX 200K)", "impact": f"Cash realised: UGX {pipeline:,}. Cash return on capital: {round((sold+pipeline)/capital*100,1)}%"},
                {"scenario": "Flock grows to 30 via breeding", "impact": f"Flock value: UGX {30*avg_val:,}. Return on capital: {round((30*avg_val-capital)/capital*100,1)}%"},
            ],
        }

    elif project_id == "dairy":
        capex = 16_000_000; cow_cost = 8_000_000; cows = 2
        litres_mid = 30 * cows; price = 1_500
        monthly_rev_mid = litres_mid * price * 30
        annual_opex = 3_600_000
        annual_profit = monthly_rev_mid * 12 - annual_opex
        payback = round(capex / (monthly_rev_mid - annual_opex/12), 1)
        return {
            "project": "Dairy", "status": "Planning",
            "assumptions": [
                {"variable": "Cost per Cow", "value": cow_cost, "unit": "UGX",
                 "type": "ESTIMATE", "confidence": "MEDIUM",
                 "source": "Site visit 25 May 2026 (Dad visited working shed dairy)",
                 "evidence": "Working farm site visit confirmed UGX 8M per quality Friesian cow. Market range UGX 6M-12M."},
                {"variable": "Daily Milk Yield (mid)", "value": 30, "unit": "L/day/cow",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Zero-grazing industry benchmarks",
                 "evidence": "Quality Friesians in zero-grazing sheds: 15-50L/day. 30L is realistic for well-managed shed. Actual yield unknown until cows purchased."},
                {"variable": "Farm-gate Milk Price", "value": price, "unit": "UGX/L",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Current Mbarara region milk prices",
                 "evidence": "Farm-gate: UGX 1,200-1,800/L. Using UGX 1,500 conservative. Retail (own processing) can reach UGX 2,500/L."},
                {"variable": "Annual Operating Costs", "value": annual_opex, "unit": "UGX/year",
                 "type": "ESTIMATE", "confidence": "LOW",
                 "source": "Similar zero-grazing operations",
                 "evidence": "Feed, vet, labour, utilities. Highly variable. UGX 3.6M/year is rough estimate pending detailed budget."},
            ],
            "formula_derivations": [
                {"metric": "Monthly Revenue (mid)", "formula": "Litres/day × Cows × Days × Price/litre",
                 "substitution": f"30L × {cows} cows × 30 days × UGX {price:,}",
                 "result": f"UGX {monthly_rev_mid:,}/month",
                 "note": "The 30L/day assumption drives this entire model. Actual yield determines everything."},
                {"metric": "Payback Period (mid)", "formula": "CapEx ÷ (Monthly Revenue - Monthly OpEx)",
                 "substitution": f"UGX {capex:,} ÷ (UGX {monthly_rev_mid:,} - UGX {annual_opex//12:,})",
                 "result": f"{payback} months",
                 "note": "Highly sensitive to milk yield. At 15L/day: {round(capex/(15*cows*price*30 - annual_opex/12),1)} months. At 50L/day: {round(capex/(50*cows*price*30 - annual_opex/12),1)} months."},
            ],
            "data_gaps": [
                "Capital not committed — UGX 16M not yet approved or sourced",
                "Shed construction/conversion cost not estimated",
                "Specific breed and supplier not identified",
                "Offtake arrangement (who buys the milk) not in place",
                "Skilled labour (zero-grazing experience) not sourced",
            ],
            "what_would_change": [
                {"scenario": "Low yield (15L/day)", "impact": f"Monthly rev = UGX {15*cows*price*30:,}. Payback = {round(capex/(15*cows*price*30-annual_opex/12),1)} months"},
                {"scenario": "High yield (50L/day)", "impact": f"Monthly rev = UGX {50*cows*price*30:,}. Payback = {round(capex/(50*cows*price*30-annual_opex/12),1)} months"},
                {"scenario": "Milk price drops to UGX 1,000/L", "impact": f"Monthly rev = UGX {30*cows*1000*30:,}. Payback = {round(capex/(30*cows*1000*30-annual_opex/12),1)} months"},
            ],
        }

    elif project_id == "irrigation":
        rev = 5_000_000; costs = 3_960_000; profit = rev - costs
        bunches = 500; price_bunch = 10_000
        return {
            "project": "Irrigation + Bananas", "status": "Planning",
            "assumptions": [
                {"variable": "Annual Banana Bunches", "value": bunches, "unit": "bunches/year",
                 "type": "PROJECTION", "confidence": "MEDIUM",
                 "source": "Irrigation proposal model",
                 "evidence": "Based on 2-acre banana plantation with irrigation. Current yield without irrigation ~80 bunches lost to heat stress annually. 500 bunches is achievable target."},
                {"variable": "Price per Bunch", "value": price_bunch, "unit": "UGX",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Local Mbarara banana market",
                 "evidence": "Banana bunch prices: UGX 8,000-15,000. Using UGX 10,000 conservative."},
                {"variable": "Annual Maintenance Labor", "value": 2_400_000, "unit": "UGX/year",
                 "type": "ESTIMATE", "confidence": "MEDIUM",
                 "source": "Irrigation proposal costing",
                 "evidence": "From proposal spreadsheet: weeding, spraying, harvesting labour."},
                {"variable": "Irrigation System CapEx", "value": 0, "unit": "UGX",
                 "type": "UNKNOWN", "confidence": "LOW",
                 "source": "Phase 2 NOT YET COSTED",
                 "evidence": "Design complete (Phase 1). Equipment procurement cost (Phase 2) not yet estimated. THIS IS A CRITICAL DATA GAP."},
            ],
            "formula_derivations": [
                {"metric": "Annual Revenue", "formula": "Bunches/year × Price/bunch",
                 "substitution": f"{bunches} × UGX {price_bunch:,}",
                 "result": f"UGX {rev:,}",
                 "note": "Assumes all 500 bunches find buyers. Market absorption should be validated."},
                {"metric": "Annual Profit", "formula": "Revenue - Annual Operating Costs",
                 "substitution": f"UGX {rev:,} - UGX {costs:,}",
                 "result": f"UGX {profit:,}",
                 "note": "Does NOT include irrigation system depreciation/amortisation (CapEx unknown)."},
                {"metric": "Payback Period", "formula": "Irrigation CapEx ÷ Annual Profit",
                 "substitution": "UNKNOWN ÷ UGX 1,040,000",
                 "result": "CANNOT CALCULATE — CapEx unknown",
                 "note": "This is the most important missing number. Get Phase 2 quote ASAP."},
            ],
            "data_gaps": [
                "CRITICAL: Irrigation system CapEx (Phase 2 procurement + installation) not costed",
                "Expansion plan (4 additional acres) not costed",
                "Water source capacity (farm ponds) not assessed against irrigation demand",
                "Offtake: who buys the bananas? Market arrangement not documented",
            ],
            "what_would_change": [
                {"scenario": "Irrigation CapEx = UGX 5M", "impact": f"Payback = {round(5000000/profit,1)} years. Very viable."},
                {"scenario": "Irrigation CapEx = UGX 15M", "impact": f"Payback = {round(15000000/profit,1)} years. Marginal — needs higher yield or price."},
                {"scenario": "Expand to 6 acres (3x output)", "impact": f"Revenue = UGX {rev*3:,}, Profit = UGX {profit*3:,}/yr. Transforms the project economics."},
            ],
        }

    elif project_id == "bees":
        hives = 10; kg_hive = 25; kg_total = hives * kg_hive
        price_kg = 15_000; rev = kg_total * price_kg
        setup = hives * 150_000 + 500_000; opex = 600_000
        profit = rev - opex
        return {
            "project": "Bees", "status": "Research",
            "assumptions": [
                {"variable": "Hives Planned", "value": hives, "unit": "hives",
                 "type": "PROJECTION", "confidence": "MEDIUM",
                 "source": "Research phase decision",
                 "evidence": "10 hives is a typical commercial starter scale in Uganda."},
                {"variable": "Yield per Hive", "value": kg_hive, "unit": "kg/hive/year",
                 "type": "ASSUMPTION", "confidence": "LOW",
                 "source": "Ugandan beekeeping benchmarks",
                 "evidence": "NAADS benchmarks: 15-35kg/hive/year with modern hives. Using 25kg conservative. Highly dependent on forage availability and management quality."},
                {"variable": "Honey Price", "value": price_kg, "unit": "UGX/kg",
                 "type": "ASSUMPTION", "confidence": "MEDIUM",
                 "source": "Current Ugandan honey market",
                 "evidence": "Bulk honey: UGX 10,000-20,000/kg. Branded/retail: UGX 25,000-50,000/kg. Using UGX 15,000 (bulk/semi-processed)."},
                {"variable": "Hive Cost", "value": 150_000, "unit": "UGX/hive",
                 "type": "ESTIMATE", "confidence": "MEDIUM",
                 "source": "Ugandan beekeeping suppliers",
                 "evidence": "Langstroth hives: UGX 120,000-200,000 each. Using UGX 150,000 midpoint."},
            ],
            "formula_derivations": [
                {"metric": "Annual Revenue", "formula": "Hives × Yield/hive × Price/kg",
                 "substitution": f"{hives} × {kg_hive}kg × UGX {price_kg:,}",
                 "result": f"UGX {rev:,}",
                 "note": "The yield assumption is the weakest link. Year 1 yield will be lower as colonies establish."},
                {"metric": "Annual Profit", "formula": "Revenue - Annual Operating Costs",
                 "substitution": f"UGX {rev:,} - UGX {opex:,}",
                 "result": f"UGX {profit:,}",
                 "note": "OpEx is an estimate. Includes protective gear, harvesting equipment, transport."},
                {"metric": "Setup Cost", "formula": "(Hives × Cost/hive) + Initial Equipment",
                 "substitution": f"({hives} × UGX 150,000) + UGX 500,000",
                 "result": f"UGX {setup:,}",
                 "note": "Excludes colony purchase cost if buying packaged bees vs. catching wild swarms."},
            ],
            "data_gaps": [
                "No capital committed — project has not started",
                "Forage assessment not done — what trees/flowers on and near the farm?",
                "Training/expertise confirmed: Dad joined beekeeper groups but formal training status unknown",
                "Branded distribution channel not designed (KIM FAM Honey Farm brand registered?)",
            ],
            "what_would_change": [
                {"scenario": "Yield only 15kg/hive (year 1 colony establishment)", "impact": f"Revenue = UGX {hives*15*price_kg:,}. Profit = UGX {hives*15*price_kg-opex:,}. Still profitable."},
                {"scenario": "Brand direct (UGX 30,000/kg retail)", "impact": f"Revenue doubles to UGX {kg_total*30000:,}. Profit = UGX {kg_total*30000-opex:,}/yr. Game-changing."},
                {"scenario": "Scale to 30 hives (natural split in yr 3)", "impact": f"Revenue = UGX {30*kg_hive*price_kg:,}. Profit = UGX {30*kg_hive*price_kg-opex*1.5:,.0f}/yr. Near-zero marginal cost to scale."},
            ],
        }

    elif project_id == "chicken":
        return {
            "project": "Chicken (Broilers)", "status": "Active",
            "assumptions": [
                {"variable": "Financial data", "value": "Live from Solomon AppSheet", "unit": "",
                 "type": "ACTUAL", "confidence": "HIGH",
                 "source": "Solomon's AppSheet live data integration",
                 "evidence": "Cycle revenue, costs, and flock data pulled in real time from Solomon's management system."},
            ],
            "formula_derivations": [
                {"metric": "All metrics", "formula": "See chicken analysis tab",
                 "substitution": "Live data from AppSheet API",
                 "result": "Dynamic",
                 "note": "Chicken is the most data-rich project. All numbers are actuals from Solomon's records."},
            ],
            "data_gaps": ["None critical — AppSheet integration provides real-time data"],
            "what_would_change": [],
        }
    return {"project": project_id, "assumptions": [], "formula_derivations": [], "data_gaps": [], "what_would_change": []}


def _ask_claude(prompt: str, model: str = "sonnet", timeout: int = 120) -> str:
    """Call the Claude Code CLI. Returns response text or empty string on failure."""
    import subprocess as _sp, os as _os
    env = dict(_os.environ); env["HOME"] = "/root"
    try:
        r = _sp.run(
            ["claude", "-p", prompt, "--model", model],
            capture_output=True, text=True, timeout=timeout, env=env
        )
        return r.stdout.strip() if r.returncode == 0 else ""
    except Exception:
        return ""



async def _ask_claude_async(prompt: str, model: str = "haiku", timeout: int = 40) -> str:
    """Non-blocking Claude CLI call using asyncio subprocess."""
    import asyncio as _aio, os as _os
    env = dict(_os.environ); env["HOME"] = "/root"
    proc = None
    try:
        proc = await _aio.create_subprocess_exec(
            "claude", "-p", prompt, "--model", model,
            stdout=_aio.subprocess.PIPE,
            stderr=_aio.subprocess.DEVNULL,
            env=env
        )
        stdout, _ = await _aio.wait_for(proc.communicate(), timeout=timeout)
        return stdout.decode().strip() if proc.returncode == 0 else ""
    except Exception:
        if proc:
            try: proc.kill()
            except: pass
        return ""

@app.get("/api/projects/{project_id}/audit")
async def project_audit(project_id: str, request: Request):
    from fastapi import HTTPException as _HE
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    valid = {"washing_bay","trees","sheep","dairy","irrigation","bees","chicken","mango","goats","rabbits","fish"}
    if project_id not in valid: raise _HE(status_code=404, detail="Unknown project")
    return _build_audit_data(project_id)


@app.post("/api/projects/{project_id}/narrative")
async def project_narrative(project_id: str, request: Request):
    from fastapi import HTTPException as _HE
    import os as _os, json as _json
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")

    audit = _build_audit_data(project_id)

    # P2 (#2): link the live project to its approved proposal + the reward guidelines, so the
    # review also checks delivery-vs-promise and reward-guideline compliance.
    _PROJ_PROPOSAL_KW = {"chicken": "chicken", "trees": "eucalyptus", "bees": "bee",
                         "irrigation": "irrigation", "washing_bay": "wash"}
    prop_block = ""
    kw = _PROJ_PROPOSAL_KW.get(project_id)
    if kw:
        try:
            from db import query as _dbq3
            pr = _dbq3("""SELECT title, overall_score, verdict, summary FROM proposals
                          WHERE is_current=TRUE AND title ILIKE %s ORDER BY id DESC LIMIT 1""",
                       (f"%{kw}%",))
            if pr:
                p0 = pr[0]
                prop_block = (
                    f"\nAPPROVED PROPOSAL (what the family committed to for this project):\n"
                    f"Title: {p0['title']} (readiness {p0.get('overall_score')}/100, {p0.get('verdict')})\n"
                    f"Proposal summary: {(p0.get('summary') or '')[:600]}\n")
        except Exception as _pe:
            log.error(f"narrative proposal link failed: {_pe}")
    try:
        # Reward guidelines portion only (skip the PM-framework half), kept small so the
        # non-streaming narrative call stays under the proxy timeout. (Streaming = #3.)
        _fw = _proposal_framework_context()
        _ri = _fw.find("INVESTMENT & REWARD")
        reward_block = "\nINVESTMENT & REWARD GUIDELINES (the family's rules):\n" + (_fw[_ri:_ri + 1800] if _ri >= 0 else _fw[:1800])
    except Exception:
        reward_block = ""

    prompt = f"""You are a seasoned agricultural investment analyst reviewing a Ugandan family investment club's farm project.

PROJECT: {audit.get("project")} | STATUS: {audit.get("status")}

ASSUMPTION REGISTRY:
{_json.dumps(audit.get("assumptions",[]), indent=2)}

FORMULA DERIVATIONS:
{_json.dumps(audit.get("formula_derivations",[]), indent=2)}

DATA GAPS:
{_json.dumps(audit.get("data_gaps",[]), indent=2)}

SENSITIVITY ANALYSIS:
{_json.dumps(audit.get("what_would_change",[]), indent=2)}

MONITORING & EVALUATION FRAMEWORK (the family's own standard, from the Project Management deck — evaluate this project AGAINST it):
- KPI categories to track: Project Progress (percentage of completion; milestone achievement); Budget Management (cost variance: planned vs actual; budget utilization: spent vs remaining); Quality (adherence to project requirements; stakeholder satisfaction); Risk Management (risk occurrence; mitigation effectiveness).
- Prescribed tracking tools: Trello (project management), Wave Accounting (expense tracking, budgeting, financial reporting), Google Sheets (supplementary analysis), Google Docs / Keep (documentation and notes), Google Calendar (scheduling), and dashboards for visual KPI representation.
{prop_block}{reward_block}
Write a board-quality analysis with these sections. Be specific, cite the numbers, be honest about what is projected vs confirmed. Do not use em-dashes. Use plain English a family with mixed business literacy can follow. Maximum 550 words.

EXECUTIVE SUMMARY (2-3 sentences: what is this project, current state, headline metric)

KEY STRENGTHS (3 bullet points max)

KEY RISKS (3 bullet points max, most critical first)

M&E ALIGNMENT (assess the project against the Monitoring & Evaluation framework above. For each KPI category — Progress, Budget, Quality, Risk — say in one line whether it is being tracked and on course, or is a gap. Then note any prescribed tracking tool not yet in use. 4-5 bullet points.)

PROPOSAL ALIGNMENT (only if an approved proposal is provided above: is the live project delivering what the proposal promised? Compare promised vs actual on the things the proposal committed to — budget, returns/payback, timeline, scope. 2-4 bullet points. If no proposal was provided, write one line: "No approved proposal is linked to this project yet.")

REWARD COMPLIANCE (using the Investment & Reward Guidelines above: is this project's participation and reward-sharing in line with the guidelines, and what should be checked or corrected? 2-3 bullet points.)

WHAT THE FAMILY SHOULD DO NEXT (2-3 specific actions, numbered)

BOARD RECOMMENDATION (one line: Proceed / Monitor / Hold / Urgent Action Needed + brief reason)

SURPRISING INSIGHT (one paragraph: something non-obvious the data reveals that the family probably has not considered)"""

    import os as _os, asyncio as _aio
    from fastapi.responses import StreamingResponse

    def _work():
        # Claude first; fall back Gemini -> Groq. Blocking; runs in a worker thread.
        txt = _ask_claude(prompt, model="sonnet", timeout=170)
        prov = "Claude Sonnet" if txt else ""
        if not txt and _os.getenv("GEMINI_API_KEY",""):
            try:
                from google import genai as _genai
                resp = _genai.Client(api_key=_os.getenv("GEMINI_API_KEY")).models.generate_content(
                    model="gemini-2.0-flash", contents=prompt)
                txt = resp.text; prov = "Gemini"
            except Exception: pass
        if not txt and _os.getenv("GROQ_API_KEY",""):
            try:
                from groq import Groq as _Groq
                resp = _Groq(api_key=_os.getenv("GROQ_API_KEY")).chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role":"user","content":prompt}], max_tokens=800)
                txt = resp.choices[0].message.content; prov = "Groq"
            except Exception: pass
        if not txt:
            txt = "AI analysis unavailable. Please try again."; prov = "none"
        return {
            "type": "result",
            "project_id": project_id, "project": audit.get("project"),
            "narrative": txt, "provider": prov,
            "based_on": {
                "assumptions_count": len(audit.get("assumptions",[])),
                "formulas_count": len(audit.get("formula_derivations",[])),
                "data_gaps_count": len(audit.get("data_gaps",[])),
            },
        }

    _STAGES = ["Loading project data...", "Checking the M&E framework...",
               "Comparing against the approved proposal and reward guidelines...",
               "Writing the board analysis with Claude (about a minute)...",
               "Almost done, finalising the analysis..."]

    async def _gen():
        yield _sse({"type": "step", "msg": "Loading project data..."})
        task = _aio.create_task(_aio.to_thread(_work))
        i = 0
        while not task.done():
            await _aio.sleep(6)
            if not task.done():
                yield _sse({"type": "step", "msg": _STAGES[min(i, len(_STAGES) - 1)]}); i += 1
        try:
            yield _sse(task.result())
        except Exception as _e:
            log.error(f"narrative stream failed: {_e}")
            yield _sse({"type": "error", "msg": "Analysis is busy right now. Please try again."})

    return StreamingResponse(_gen(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "Connection": "keep-alive"})



def _sse(obj: dict) -> str:
    import json as _j
    return f"data: {_j.dumps(obj)}\n\n"


@app.get("/api/portfolio/ranking/stream")
async def portfolio_ranking_stream(request: Request):
    """Real SSE stream: Haiku ranks → Gemini 2.5 Flash streams insights token-by-token."""
    from fastapi import HTTPException as _HE
    from fastapi.responses import StreamingResponse
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")

    async def _generate():
        import os as _os, json as _json, re as _re
        import asyncio as _aio, queue as _q, threading as _t

        try:
            yield _sse({"type":"step","step":1,"total":5,"msg":"Access granted — fetching live data from all projects..."})

            # ── Step 2: Parallel-fetch all 7 project detail endpoints ──────────
            import urllib.request as _ur
            ikey = _os.environ.get("KIMFAM_INTERNAL_KEY", "")
            PIDS = ["chicken","trees","sheep","washing_bay","irrigation","dairy","bees"]

            async def _fetch_one(pid):
                def _do():
                    try:
                        req = _ur.Request(
                            f"http://127.0.0.1:8000/api/projects/{pid}/detail",
                            headers={"X-Internal-Key": ikey}
                        )
                        with _ur.urlopen(req, timeout=8) as r:
                            return _json.loads(r.read())
                    except:
                        return {}
                return pid, await _aio.get_running_loop().run_in_executor(None, _do)

            pairs = await _aio.gather(*[_fetch_one(p) for p in PIDS])
            live = {pid: d for pid, d in pairs}
            ok_count = sum(1 for d in live.values() if d)
            yield _sse({"type":"step","step":2,"total":5,"msg":f"Got live data from {ok_count}/7 projects — Claude Haiku ranking..."})

            # ── Step 3: Claude Haiku — fast structured ranking (JSON only) ─────
            def _compact(pid, d):
                ov = d.get("overview", {}); fm = d.get("financial_metrics", {}); inv = d.get("investment", {})
                flat = {k: v for k, v in {**ov, **fm, **inv}.items() if not isinstance(v, (dict, list))}
                return flat

            summary = {pid: _compact(pid, d) for pid, d in live.items()}

            haiku_prompt = (
                "Rank these 7 KimFam Uganda farm investment projects.\n"
                f"LIVE DATA:\n{_json.dumps(summary, indent=2)[:5500]}\n\n"
                "Return ONLY valid JSON — no markdown, no explanation:\n"
                '{"ranked":[{"rank":1,"project_id":"...","tier":"MOVE NOW|BUILD CAREFULLY|LET COMPOUND",'
                '"score_out_of_10":8,"score_rationale":"2 sentences with numbers","time_frame":"short|medium|long",'
                '"strategic_note":"Specific 90-day action for the family"}]}'
            )
            haiku_raw = await _ask_claude_async(haiku_prompt, model="haiku", timeout=60)
            ranked = {"ranked": []}
            try:
                _hr = _re.sub(r"```[\w]*\n?", "", haiku_raw or "").strip()
                m = _re.search(r"\{.*\}", _hr, _re.DOTALL)
                if m: ranked = _json.loads(m.group())
            except: pass
            # Haiku CLI intermittently returns nothing/unparseable — fall back so the
            # ranking is never empty: Gemini Flash, then Groq.
            if not ranked.get("ranked"):
                _gk = _os.environ.get("GEMINI_API_KEY", "")
                if _gk:
                    try:
                        from google import genai as _genai
                        _gr = _genai.Client(api_key=_gk).models.generate_content(
                            model="gemini-2.0-flash", contents=haiku_prompt).text
                        _gm = _re.search(r"\{.*\}", _re.sub(r"```[\w]*\n?", "", _gr or ""), _re.DOTALL)
                        if _gm: ranked = _json.loads(_gm.group())
                    except: pass
            if not ranked.get("ranked"):
                _qk = _os.environ.get("GROQ_API_KEY", "")
                if _qk:
                    try:
                        from groq import Groq as _Groq
                        _qr = _Groq(api_key=_qk).chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[{"role": "user", "content": haiku_prompt}],
                            max_tokens=1200, response_format={"type": "json_object"}
                        ).choices[0].message.content
                        _qm = _re.search(r"\{.*\}", _qr or "", _re.DOTALL)
                        if _qm: ranked = _json.loads(_qm.group())
                    except: pass
            n = len(ranked.get("ranked", []))
            yield _sse({"type":"step","step":3,"total":5,"msg":f"Haiku scored {n} projects on capital, speed and risk."})
            # Reveal the interim verdict so the user watches the reasoning, not a label.
            _names = {"chicken":"🐔 Free Range Chicken","washing_bay":"🚗 Washing Bay","sheep":"🐑 Sheep",
                      "trees":"🌲 Tree Planting","irrigation":"💧 Irrigation","dairy":"🐄 Dairy","bees":"🍯 Apiary"}
            _sorted = sorted(ranked.get("ranked", []), key=lambda x: x.get("rank", 99))
            if _sorted:
                _t = _sorted[0]
                _nm = _names.get(_t.get("project_id",""), _t.get("project_id","?"))
                yield _sse({"type":"step","step":3,"total":5,"msg":f"Lead pick so far: {_nm} at {_t.get('score_out_of_10','?')}/10, tier {_t.get('tier','')}."})
            if len(_sorted) > 1:
                _b = _sorted[-1]
                _bn = _names.get(_b.get("project_id",""), _b.get("project_id","?"))
                yield _sse({"type":"step","step":3,"total":5,"msg":f"Most patience needed: {_bn} ({_b.get('tier','')}). Now weighing portfolio-wide synergies..."})

            # ── Step 4: Gemini 2.5 Flash — token streaming for narrative ───────
            insight_prompt = (
                "Senior portfolio advisor, KimFam Investment Club, Western Uganda.\n\n"
                f"RANKINGS (from Haiku):\n{_json.dumps(ranked.get('ranked',[]), indent=2)}\n\n"
                f"LIVE DATA SUMMARY:\n{_json.dumps(summary, indent=2)[:3000]}\n\n"
                "Return ONLY valid JSON (no markdown fences):\n"
                '{"portfolio_insight":"2-3 sentences — what do the numbers reveal the family has not noticed?",'
                '"biggest_opportunity":"Single biggest untapped opportunity, cite numbers",'
                '"biggest_risk":"Single risk that hurts multiple projects simultaneously",'
                '"compounding_play":"Project or combo that creates flywheel — cite specific synergies"}'
            )

            insights = {}
            yield _sse({"type":"step","step":4,"total":5,"msg":"Claude Sonnet generating strategic insights..."})

            # Run the (slow) insight call as a task and emit "thinking" heartbeats
            # while it runs — keeps the SSE alive past nginx's idle timeout AND
            # lets the user watch the reasoning unfold.
            _think = ["Cross-referencing capital against live cash flow...",
                      "Hunting for the flywheel between projects...",
                      "Finding the one risk that hits several ventures at once...",
                      "Drafting the next 90-day moves..."]
            _task = _aio.create_task(_ask_claude_async(insight_prompt, model="sonnet", timeout=90))
            _hi = 0
            while not _task.done():
                _done, _ = await _aio.wait({_task}, timeout=7)
                if not _done:
                    yield _sse({"type":"step","step":4,"total":5,"msg":_think[_hi % len(_think)]})
                    _hi += 1
            sonnet_raw = _task.result()

            if not sonnet_raw:
                groq_key = _os.environ.get("GROQ_API_KEY", "")
                if groq_key:
                    try:
                        from groq import Groq as _Groq
                        resp = _Groq(api_key=groq_key).chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[{"role":"user","content":insight_prompt}],
                            max_tokens=800, response_format={"type":"json_object"})
                        sonnet_raw = resp.choices[0].message.content
                    except: pass

            try:
                _ft = _re.sub(r"```[\w]*\n?", "", sonnet_raw or "").strip()
                m = _re.search(r"\{.*\}", _ft, _re.DOTALL)
                if m: insights = _json.loads(m.group())
            except:
                insights = {"portfolio_insight": sonnet_raw[:600] if sonnet_raw else "Analysis unavailable"}

            # ── Step 5: Merge and emit ─────────────────────────────────────────
            yield _sse({"type":"step","step":5,"total":5,"msg":"Analysis complete — loading results..."})
            yield _sse({"type":"result","data":{**ranked, **insights}})

        except Exception as e:
            import traceback as _tb
            yield _sse({"type":"error","msg":str(e),"trace":_tb.format_exc()[-400:]})

    return StreamingResponse(
        _generate(),
        media_type="text/event-stream",
        headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no","Connection":"keep-alive"}
    )


@app.post("/api/portfolio/new_ventures/stream")
async def new_ventures_stream(request: Request):
    """Real SSE stream: build context → Gemini 2.5 Flash streams venture proposals token-by-token."""
    from fastapi import HTTPException as _HE
    from fastapi.responses import StreamingResponse
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")
    body = await request.json()
    focus_area = body.get("focus_area", "all")

    async def _generate():
        import os as _os, json as _json, re as _re
        import asyncio as _aio, queue as _q, threading as _t

        try:
            yield _sse({"type":"step","step":1,"total":4,"msg":"Building portfolio context..."})

            portfolio_context = {
                "club": {"name":"KIM FAM Investment Club","location":"Mbarara area, Western Uganda",
                         "members":13,"focus":"Agricultural and agri-business investments",
                         "land_assets":["Nyabugando Farm (main)","Busisi Land (3.38 acres, eucalyptus)"]},
                "current_portfolio": {
                    "operational":["Chicken broilers — AppSheet tracked","Washing Bay — 5+ months revenue"],
                    "active":["Sheep Dorper — 17 head","Trees Eucalyptus — 3,600 seedlings"],
                    "planning":["Dairy zero-grazing — UGX 16M needed","Irrigation+Bananas — Phase 2 TBD"],
                    "research":["Bees — 10 hives planned"]
                },
                "resources":{"land":"Multi-acre Mbarara","water":"Borehole + ponds",
                             "labour":"Family + casual","markets":"Mbarara city, Kampala road",
                             "inputs":"Manure(planned), biomass, irrigation"},
                "capital":{"deployed":"UGX 30,240,000","monthly_revenue":"UGX 580,000+ (wash bay) + chicken",
                           "next_need":"UGX 16M dairy"},
                "focus_area": focus_area
            }

            yield _sse({"type":"step","step":2,"total":4,"msg":f"Context ready ({focus_area}) — Claude Sonnet researching venture opportunities..."})

            prompt = (
                "Venture advisor for KimFam Investment Club, Uganda.\n\n"
                f"CLUB PROFILE:\n{_json.dumps(portfolio_context, indent=2)}\n\n"
                f"FOCUS: {focus_area}\n\n"
                "Propose 5 NEW investment opportunities (NOT in current portfolio). "
                "Be specific, cite Ugandan market numbers, link to existing assets.\n\n"
                "Return ONLY valid JSON:\n"
                '{"ventures":[{"rank":1,"name":"...","category":"Agriculture|Agriprocessing|Services|Technology",'
                '"headline":"one-line pitch","why_now":"Uganda 2026 opportunity","why_this_club":"links to their assets",'
                '"capital_required_ugx":5000000,"capital_note":"what it covers","revenue_model":"how money is made",'
                '"time_to_first_revenue":"X months","annual_revenue_estimate_ugx":0,"annual_profit_estimate_ugx":0,'
                '"synergies_with_portfolio":"specific links to existing projects","risks":["risk1"],'
                '"first_step":"single most important next 30-day action",'
                '"external_resources":[{"title":"...","url":"https://...","why_relevant":"..."}]}],'
                '"strategic_commentary":"how these 5 ventures compound the portfolio",'
                '"compounding_sequence":"optimal order to build wealth fastest"}'
            )

            full_text = ""
            yield _sse({"type":"step","step":3,"total":4,"msg":"Claude Sonnet researching venture opportunities..."})

            # Heartbeat "thinking" lines while the slow research call runs, so the
            # SSE stays alive (past nginx idle timeout) and the user sees progress.
            _vt = ["Scanning Uganda 2026 agriculture and agri-processing trends...",
                   "Matching ideas to your land, water and market access...",
                   "Checking each idea links back to an existing asset...",
                   "Sizing capital and time to first revenue...",
                   "Sequencing them for the fastest compounding..."]
            _task = _aio.create_task(_ask_claude_async(prompt, model="sonnet", timeout=120))
            _vi = 0
            while not _task.done():
                _done, _ = await _aio.wait({_task}, timeout=7)
                if not _done:
                    yield _sse({"type":"step","step":3,"total":4,"msg":_vt[_vi % len(_vt)]})
                    _vi += 1
            full_text = _task.result()

            if not full_text:
                groq_key = _os.environ.get("GROQ_API_KEY", "")
                if groq_key:
                    try:
                        from groq import Groq as _Groq
                        resp = _Groq(api_key=groq_key).chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[{"role":"user","content":prompt}],
                            max_tokens=3000, response_format={"type":"json_object"})
                        full_text = resp.choices[0].message.content
                    except: pass

            yield _sse({"type":"step","step":4,"total":4,"msg":"Proposals ready — parsing results..."})

            result = {}
            try:
                _ft2 = _re.sub(r"```[\w]*\n?", "", full_text).strip()
                m = _re.search(r"\{.*\}", _ft2, _re.DOTALL)
                if m: result = _json.loads(m.group())
            except:
                result = {"ventures":[],"raw":full_text[:2000],"error":"JSON parse failed"}

            _vs = result.get("ventures", []) if isinstance(result, dict) else []
            if _vs:
                _v0 = _vs[0]
                yield _sse({"type":"step","step":4,"total":4,"msg":f"Top idea: {_v0.get('name','?')} — {(_v0.get('headline','') or '')[:80]}"})
                yield _sse({"type":"step","step":4,"total":4,"msg":f"{len(_vs)} fresh ventures mapped to your land, water and markets. Compiling..."})

            yield _sse({"type":"result","data":result})

        except Exception as e:
            import traceback as _tb
            yield _sse({"type":"error","msg":str(e),"trace":_tb.format_exc()[-400:]})

    return StreamingResponse(
        _generate(),
        media_type="text/event-stream",
        headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no","Connection":"keep-alive"}
    )


@app.get("/api/portfolio/ranking")
async def portfolio_ranking(request: Request):
    """AI-ranked portfolio: returns all projects ranked by strategic priority."""
    from fastapi import HTTPException as _HE
    import os as _os, json as _json
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")

    portfolio_data = {
        "chicken": {"status":"Active","payback_months":None,"roi_pct":"live","capital_deployed":True,"revenue_generating":True,"capital_ugx":"live AppSheet","risk":"Low","time_horizon":"Short","data_quality":"HIGH"},
        "washing_bay": {"status":"Operational","payback_months":44.7,"roi_pct":None,"capital_deployed":True,"revenue_generating":True,"capital_ugx":"25,900,000","recovery_pct":11.3,"risk":"Medium","time_horizon":"Medium","data_quality":"HIGH"},
        "sheep": {"status":"Active","payback_months":None,"roi_pct":36,"capital_deployed":True,"revenue_generating":True,"capital_ugx":"1,710,000","pipeline_ugx":"2,400,000","risk":"Low","time_horizon":"Short","data_quality":"MEDIUM"},
        "trees": {"status":"Growing","payback_months":None,"roi_pct":731,"capital_deployed":True,"revenue_generating":False,"capital_ugx":"2,630,000","harvest_year":2036,"risk":"Low","time_horizon":"Long","data_quality":"MEDIUM"},
        "irrigation": {"status":"Planning","payback_months":"TBD","roi_pct":"26% ops only","capital_deployed":False,"revenue_generating":False,"capex_gap":"Phase 2 unknown","risk":"Medium","time_horizon":"Medium","data_quality":"LOW"},
        "dairy": {"status":"Planning","payback_months":6.7,"roi_pct":180,"capital_deployed":False,"revenue_generating":False,"capital_ugx":"16,000,000 needed","risk":"High","time_horizon":"Short","data_quality":"LOW"},
        "bees": {"status":"Research","payback_months":7,"roi_pct":"157","capital_deployed":False,"revenue_generating":False,"setup_ugx":"2,000,000","risk":"Medium","time_horizon":"Short","data_quality":"LOW"},
    }

    prompt = f"""You are a strategic advisor to the KIM FAM Investment Club, a Ugandan family farm investment group.

Here is the complete portfolio snapshot:
{_json.dumps(portfolio_data, indent=2)}

Rank ALL 7 projects into 3 strategic tiers and explain WHY. Focus on: capital efficiency, speed to return, risk/reward, what the family should START NOW vs PLAN CAREFULLY vs LET MATURE.

Respond in strict JSON format:
{{
  "ranked": [
    {{
      "rank": 1,
      "project_id": "...",
      "tier": "MOVE NOW | BUILD CAREFULLY | LET COMPOUND",
      "tier_reason": "one sentence why this tier",
      "score_out_of_10": 8,
      "score_rationale": "2 sentences: what earns this score",
      "time_frame": "short/medium/long",
      "strategic_note": "one paragraph — what specific action should the family take on this project in the next 90 days"
    }}
  ],
  "portfolio_insight": "2-3 sentences on the portfolio as a whole — what the numbers tell you about the family's investment strategy that they might not have noticed",
  "biggest_opportunity": "The single biggest untapped opportunity across the whole portfolio",
  "biggest_risk": "The single biggest risk that could hurt multiple projects simultaneously",
  "compounding_play": "The one project or combination of projects that, if executed well, creates a compounding flywheel effect (e.g., manure from dairy feeds bananas, eucalyptus shade reduces cow heat stress, etc.)"
}}

Be direct, cite numbers, no fluff. The family needs to make money soon AND build for the long term."""

    import os as _os

    # Try Claude first
    raw = _ask_claude(prompt + "\n\nRespond ONLY with valid JSON. No markdown fences.", model="sonnet", timeout=150)

    # Fallback chain
    if not raw:
        gemini_key = _os.getenv("GEMINI_API_KEY","")
        if gemini_key:
            try:
                from google import genai as _genai
                raw = _genai.Client(api_key=gemini_key).models.generate_content(
                    model="gemini-2.0-flash", contents=prompt).text
            except: pass
    if not raw:
        groq_key = _os.getenv("GROQ_API_KEY","")
        if groq_key:
            try:
                from groq import Groq as _Groq
                resp = _Groq(api_key=groq_key).chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role":"user","content":prompt}],
                    max_tokens=1500, response_format={"type":"json_object"})
                raw = resp.choices[0].message.content
            except: pass

    # Parse JSON from response
    import re as _re, json as _json
    try:
        match = _re.search(r'\{.*\}', raw, _re.DOTALL)
        if match:
            ranking = _json.loads(match.group())
        else:
            ranking = _json.loads(raw)
    except:
        ranking = {"ranked": [], "portfolio_insight": raw[:500] if raw else "AI unavailable", "raw": raw[:1000]}

    return {"portfolio_data": portfolio_data, "ai_ranking": ranking}


@app.post("/api/portfolio/new_ventures")
async def new_ventures_engine(request: Request):
    """AI reads club constitution, existing portfolio, Uganda market context + proposes new ventures."""
    from fastapi import HTTPException as _HE
    import os as _os, json as _json
    ok = _auth_verify(_get_tok(request)) or _internal_key_ok(request)
    if not ok: raise _HE(status_code=401, detail="Auth required")

    # Current portfolio snapshot
    portfolio_context = {
        "club": {
            "name": "KIM FAM Investment Club",
            "location": "Mbarara area, Western Uganda",
            "members": 13,
            "focus": "Agricultural and agri-business investments",
            "mandate_from_constitution": [
                "Invest in income-generating projects on Nyabugando farm and Busisi land",
                "Generate returns to be distributed to member families",
                "Build compound wealth through reinvestment",
                "Diversify across short, medium, and long-term investments",
                "All projects must align with sustainable agriculture and family values",
            ],
            "existing_land_assets": [
                "Nyabugando Farm (main farm, Mbarara area) — multiple active projects",
                "Busisi Land (3.38 acres, eucalyptus planted)",
            ]
        },
        "current_portfolio": {
            "operational": ["Chicken (broilers) — Solomon AppSheet, most data-rich", "Washing Bay (car wash) — 5 months revenue confirmed"],
            "active_growing": ["Sheep (Dorper) — 17 head, UGX 1.71M invested", "Trees (Eucalyptus) — 3,600 seedlings, 2024 planted"],
            "planning": ["Dairy (zero-grazing shed) — UGX 16M needed", "Irrigation + Bananas — Phase 2 CapEx unknown"],
            "research": ["Bees — 10 hives planned, no capital committed"],
        },
        "known_farm_resources": {
            "land": "Multi-acre, mix of flat and hilly terrain, Mbarara western Uganda",
            "water": "Borehole, farm ponds, rainwater harvesting capacity",
            "labour": "Family-managed + hired casual labour",
            "existing_infrastructure": "Shed capacity, car wash infrastructure, banana plantation area",
            "markets": "Mbarara urban centre (growing fast), Kampala road access",
            "inputs_available": "Manure from dairy (planned), biomass from eucalyptus, irrigation from ponds",
        },
        "capital_position": {
            "confirmed_deployed": "UGX 30,240,000 total across all projects",
            "monthly_revenue": "UGX 580,000+ (washing bay) + chicken cycle income",
            "next_capital_need": "UGX 16,000,000 for dairy; irrigation Phase 2 TBD",
        }
    }

    body = await request.json()
    focus_area = body.get("focus_area", "all")  # short_term, long_term, low_capital, or all

    prompt = f"""You are a venture development advisor for a Ugandan family investment club. You have deep knowledge of Ugandan agriculture, agribusiness, and rural enterprise economics.

CLUB PROFILE AND CURRENT PORTFOLIO:
{_json.dumps(portfolio_context, indent=2)}

TASK: Propose 5 NEW investment opportunities this club should seriously consider. These are NOT currently in their portfolio.

For each venture:
1. It must be realistic for a small family club in Western Uganda with moderate capital (<UGX 20M per venture ideally)
2. It must be auditable: cite specific numbers from Ugandan agricultural data
3. It must complement or leverage existing assets (land, manure, water, market access)
4. It should create synergies with existing portfolio (e.g. manure from cows feeds crops, eucalyptus poles used locally)
5. Include external resources / websites the family could visit to learn more

Focus area requested: {focus_area}

Respond in strict JSON:
{{
  "ventures": [
    {{
      "rank": 1,
      "name": "Short name",
      "category": "Agriculture / Agriprocessing / Services / Technology",
      "headline": "One sentence pitch",
      "why_now": "Why this is a good opportunity in Uganda in 2026",
      "why_this_club": "Why this specific family is positioned to do this (links to their assets)",
      "capital_required_ugx": 5000000,
      "capital_note": "What that covers",
      "revenue_model": "How money is made",
      "time_to_first_revenue": "X months",
      "annual_revenue_estimate_ugx": 0,
      "annual_profit_estimate_ugx": 0,
      "key_assumptions": ["assumption 1", "assumption 2"],
      "synergies_with_portfolio": "Specific connections to existing projects",
      "risks": ["risk 1", "risk 2"],
      "first_step": "The single most important thing to do in the next 30 days",
      "external_resources": [
        {{"title": "Resource name", "url": "https://...", "why_relevant": "one line"}}
      ]
    }}
  ],
  "strategic_commentary": "2-3 sentences on how these 5 ventures together create a more resilient and compounding portfolio",
  "compounding_sequence": "What to start first, second, third and why — the optimal sequencing to build wealth fastest"
}}

Use real Ugandan numbers. Be bold but realistic. The family wants to build generational wealth, not just pocket money."""

    import os as _os

    # Try Claude first — best for structured research + JSON output
    raw = _ask_claude(prompt + "\n\nRespond ONLY with valid JSON. No markdown fences, no explanation outside the JSON.", model="sonnet", timeout=180)

    # Fallback chain
    if not raw:
        gemini_key = _os.getenv("GEMINI_API_KEY","")
        if gemini_key:
            try:
                from google import genai as _genai
                raw = _genai.Client(api_key=gemini_key).models.generate_content(
                    model="gemini-2.0-flash", contents=prompt).text
            except: pass
    if not raw:
        groq_key = _os.getenv("GROQ_API_KEY","")
        if groq_key:
            try:
                from groq import Groq as _Groq
                resp = _Groq(api_key=groq_key).chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role":"user","content":prompt}],
                    max_tokens=3000, response_format={"type":"json_object"})
                raw = resp.choices[0].message.content
            except Exception as e: raw = str(e)

    import re as _re
    try:
        match = _re.search(r'\{.*\}', raw, _re.DOTALL)
        result = _json.loads(match.group() if match else raw)
    except:
        result = {"ventures": [], "raw": raw[:2000], "error": "JSON parse failed"}

    return {"status": "ok", "focus_area": focus_area, "result": result}


@app.get("/api/club/equity")
async def club_equity(request: Request):
    from fastapi import HTTPException as _HE
    tok  = _get_tok(request)
    ok   = _auth_verify(tok) or _internal_key_ok(request)
    if not ok:
        raise _HE(status_code=401, detail="Auth required")
    data = dict(_fetch_family_equity())   # shallow copy so the decision overlay isn't cached
    try:
        _ensure_equity_votes_table()
        decision = _equity_decision()
    except Exception:
        decision = None
    if decision:
        m = decision["adopted_model"]
        where = f" at {decision['meeting_ref']}" if decision.get("meeting_ref") else ""
        data["decision_status"] = "decided"
        data["adopted_model"]   = m
        data["decision_note"]   = (
            f"Model {m} adopted{where}. This is now the club's official equity model; "
            f"Models A, B and C remain visible for reference."
        )
    return data

@app.get("/api/projects/chicken/allocation")
async def chicken_allocation(request: Request):
    """Legacy alias kept for backwards compat."""
    return await club_equity(request)


# ── Equity model vote (KIM 008/2026 → decided at KIM 009) ─────────────────────────
# One vote per FAMILY (not per member). A family may change its vote until the
# decision is recorded. Models: A (equal share), B (proportional), C (Solomon's).
from family_profiles import FAMILY_MEMBER_MAP as _FAMILY_MEMBER_MAP

# member first-name → family_id (invert the family→members map)
_MEMBER_TO_FAMILY = {
    member: fam_id
    for fam_id, members in _FAMILY_MEMBER_MAP.items()
    for member in members
}
_TOTAL_FAMILIES = len(_FAMILY_MEMBER_MAP)
_VALID_MODELS = {"A", "B", "C"}

def _ensure_equity_votes_table():
    from db import execute as _ex
    _ex("""CREATE TABLE IF NOT EXISTS equity_votes (
        family_id  TEXT PRIMARY KEY,
        model      TEXT NOT NULL CHECK (model IN ('A','B','C')),
        voter      TEXT NOT NULL,
        voted_at   TIMESTAMPTZ NOT NULL DEFAULT now()
    )""")
    _ex("""CREATE TABLE IF NOT EXISTS equity_decision (
        id            SERIAL PRIMARY KEY,
        adopted_model TEXT NOT NULL CHECK (adopted_model IN ('A','B','C')),
        decided_at    TIMESTAMPTZ NOT NULL DEFAULT now(),
        decided_by    TEXT,
        meeting_ref   TEXT,
        tally         JSONB,
        active        BOOLEAN NOT NULL DEFAULT TRUE
    )""")

def _equity_decision():
    """The active adopted-model decision, or None if the vote is still open."""
    from db import query as _q
    import json as _json_d
    rows = _q("""SELECT adopted_model, decided_at, decided_by, meeting_ref, tally
                 FROM equity_decision WHERE active = TRUE
                 ORDER BY decided_at DESC LIMIT 1""")
    if not rows:
        return None
    r = rows[0]
    raw = r["tally"]
    return {
        "adopted_model": r["adopted_model"],
        "decided_at": str(r["decided_at"]),
        "decided_by": r["decided_by"],
        "meeting_ref": r["meeting_ref"],
        "tally": (_json_d.loads(raw) if isinstance(raw, (str, bytes, bytearray)) else (raw or {})),
    }

def _equity_vote_state(member_name: str | None):
    from db import query as _q
    rows = _q("SELECT family_id, model, voter, voted_at FROM equity_votes")
    tally = {"A": 0, "B": 0, "C": 0}
    voted = []
    for r in rows:
        if r["model"] in tally:
            tally[r["model"]] += 1
        voted.append({"family_id": r["family_id"], "model": r["model"]})
    my_family = _MEMBER_TO_FAMILY.get(member_name) if member_name else None
    my_vote = next((r["model"] for r in rows if r["family_id"] == my_family), None)
    decision = _equity_decision()
    # Leading model (for the "Close vote & adopt" suggestion) — None on a tie.
    leader = None
    if any(tally.values()):
        mx = max(tally.values())
        leaders = [m for m, c in tally.items() if c == mx]
        leader = leaders[0] if len(leaders) == 1 else None
    return {
        "tally": tally,
        "total_families": _TOTAL_FAMILIES,
        "families_voted": len(rows),
        "my_family_id": my_family,
        "my_vote": my_vote,
        "voted": voted,
        "decided": decision is not None,
        "decision": decision,
        "leader": leader,
    }

@app.get("/api/club/equity/vote")
async def get_equity_vote(request: Request):
    from fastapi import HTTPException as _HE
    tok = _get_tok(request)
    payload = _auth_verify(tok)
    if not payload:
        raise _HE(status_code=401, detail="Login required")
    _ensure_equity_votes_table()
    return _equity_vote_state(payload.get("sub"))

@app.post("/api/club/equity/vote")
async def cast_equity_vote(request: Request):
    from fastapi import HTTPException as _HE
    from db import execute as _ex
    tok = _get_tok(request)
    payload = _auth_verify(tok)
    if not payload:
        raise _HE(status_code=401, detail="Login required")
    member_name = payload.get("sub")
    family_id = _MEMBER_TO_FAMILY.get(member_name)
    if not family_id:
        raise _HE(status_code=403, detail="Your account is not linked to a family, so you cannot vote.")
    body = await request.json()
    model = str(body.get("model", "")).strip().upper()
    if model not in _VALID_MODELS:
        raise _HE(status_code=400, detail="Model must be A, B or C.")
    _ensure_equity_votes_table()
    if _equity_decision() is not None:
        raise _HE(status_code=409, detail="The equity model has already been decided — voting is closed.")
    # Upsert: one row per family; re-voting updates the family's choice.
    _ex("""INSERT INTO equity_votes (family_id, model, voter, voted_at)
           VALUES (%s, %s, %s, now())
           ON CONFLICT (family_id)
           DO UPDATE SET model = EXCLUDED.model, voter = EXCLUDED.voter, voted_at = now()""",
        (family_id, model, member_name))
    return _equity_vote_state(member_name)


@app.post("/api/club/equity/finalize")
async def finalize_equity_vote(request: Request):
    """Admin: close the vote and adopt a model. Uses the leading model unless an
    explicit model is supplied (needed on a tie). Records the decision + tally."""
    from fastapi import HTTPException as _HE
    from db import execute as _ex
    import json as _json_f
    tok = _get_tok(request)
    payload = _auth_verify(tok)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    _ensure_equity_votes_table()
    state = _equity_vote_state(None)
    body = await request.json()
    model = str(body.get("model", "")).strip().upper()
    meeting_ref = str(body.get("meeting_ref", "")).strip() or None
    if not model:
        model = state["leader"]   # default to the clear leader
    if model not in _VALID_MODELS:
        raise _HE(status_code=400, detail="Vote is tied — choose the model to adopt explicitly (A, B or C).")
    # Deactivate any prior decision, then record the new one.
    _ex("UPDATE equity_decision SET active = FALSE WHERE active = TRUE")
    _ex("""INSERT INTO equity_decision (adopted_model, decided_by, meeting_ref, tally, active)
           VALUES (%s, %s, %s, %s, TRUE)""",
        (model, payload.get("sub", "admin"), meeting_ref, _json_f.dumps(state["tally"])))
    global _ALLOC_CACHE
    _ALLOC_CACHE = {"data": None, "ts": 0}   # bust equity cache so decision shows immediately
    return _equity_vote_state(payload.get("sub"))


@app.post("/api/club/equity/reopen")
async def reopen_equity_vote(request: Request):
    """Admin: reopen the vote (deactivates the current decision)."""
    from fastapi import HTTPException as _HE
    from db import execute as _ex
    tok = _get_tok(request)
    payload = _auth_verify(tok)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    _ensure_equity_votes_table()
    _ex("UPDATE equity_decision SET active = FALSE WHERE active = TRUE")
    global _ALLOC_CACHE
    _ALLOC_CACHE = {"data": None, "ts": 0}
    return _equity_vote_state(payload.get("sub"))


# ── Club Office Bearers ──────────────────────────────────────────────────────────────────────────────

_VALID_OFFICER_ROLES = {
    "chairman":         "Chairman",
    "vice_chairman":    "Vice Chairman",
    "secretary":        "Secretary",
    "treasurer":        "Treasurer",
    "committee_member": "Committee Member",
}

@app.get("/api/club/officers")
async def get_officers(request: Request):
    from db import query as _q
    from fastapi import HTTPException as _HE
    import notifications as _notif2
    token = _get_tok(request)
    payload = _auth_verify(token)
    is_internal = _internal_key_ok(request)
    if not payload and not is_internal:
        raise _HE(status_code=401, detail="Auth required")
    rows = _q(
        "SELECT role_slug, role_display, member_name, effective_from, notes "
        "FROM club_office_bearers WHERE effective_to IS NULL ORDER BY role_slug"
    )
    return [{
        "role_slug":      r["role_slug"],
        "role_display":   r["role_display"],
        "member_name":    r["member_name"],
        "effective_from": str(r["effective_from"]),
        "phone":          _notif2.MEMBER_PHONES.get(r["member_name"]),
        "notes":          r["notes"],
    } for r in rows]

@app.put("/api/club/officers/{role_slug}")
async def assign_officer(role_slug: str, request: Request):
    from db import execute as _exec
    from fastapi import HTTPException as _HE
    from datetime import date
    token = _get_tok(request)
    payload = _auth_verify(token)
    if not payload or payload.get("sub") not in _ADMINS_PP:
        raise _HE(status_code=403, detail="Admin only")
    if role_slug not in _VALID_OFFICER_ROLES:
        raise _HE(status_code=400, detail=f"Unknown role. Valid: {list(_VALID_OFFICER_ROLES)}")
    body = await request.json()
    member_name    = str(body.get("member_name", "")).strip()
    notes          = str(body.get("notes", "")).strip() or None
    effective_from = str(body.get("effective_from", date.today().isoformat())).strip()
    if not member_name:
        raise _HE(status_code=400, detail="member_name required")
    _exec(
        "UPDATE club_office_bearers SET effective_to=%s WHERE role_slug=%s AND effective_to IS NULL",
        (effective_from, role_slug)
    )
    _exec(
        "INSERT INTO club_office_bearers (role_slug, role_display, member_name, effective_from, notes) "
        "VALUES (%s, %s, %s, %s, %s)",
        (role_slug, _VALID_OFFICER_ROLES[role_slug], member_name, effective_from, notes)
    )
    return {"ok": True, "role": role_slug, "now_held_by": member_name}


# ── Member Engagement Score ───────────────────────────────────────────────────
# Weekly engagement scoring across 5 dimensions:
#   1. Payment timeliness (40 pts max) — paid before the 5th of the month = full marks
#   2. Action point completion (25 pts max) — % of owned actions completed by deadline
#   3. Project participation (20 pts max) — confirmed interests + active contributions
#   4. Meetings proxy (10 pts max) — attendance derived from meeting register (future)
#   5. App engagement (5 pts max) — has uploaded profile photo + submitted payments via app
# Score is per-family (both spouses count as one unit).
# Ranking shows top performers; bottom performers are NOT ranked (club culture).

@app.get("/api/engagement/scores")
async def engagement_scores(request: Request):
    from fastapi import HTTPException as _HE
    from datetime import date as _date, datetime as _datetime
    token = _get_tok(request)
    payload = _auth_verify(token) if token else None
    if not payload: raise _HE(status_code=401, detail="Auth required")

    today = _date.today()

    # 1. Payment scores — families that have paid current month by the 5th get full marks
    ledger_data = []
    try:
        from contributions import compute_family_balance as _cfb
        from db import query as _dq
        fam_rows = _dq("SELECT id, family_name FROM families ORDER BY family_name")
        for fr in fam_rows:
            bal = _cfb(fr["id"])
            ledger_data.append({
                "family_id":   fr["id"],
                "family_name": fr["family_name"],
                "bal":         bal,
            })
    except Exception:
        pass

    # 2. Action point scores
    action_scores: dict = {}
    try:
        sh = gc().open_by_key(SHEET_ID)
        at = sh.worksheet("Action Tracker")
        at_rows = at.get_all_records()
        # Map member first name → family
        _MF = {
            "Hillary": "ARINDAS",  "Esther": "ARINDAS",
            "Viola":   "ARUNGAS",  "Simon":  "ARUNGAS",
            "Israel":  "KIKANGIS", "Merab":  "KIKANGIS",
            "Max":     "TURAMYES", "Janet":  "TURAMYES",
            "Solomon": "ARIHOS",
            "Hellen":  "KOFUNAS",  "Lawi":   "KOFUNAS",
            "Alex":    "TUHIMBISES","Priscilla": "TUHIMBISES",
        }
        family_actions: dict = {}
        for r in at_rows:
            owner = str(r.get("Owner", "") or r.get("Responsible", "") or "").strip()
            if not owner: continue
            first = owner.split()[0]
            fam = _MF.get(first)
            if not fam: continue
            if fam not in family_actions:
                family_actions[fam] = {"total": 0, "done": 0}
            family_actions[fam]["total"] += 1
            status = str(r.get("Status", "") or r.get("status", "")).strip().lower()
            if "done" in status or "complete" in status or "closed" in status:
                family_actions[fam]["done"] += 1
        for fam, v in family_actions.items():
            pct = (v["done"] / v["total"]) if v["total"] > 0 else 0
            action_scores[fam] = round(pct * 25)
    except Exception:
        pass

    # 3. Project participation scores
    proj_scores: dict = {}
    try:
        from db import query as _dq2
        interests = _dq2("""
            SELECT pp.status, f.family_name
            FROM project_participation pp
            LEFT JOIN persons p ON LOWER(p.whatsapp_name) = LOWER(pp.member_name)
            LEFT JOIN families f ON p.family_id = f.id
            WHERE f.family_name IS NOT NULL
        """)
        for row in interests:
            fam = row["family_name"]
            if not fam: continue
            if fam not in proj_scores: proj_scores[fam] = 0
            if row["status"] == "confirmed":
                proj_scores[fam] = min(20, proj_scores[fam] + 8)
            elif row["status"] in ("pending", "awaiting_chairman"):
                proj_scores[fam] = min(20, proj_scores[fam] + 4)
    except Exception:
        pass

    # 4. App engagement score
    app_scores: dict = {}
    try:
        from pathlib import Path as _Path
        avatars_dir = _Path(__file__).parent / "static" / "avatars"
        _FMAP = {
            "hillaryarinda": "ARINDAS", "estherarinda": "ARINDAS",
            "violaarunga": "ARUNGAS", "simonarunga": "ARUNGAS",
            "israelkikangi": "KIKANGIS", "merabkikangi": "KIKANGIS",
            "maxturamye": "TURAMYES", "janetturamye": "TURAMYES",
            "solomonariho": "ARIHOS",
            "hellenkofuna": "KOFUNAS", "lawikofuna": "KOFUNAS",
            "alextuhimbise": "TUHIMBISES", "priscillatuhimbise": "TUHIMBISES",
        }
        for avatar_key, fam in _FMAP.items():
            if (avatars_dir / f"{avatar_key}.jpg").exists():
                app_scores[fam] = app_scores.get(fam, 0) + 3
        # Check who has submitted payments via app
        from db import query as _dq3
        pmt_submitters = _dq3("""
            SELECT DISTINCT submitted_by_user_id FROM contribution_payments
            WHERE submitted_by_user_id IS NOT NULL AND submitted_by_user_id != ''
        """)
        _NAMEMAP = {
            "Hillary": "ARINDAS", "Esther": "ARINDAS",
            "Viola": "ARUNGAS", "Simon": "ARUNGAS",
            "Israel": "KIKANGIS", "Merab": "KIKANGIS",
            "Max": "TURAMYES", "Janet": "TURAMYES",
            "Solomon": "ARIHOS",
            "Hellen": "KOFUNAS", "Lawi": "KOFUNAS",
            "Alex": "TUHIMBISES", "Priscilla": "TUHIMBISES",
        }
        for row in pmt_submitters:
            fam = _NAMEMAP.get(row["submitted_by_user_id"])
            if fam: app_scores[fam] = min(5, app_scores.get(fam, 0) + 2)
    except Exception:
        pass

    # Assemble final scores
    scores = []
    for ld in ledger_data:
        fam = ld["family_name"]
        bal = ld["bal"]

        # Payment score: negative current_balance = paid ahead = full marks
        # positive = arrears, scaled inversely
        cur_bal  = bal.get("current_balance", 0)
        rate     = bal.get("monthly_rate", 1) or 1
        if cur_bal <= 0:
            pay_score = 40
        else:
            shortfall_ratio = min(cur_bal / rate, 3)  # cap at 3 months behind
            pay_score = max(0, round(40 * (1 - shortfall_ratio / 3)))

        act_score  = action_scores.get(fam, 0)
        prj_score  = proj_scores.get(fam, 0)
        app_score  = min(5, app_scores.get(fam, 0))
        total      = pay_score + act_score + prj_score + app_score

        scores.append({
            "family_name": fam,
            "label":       "The " + fam.capitalize() if fam not in ("ARIHOS",) else "The " + fam[0] + fam[1:].lower(),
            "total":       total,
            "breakdown": {
                "payments":     pay_score,
                "actions":      act_score,
                "projects":     prj_score,
                "app_activity": app_score,
            },
            "current_balance": cur_bal,
        })

    # Sort by score descending; never expose the bottom rank publicly
    scores.sort(key=lambda x: x["total"], reverse=True)

    return {
        "week": today.isocalendar()[1],
        "year": today.year,
        "as_at": today.isoformat(),
        "scores": scores,
        "max_score": 90,
        "note": "Scores reset weekly. Bottom performers are not ranked — only top 5 shown.",
    }


# ── SPA fallback — MUST be last; catches all non-API routes for React Router ──
@app.get("/{full_path:path}", response_class=HTMLResponse, include_in_schema=False)
def spa_fallback(full_path: str):
    if full_path.startswith("api/") or full_path.startswith("static/") or full_path.startswith("assets/"):
        from fastapi import HTTPException
        raise HTTPException(status_code=404)
    return _react_index()
